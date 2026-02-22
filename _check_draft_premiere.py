# -*- coding: utf-8 -*-
"""Check draft premiere charts."""
from pptx import Presentation
from lxml import etree
import tempfile, os, re

from data_reader import read_excel_data
from generate_report import generate_report

data = read_excel_data("input_data/20260215CCTV-17日数据（简版）.xlsx")
tmp = tempfile.mktemp(suffix=".pptx")
generate_report("input_data/20260215CCTV-17日数据（简版）.xlsx", output_path=tmp)
prs = Presentation(tmp)

ns_c = "http://schemas.openxmlformats.org/drawingml/2006/chart"
for pg in range(len(prs.slides)):
    slide = prs.slides[pg]
    for shp in slide.shapes:
        if getattr(shp, "has_chart", False):
            for i, s in enumerate(shp.chart.series):
                tx = s._element.find(f".//{{{ns_c}}}tx//{{{ns_c}}}v")
                name = tx.text if tx is not None else "N/A"
                if "首播" in name or "前一日" in name or "2月15" in name or "前一" in name:
                    ser_xml = etree.tostring(s._element).decode()
                    fills = re.findall(r'srgbClr val="([^"]+)"', ser_xml)
                    print(f"Draft p{pg+1} s[{i}] name={name} fills={fills[:3]}")

# Also check series count
for pg in range(len(prs.slides)):
    slide = prs.slides[pg]
    for shp in slide.shapes:
        if getattr(shp, "has_chart", False):
            n_ser = len(list(shp.chart.series))
            cats = list(shp.chart.plots[0].categories)
            if cats and ("大地" in str(cats[0]) or "共富" in str(cats[0])):
                print(f"Draft p{pg+1}: {n_ser} series, {len(cats)} cats, first={cats[0]}")

os.unlink(tmp)
