# -*- coding: utf-8 -*-
"""Full XML dump of demo premiere series fills."""
from pptx import Presentation
from lxml import etree

C_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

for label, fname, pg_idx in [
    ("DEMO", "output_ppt/CCTV-17频道收视日报 0210 demo（简版）.pptx", 10),
    ("GEN", "output_ppt/result_0215_v6.pptx", 10),
]:
    prs = Presentation(fname)
    slide = prs.slides[pg_idx]
    print(f"\n=== {label} page {pg_idx+1} ===")
    for shp in slide.shapes:
        if getattr(shp, "has_chart", False):
            for i, s in enumerate(shp.chart.series):
                ser = s._element
                sp = ser.find(f"{{{C_NS}}}spPr")
                if sp is not None:
                    solid = sp.find(f"{{{A_NS}}}solidFill")
                    if solid is not None:
                        print(f"\n  series[{i}] solidFill XML:")
                        print("    " + etree.tostring(solid, pretty_print=True).decode().strip())
