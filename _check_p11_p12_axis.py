"""检查 demo 模板 P11/P12 图表坐标轴设置"""
from pptx import Presentation
from lxml import etree

C = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
prs = Presentation('output_ppt/CCTV-17频道收视日报 0210 demo（简版）.pptx')

for idx in [10, 11]:  # P11, P12
    slide = prs.slides[idx]
    print(f'=== Slide {idx+1} ===')
    for shp in slide.shapes:
        if not getattr(shp, 'has_chart', False):
            continue
        cs = shp.chart._chartSpace
        pa = cs.find(f'.//{{{C}}}plotArea')
        # Check all value axes
        for vax in pa.findall(f'{{{C}}}valAx'):
            ax_id = vax.find(f'{{{C}}}axId')
            scl = vax.find(f'{{{C}}}scaling')
            if scl is not None:
                mx = scl.find(f'{{{C}}}max')
                mn = scl.find(f'{{{C}}}min')
                mx_val = mx.get('val') if mx is not None else 'auto'
                mn_val = mn.get('val') if mn is not None else 'auto'
                ax_val = ax_id.get('val') if ax_id is not None else '?'
                print(f'  valAx id={ax_val}: max={mx_val}, min={mn_val}')
                dl = vax.find(f'{{{C}}}delete')
                dl_val = dl.get('val') if dl is not None else 'unset'
                print(f'    delete={dl_val}')
            else:
                print(f'  valAx: no scaling')
        # Check chart types
        for ct in ['barChart', 'lineChart', 'bar3DChart']:
            node = pa.find(f'{{{C}}}{ct}')
            if node is not None:
                sers = node.findall(f'{{{C}}}ser')
                print(f'  {ct}: {len(sers)} series')
                for ser in sers:
                    nm = ser.find(f'{{{C}}}tx/{{{C}}}strRef/{{{C}}}strCache/{{{C}}}pt/{{{C}}}v')
                    vals = []
                    for pt in ser.findall(f'{{{C}}}val/{{{C}}}numRef/{{{C}}}numCache/{{{C}}}pt'):
                        v = pt.find(f'{{{C}}}v')
                        if v is not None and v.text:
                            try:
                                vals.append(float(v.text))
                            except ValueError:
                                pass
                    sname = nm.text if nm is not None else '?'
                    vmax = max(vals) if vals else 0
                    print(f'    ser "{sname}": max={vmax:.4f}')

# Also check the generated output
print('\n\n=== Generated output ===')
prs2 = Presentation('output_ppt/test_fix_p4_halfrow_red.pptx')
for idx in [10, 11]:
    slide = prs2.slides[idx]
    print(f'=== Slide {idx+1} ===')
    for shp in slide.shapes:
        if not getattr(shp, 'has_chart', False):
            continue
        cs = shp.chart._chartSpace
        pa = cs.find(f'.//{{{C}}}plotArea')
        for vax in pa.findall(f'{{{C}}}valAx'):
            ax_id = vax.find(f'{{{C}}}axId')
            scl = vax.find(f'{{{C}}}scaling')
            if scl is not None:
                mx = scl.find(f'{{{C}}}max')
                mn = scl.find(f'{{{C}}}min')
                mx_val = mx.get('val') if mx is not None else 'auto'
                mn_val = mn.get('val') if mn is not None else 'auto'
                ax_val = ax_id.get('val') if ax_id is not None else '?'
                print(f'  valAx id={ax_val}: max={mx_val}, min={mn_val}')
        for ct in ['barChart', 'lineChart']:
            node = pa.find(f'{{{C}}}{ct}')
            if node is not None:
                sers = node.findall(f'{{{C}}}ser')
                print(f'  {ct}: {len(sers)} series')
                for ser in sers:
                    nm = ser.find(f'{{{C}}}tx/{{{C}}}strRef/{{{C}}}strCache/{{{C}}}pt/{{{C}}}v')
                    vals = []
                    for pt in ser.findall(f'{{{C}}}val/{{{C}}}numRef/{{{C}}}numCache/{{{C}}}pt'):
                        v = pt.find(f'{{{C}}}v')
                        if v is not None and v.text:
                            try:
                                vals.append(float(v.text))
                            except ValueError:
                                pass
                    sname = nm.text if nm is not None else '?'
                    vmax = max(vals) if vals else 0
                    print(f'    ser "{sname}": max={vmax:.4f}')
