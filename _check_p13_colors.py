# -*- coding: utf-8 -*-
"""Check page 13 label colors."""
from pptx import Presentation
from lxml import etree

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

for label, fname in [
    ("DEMO", "output_ppt/CCTV-17频道收视日报 0210 demo（简版）.pptx"),
    ("GEN", "output_ppt/result_0215_v7.pptx"),
]:
    prs = Presentation(fname)
    slide = prs.slides[12]
    print(f"\n=== {label} ===")
    for shp in slide.shapes:
        if getattr(shp, "has_text_frame", False):
            t = shp.text.strip()
            if t and ("%" in t):
                txBody = shp._element.find(f".//{{{A_NS}}}txBody")
                if txBody is None:
                    continue
                runs = txBody.findall(f".//{{{A_NS}}}r")
                color = "none"
                for r in runs[:1]:
                    rPr = r.find(f"{{{A_NS}}}rPr")
                    if rPr is not None:
                        sf = rPr.find(f"{{{A_NS}}}solidFill")
                        if sf is not None:
                            ch = sf[0]
                            tag = ch.tag.split("}")[1]
                            val = ch.get("val")
                            color = f"{tag}={val}"
                print(f"  text={t!r:14s}  color={color:20s}  name={shp.name}")
