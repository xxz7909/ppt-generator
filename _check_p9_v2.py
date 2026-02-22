# -*- coding: utf-8 -*-
"""Verify page 9 program table in result_0215_v5."""
from pptx import Presentation
from pptx.util import Emu, Cm
from lxml import etree
import re

prs = Presentation("output_ppt/result_0215_v5.pptx")
slide = prs.slides[8]

print("=== Page 9 shapes ===")
for i, shp in enumerate(slide.shapes):
    tp = "table" if getattr(shp, "has_table", False) else ("chart" if getattr(shp, "has_chart", False) else "other")
    print(f"  [{i}] {tp}  left={shp.left}  top={shp.top}  w={shp.width}  h={shp.height}")
    
    if tp == "chart":
        chart = shp.chart
        cats = list(chart.plots[0].categories)
        print(f"      time: first={cats[0]}, last={cats[-1]}, count={len(cats)}")
        chart_xml = etree.tostring(chart._chartSpace).decode()
        layout_match = re.search(r'<c:plotArea>(.*?)</c:plotArea>', chart_xml, re.DOTALL)
        if layout_match:
            x_m = re.search(r'<c:x val="([^"]+)"', layout_match.group(1))
            w_m = re.search(r'<c:w val="([^"]+)"', layout_match.group(1))
            if x_m and w_m:
                pa_x = float(x_m.group(1))
                pa_w = float(w_m.group(1))
                plot_left = shp.left + int(pa_x * shp.width)
                plot_width = int(pa_w * shp.width)
                print(f"      plotArea: left={plot_left} width={plot_width} right={plot_left+plot_width}")
    
    if tp == "table":
        tbl = shp.table
        ncols = len(tbl.columns)
        if ncols == 20:
            print(f"    PROGRAM TABLE: left={shp.left} width={shp.width}")
            # Check for 国歌, 歌曲, 再见
            has_hidden = False
            for j in range(ncols):
                t = tbl.rows[0].cells[j].text[:50]
                w = tbl.columns[j].width
                if '国歌' in t or '歌曲' in t or '再见' in t:
                    has_hidden = True
                    print(f"      *** HIDDEN FOUND: col[{j}] text={repr(t)}")
                # Check alignment (center)
                tc = tbl.cell(0, j)._tc
                a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                pPr = tc.find(f'.//{{{a_ns}}}pPr')
                algn = pPr.get('algn', 'NONE') if pPr is not None else 'NONE'
                print(f"      col[{j:2d}] w={w:>8d} algn={algn:4s} text={repr(t[:40])}")
            if not has_hidden:
                print("    ✓ No hidden programs (国歌/歌曲/再见) found")
        elif ncols == 8:
            cell0 = tbl.rows[0].cells[0].text[:20]
            print(f"    8-col table: left={shp.left} cell0={repr(cell0)}")
            for j in range(ncols):
                print(f"      col[{j}] w={tbl.columns[j].width}")
