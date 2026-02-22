# -*- coding: utf-8 -*-
"""Check demo vs generated premiere pages (11-12)."""
from pptx import Presentation
from lxml import etree
import re

for label, fname in [("DEMO", "output_ppt/CCTV-17频道收视日报 0210 demo（简版）.pptx"),
                       ("GEN", "output_ppt/result_0215_v6.pptx")]:
    prs = Presentation(fname)
    for pg in [10, 11]:
        slide = prs.slides[pg]
        print(f"=== {label} page {pg+1} ===")
        for shp in slide.shapes:
            if getattr(shp, "has_chart", False):
                chart = shp.chart
                ns_c = "http://schemas.openxmlformats.org/drawingml/2006/chart"
                for i, s in enumerate(chart.series):
                    ser_xml = etree.tostring(s._element, pretty_print=True).decode()
                    fills = re.findall(r'<a:srgbClr val="([^"]+)"', ser_xml)
                    tx = s._element.find(f".//{{{ns_c}}}tx//{{{ns_c}}}v")
                    name = tx.text if tx is not None else "N/A"
                    print(f"  series[{i}] name={name}  fills={fills[:3]}")
                cats = list(chart.plots[0].categories)
                print(f"  cats: first={cats[0]}, last={cats[-1]}, count={len(cats)}")
            
            # Check text boxes for description
            if shp.has_text_frame:
                txt = shp.text_frame.text[:80]
                if "首播" in txt or "月均" in txt or "前一" in txt:
                    print(f"  TEXT: {txt[:100]}...")
