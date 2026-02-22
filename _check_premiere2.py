# -*- coding: utf-8 -*-
"""Deep check of premiere chart series XML from demo."""
from pptx import Presentation
from lxml import etree
import re

prs = Presentation("output_ppt/CCTV-17频道收视日报 0210 demo（简版）.pptx")
slide = prs.slides[10]  # page 11

for shp in slide.shapes:
    if getattr(shp, "has_chart", False):
        C_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"
        A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
        
        # Get series XML  
        for i, s in enumerate(shp.chart.series):
            ser = s._element
            tx = ser.find(f".//{{{C_NS}}}tx//{{{C_NS}}}v")
            name = tx.text if tx is not None else "N/A"
            
            # Check for fill specification
            sp = ser.find(f"{{{C_NS}}}spPr")
            fill_info = "no spPr"
            if sp is not None:
                solid = sp.find(f"{{{A_NS}}}solidFill")
                if solid is not None:
                    clr = solid.find(f"{{{A_NS}}}srgbClr")
                    if clr is not None:
                        fill_info = f"solidFill srgbClr={clr.get('val')}"
                    else:
                        fill_info = f"solidFill (no srgb): {etree.tostring(solid).decode()[:100]}"
                else:
                    fill_info = f"spPr (no solidFill): {etree.tostring(sp).decode()[:200]}"
            
            # Also check for color in dPt or other places 
            dLbls = ser.find(f"{{{C_NS}}}dLbls")
            print(f"  series[{i}] name={name}  fill={fill_info}  hasDLbls={dLbls is not None}")

# Also check generated result
print("\n=== Generated ===")
prs2 = Presentation("output_ppt/result_0215_v6.pptx")
slide2 = prs2.slides[10]
for shp in slide2.shapes:
    if getattr(shp, "has_chart", False):
        C_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"
        A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
        for i, s in enumerate(shp.chart.series):
            ser = s._element
            tx = ser.find(f".//{{{C_NS}}}tx//{{{C_NS}}}v")
            name = tx.text if tx is not None else "N/A"
            sp = ser.find(f"{{{C_NS}}}spPr")
            fill_info = "no spPr"
            if sp is not None:
                solid = sp.find(f"{{{A_NS}}}solidFill")
                if solid is not None:
                    clr = solid.find(f"{{{A_NS}}}srgbClr")
                    if clr is not None:
                        fill_info = f"solidFill srgbClr={clr.get('val')}"
                    else:
                        fill_info = f"solidFill: {etree.tostring(solid).decode()[:100]}"
                else:
                    fill_info = f"spPr: {etree.tostring(sp).decode()[:200]}"
            print(f"  series[{i}] name={name}  fill={fill_info}")
