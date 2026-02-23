"""测试极端超轴场景 - 验证 _fix_premiere_chart_axis"""
from pptx import Presentation
from lxml import etree

C = 'http://schemas.openxmlformats.org/drawingml/2006/chart'

# 打开生成的 PPT，人为把 P11 数据 x3 使之超过 0.2
prs = Presentation('output_ppt/test_axis_adaptive.pptx')
slide = prs.slides[10]  # P11

for shp in slide.shapes:
    if not getattr(shp, 'has_chart', False):
        continue
    cs = shp.chart._chartSpace
    pa = cs.find(f'.//{{{C}}}plotArea')
    bar_chart = pa.find(f'{{{C}}}barChart')
    for ser in bar_chart.findall(f'{{{C}}}ser'):
        val = ser.find(f'{{{C}}}val')
        if val is None:
            continue
        nr = val.find(f'{{{C}}}numRef')
        if nr is None:
            continue
        nc = nr.find(f'{{{C}}}numCache')
        if nc is None:
            continue
        for pt in nc.findall(f'{{{C}}}pt'):
            v = pt.find(f'{{{C}}}v')
            if v is not None and v.text:
                try:
                    v.text = str(float(v.text) * 3)
                except ValueError:
                    pass
    for vax in pa.findall(f'{{{C}}}valAx'):
        scl = vax.find(f'{{{C}}}scaling')
        mx = scl.find(f'{{{C}}}max')
        val_str = mx.get('val') if mx is not None else 'auto'
        print(f'BEFORE: axis_max={val_str}')

prs.save('output_ppt/test_overflow_p11.pptx')

# 用修复函数处理
import importlib.util, sys
spec2 = importlib.util.spec_from_file_location('gen', 'generate_report_config_driven.py')
gen = importlib.util.module_from_spec(spec2)
sys.modules['gen'] = gen
spec2.loader.exec_module(gen)

prs2 = Presentation('output_ppt/test_overflow_p11.pptx')
s = prs2.slides[10]
gen._fix_premiere_chart_axis(s)

for shp in s.shapes:
    if not getattr(shp, 'has_chart', False):
        continue
    cs = shp.chart._chartSpace
    pa = cs.find(f'.//{{{C}}}plotArea')
    for vax in pa.findall(f'{{{C}}}valAx'):
        scl = vax.find(f'{{{C}}}scaling')
        mx = scl.find(f'{{{C}}}max')
        val_str = mx.get('val') if mx is not None else 'auto'
        print(f'AFTER: axis_max={val_str}')

    bar_chart = pa.find(f'{{{C}}}barChart')
    dmax = 0
    for ser in bar_chart.findall(f'{{{C}}}ser'):
        val = ser.find(f'{{{C}}}val')
        if val is None:
            continue
        nr = val.find(f'{{{C}}}numRef')
        if nr is None:
            continue
        nc = nr.find(f'{{{C}}}numCache')
        if nc is None:
            continue
        for pt in nc.findall(f'{{{C}}}pt'):
            v = pt.find(f'{{{C}}}v')
            if v is not None and v.text:
                try:
                    dmax = max(dmax, float(v.text))
                except ValueError:
                    pass
    print(f'data_max={dmax:.4f}')
    fits = float(mx.get('val')) >= dmax
    print(f'fits_in_axis={fits}')
