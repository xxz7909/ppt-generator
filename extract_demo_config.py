# -*- coding: utf-8 -*-
"""
从demo PPT自动提取可用于生成程序的版式配置。

用法:
    python -X utf8 extract_demo_config.py \
        "output_ppt/CCTV-17频道收视日报 0210 demo（简版）.pptx" \
        --output demo_layout_config.json
"""

import argparse
import json
from pathlib import Path

from pptx import Presentation


EMU_PER_CM = 360000


def emu_to_cm(value):
    return round(value / EMU_PER_CM, 3)


def color_to_rgb(shape):
    try:
        if shape.fill and shape.fill.fore_color and shape.fill.fore_color.rgb:
            rgb = shape.fill.fore_color.rgb
            return [rgb[0], rgb[1], rgb[2]]
    except Exception:
        pass
    return None


def text_runs(paragraph):
    runs = []
    for run in paragraph.runs:
        color = None
        if run.font and run.font.color:
            try:
                rgb = run.font.color.rgb
                if rgb is not None:
                    color = [rgb[0], rgb[1], rgb[2]]
            except Exception:
                color = None
        runs.append(
            {
                "text": run.text,
                "fontName": run.font.name if run.font else None,
                "fontSizePt": float(run.font.size.pt) if run.font and run.font.size else None,
                "bold": bool(run.font.bold) if run.font else None,
                "colorRGB": color,
            }
        )
    return runs


def extract_text_shape(shape, idx):
    tf = shape.text_frame
    paragraphs = []
    for p in tf.paragraphs:
        paragraphs.append(
            {
                "text": p.text,
                "alignment": str(p.alignment) if p.alignment is not None else None,
                "runs": text_runs(p),
            }
        )

    return {
        "shapeIndex": idx,
        "name": shape.name,
        "type": "text",
        "leftCm": emu_to_cm(shape.left),
        "topCm": emu_to_cm(shape.top),
        "widthCm": emu_to_cm(shape.width),
        "heightCm": emu_to_cm(shape.height),
        "text": shape.text,
        "paragraphs": paragraphs,
    }


def extract_geometry_shape(shape, idx):
    return {
        "shapeIndex": idx,
        "name": shape.name,
        "type": "geometry",
        "autoShapeType": str(shape.auto_shape_type) if hasattr(shape, "auto_shape_type") else None,
        "leftCm": emu_to_cm(shape.left),
        "topCm": emu_to_cm(shape.top),
        "widthCm": emu_to_cm(shape.width),
        "heightCm": emu_to_cm(shape.height),
        "fillRGB": color_to_rgb(shape),
    }


def extract_slide(slide, slide_no):
    result = {
        "slideNo": slide_no,
        "textShapes": [],
        "geometryShapes": [],
    }

    for idx, shape in enumerate(slide.shapes):
        if getattr(shape, "has_text_frame", False):
            if shape.text and shape.text.strip():
                result["textShapes"].append(extract_text_shape(shape, idx))
            continue

        if getattr(shape, "has_chart", False):
            continue

        if hasattr(shape, "auto_shape_type"):
            result["geometryShapes"].append(extract_geometry_shape(shape, idx))

    return result


def extract_demo_config(demo_path):
    prs = Presentation(str(demo_path))
    data = {
        "source": str(demo_path),
        "slideWidthCm": emu_to_cm(prs.slide_width),
        "slideHeightCm": emu_to_cm(prs.slide_height),
        "slides": [],
    }

    for i, slide in enumerate(prs.slides, start=1):
        data["slides"].append(extract_slide(slide, i))

    return data


def main():
    parser = argparse.ArgumentParser(description="提取demo PPT版式配置")
    parser.add_argument("demo", help="demo pptx路径")
    parser.add_argument("--output", "-o", default="demo_layout_config.json", help="输出json路径")
    args = parser.parse_args()

    demo_path = Path(args.demo)
    out_path = Path(args.output)

    config = extract_demo_config(demo_path)
    out_path.write_text(json.dumps(config, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"✅ 已输出配置: {out_path}")


if __name__ == "__main__":
    main()
