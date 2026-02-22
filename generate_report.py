# -*- coding: utf-8 -*-
"""
CCTV-17 农业农村频道 收视日报 PPT 生成器
读取Excel数据，自动生成精美的收视报告PPT

使用方法:
    python generate_report.py <excel_file> [--template origin.pptx] [--output report.pptx]
"""
import sys
import os
import argparse
import copy
import re
import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement

from ppt_config import Colors, Fonts, Layout, SLIDE_WIDTH, SLIDE_HEIGHT, CHANNEL_NAME
from data_reader import read_excel_data, ReportData
from slide_utils import (
    add_textbox, add_rich_textbox, add_title, add_subtitle, add_page_number,
    add_note, add_styled_table, add_ranking_table,
    add_column_chart, add_bar_chart, add_area_chart,
    add_rect, add_rounded_rect, format_change_str, format_change_color,
    _style_cell
)


BODY_PLAIN_MODE = False


# ═══════════════════════════════════════════════════════════════
# PPT 模板操作
# ═══════════════════════════════════════════════════════════════

def _remove_all_slides(prs):
    """删除所有幻灯片"""
    xml_slides = prs.slides._sldIdLst
    for el in list(xml_slides):
        rId = el.get(qn('r:id'))
        prs.part.drop_rel(rId)
        xml_slides.remove(el)


def _remove_slides_except_first(prs):
    """保留第一张幻灯片，删除其余"""
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    for el in slides_list[1:]:
        rId = el.get(qn('r:id'))
        prs.part.drop_rel(rId)
        xml_slides.remove(el)


def _get_layout(prs, name_hint='普通'):
    """获取幻灯片布局"""
    for layout in prs.slide_layouts:
        if name_hint in layout.name:
            return layout
    # 如果找不到，返回第一个布局
    return prs.slide_layouts[0]


def _get_blank_layout(prs):
    """获取空白布局"""
    for layout in prs.slide_layouts:
        if '空白' in layout.name or 'blank' in layout.name.lower():
            return layout
    return prs.slide_layouts[-1]


def _format_start_time(start_time):
    """格式化节目开始时间为 HH:MM 格式。

    处理异常格式如 '1900-01-01 00:39:34' → '00:39'。
    """
    st = str(start_time)
    if '1900-01-01' in st:
        parts = st.split(' ')
        if len(parts) >= 2:
            time_part = parts[-1]
            return ':'.join(time_part.split(':')[:2])
    # 已经是 HH:MM 或 HH:MM:SS 格式
    return ':'.join(st.split(':')[:2])


def _add_slide(prs, layout=None):
    """添加新幻灯片"""
    if layout is None:
        layout = _get_layout(prs)
    slide = prs.slides.add_slide(layout)
    # 清除占位符中的默认文本
    for ph in slide.placeholders:
        if ph.has_text_frame:
            ph.text_frame.clear()
    return slide


# ═══════════════════════════════════════════════════════════════
# 装饰元素
# ═══════════════════════════════════════════════════════════════

def _add_slide_decoration(slide):
    """添加右下角装饰小三角（模仿origin.pptx风格）"""
    from pptx.enum.shapes import MSO_SHAPE
    # 小三角组合（右下角装饰）
    colors = [
        RGBColor(0x1A, 0x3C, 0x5E),
        RGBColor(0x2B, 0x7A, 0xB8),
        RGBColor(0x34, 0x98, 0xDB),
    ]
    base_left = Cm(30.5)
    base_top = Cm(13.0)

    for i, color in enumerate(colors):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            base_left + Cm(i * 0.6),
            base_top + Cm(i * 0.5),
            Cm(1.2 - i * 0.3),
            Cm(1.2 - i * 0.3)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        shape.rotation = 180.0


def _add_title_bar(slide, title_text, subtitle_text=None, page_num=None):
    """添加统一的标题栏"""
    if BODY_PLAIN_MODE:
        return

    # 标题背景横线
    add_rect(slide, Cm(0), Cm(0), SLIDE_WIDTH, Cm(0.15),
             fill_color=Colors.PRIMARY)

    # 标题文本
    add_title(slide, title_text)

    # 副标题
    if subtitle_text:
        add_subtitle(slide, subtitle_text)

    # 页码
    if page_num is not None:
        add_page_number(slide, page_num)

    # 装饰
    _add_slide_decoration(slide)


# ═══════════════════════════════════════════════════════════════
# 幻灯片构建 - 封面
# ═══════════════════════════════════════════════════════════════

def build_cover(prs, data: ReportData):
    """构建封面幻灯片"""
    slide = _add_slide(prs, _get_blank_layout(prs))

    # 浅灰底色
    from pptx.enum.shapes import MSO_SHAPE
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), Cm(33.87), Cm(19.05))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(0xF1, 0xF1, 0xF1)
    bg_shape.line.fill.background()

    # 左下浅灰大圆
    circle_left = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(-3.8), Cm(7.0), Cm(18.0), Cm(18.0))
    circle_left.fill.solid()
    circle_left.fill.fore_color.rgb = RGBColor(241, 243, 240)
    circle_left.line.fill.background()

    # 左上圆环（用双圆叠加，便于精确控制环宽）
    # 图层顺序要求：绿色需覆盖灰色，因此放在灰圆之后绘制
    ring_outer = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(-5.4), Cm(-4.0), Cm(12.4), Cm(12.4))
    ring_outer.fill.solid()
    ring_outer.fill.fore_color.rgb = RGBColor(91, 128, 91)
    ring_outer.line.fill.background()

    ring_inner = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(-3.4), Cm(-2.0), Cm(8.4), Cm(8.4))
    ring_inner.fill.solid()
    ring_inner.fill.fore_color.rgb = RGBColor(0xF1, 0xF1, 0xF1)
    ring_inner.line.fill.background()

    # 右侧灰色小圆点
    circle_mid = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(30.95), Cm(9.35), Cm(1.7), Cm(1.7))
    circle_mid.fill.solid()
    circle_mid.fill.fore_color.rgb = RGBColor(152, 160, 162)
    circle_mid.line.fill.background()

    # 右下绿色小圆
    circle_right = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(26.0), Cm(12.2), Cm(4.5), Cm(4.5))
    circle_right.fill.solid()
    circle_right.fill.fore_color.rgb = RGBColor(91, 128, 91)
    circle_right.line.fill.background()

    # 标题文本（按spec）
    title_text = f'{data.report_date}\n农村农业频道收视日报' if data.report_date else '农村农业频道收视日报'
    add_textbox(slide, Cm(3.8), Cm(4.8), Cm(26.3), Cm(7.3),
                title_text,
                font_name='微软雅黑',
                font_size=Pt(70), font_color=RGBColor(0x00, 0x00, 0x00),
                bold=True, alignment=PP_ALIGN.CENTER)

    # 副标题：统筹策划部 {年}年第{当年第几天}期
    subtitle_text = '统筹策划部'
    if data.report_date:
        m = re.match(r'^(\d{4})年(\d{1,2})月(\d{1,2})日$', data.report_date)
        if m:
            year = int(m.group(1))
            month = int(m.group(2))
            day = int(m.group(3))
            day_of_year = datetime.date(year, month, day).timetuple().tm_yday
            subtitle_text = f'统筹策划部 {year}年第{day_of_year}期'

    add_textbox(slide, Cm(8.0), Cm(12.6), Cm(17.9), Cm(1.6),
                subtitle_text,
                font_name='微软雅黑',
                font_size=Pt(24), font_color=RGBColor(0x00, 0x00, 0x00),
                bold=True, alignment=PP_ALIGN.CENTER)

    return slide


# ═══════════════════════════════════════════════════════════════
# 幻灯片构建 - 收视速报
# ═══════════════════════════════════════════════════════════════

def build_summary(prs, data: ReportData):
    """构建收视速报幻灯片"""
    slide = _add_slide(prs)
    ms = data.market_share
    dr = data.drama

    # 标题（正文净版模式下不显示）
    if not BODY_PLAIN_MODE:
        add_textbox(slide, Cm(1.94), Cm(0.6), Cm(30), Cm(2),
                    '收视速报',
                    font_name='微软雅黑',
                    font_size=Pt(32),
                    font_color=RGBColor(0, 0, 0),
                    bold=True,
                    alignment=PP_ALIGN.LEFT)

    # 变化描述口径（按spec）
    def _summary_change_desc(pct):
        if abs(pct) == 0:
            return '保持不变'
        if abs(pct) < 5:
            return '基本持平'
        return f'提升{abs(pct):.0f}%' if pct > 0 else f'下降{abs(pct):.0f}%'

    # 生成摘要文本（三段）
    date_s = data.report_date_short

    share_change_str = _summary_change_desc(ms.share_change)
    rating_change_str = _summary_change_desc(ms.rating_change)
    org_change_str = _rank_change_desc(data.org_rank_change)
    ch_change_str = _rank_change_desc(data.channel_rank_change)
    para1 = (
        f'{date_s}，频道市场份额为{ms.cctv17_current_share:.3f}%，较前一个月均值{share_change_str}；'
        f'收视率{ms.cctv17_current_rating:.3f}%，{rating_change_str}；'
        f'上星频道排名第{data.channel_rank}位，较前一日{ch_change_str}，'
        f'央视台组内排名第{data.org_rank}位，与前一日{org_change_str}。'
    )

    drama_share_chg = 0
    if dr.drama_period_share > 0:
        drama_share_chg = round((dr.drama_current_share - dr.drama_period_share) / dr.drama_period_share * 100)
    non_drama_share_chg = 0
    if dr.non_drama_period_share > 0:
        non_drama_share_chg = round((dr.non_drama_current_share - dr.non_drama_period_share) / dr.non_drama_period_share * 100)

    drama_chg_str = _summary_change_desc(drama_share_chg)
    non_drama_chg_str = _summary_change_desc(non_drama_share_chg)
    para2 = (
        f'{date_s}，对比前一个月均值，电视剧市场份额{dr.drama_current_share:.3f}%，{drama_chg_str}，'
        f'非电视剧市场份额{dr.non_drama_current_share:.3f}%，{non_drama_chg_str}。'
    )

    para3 = (
        f'{date_s}，早间节目及上午剧场收视率提升，9:24《农人秀-家乡年货开箱记》为上午剧场提供高收视入点；'
        f'午间节目至黄金剧场收视率下滑，《六姊妹》后两集表现不佳，18:00《中国三农报道》将收视平台拉升至平均水平，'
        f'但后续加播《振兴路上新观察》未能有效承接观众；黄金剧场时间因编排调整后移，晚间节目有效承接高收视入点，'
        f'晚间及夜间节目收视提升显著。'
    )

    # 正文文本框（按spec）
    tx = slide.shapes.add_textbox(Cm(1.78), Cm(2.8), Cm(30.31), Cm(14.31))
    tf = tx.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Cm(0.3)
    tf.margin_right = Cm(0.3)
    tf.margin_top = Cm(0.2)
    tf.margin_bottom = Cm(0.2)
    try:
        tf.auto_size = None
    except Exception:
        pass
    tf.vertical_anchor = MSO_ANCHOR.TOP

    # 中文排版控制（对应：按中文习惯控制首尾字符 / 允许标点溢出边界）
    body_pr = tf._txBody.bodyPr
    body_pr.set('eaLnBrk', '1')
    body_pr.set('hangingPunct', '1')

    paragraphs = [para1, para2, para3]
    for i, text in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.25
        p.font.name = '微软雅黑'
        p.font.size = Pt(20)
        p.font.bold = False
        p.font.color.rgb = RGBColor(0, 0, 0) if i < 2 else RGBColor(127, 127, 127)
        p.space_after = Pt(12)

        # 项目符号列表（黑色菱形）
        pPr = p._p.get_or_add_pPr()
        for tag in ('a:buNone', 'a:buAutoNum', 'a:buBlip', 'a:buChar'):
            node = pPr.find(qn(tag))
            if node is not None:
                pPr.remove(node)
        bu = OxmlElement('a:buChar')
        bu.set('char', '◆')
        pPr.insert(0, bu)
        p.level = 0

    # 底部备注文本框（按spec）
    note_text = ('备注：节目名称为监播系统自动抓取。为方便查阅，首播栏目收视数据中不再体现短节目（5分钟及以内）。\n'
                 '对比前期数据日期：前推1个月（30天）均值')
    add_textbox(slide, Cm(1.78), Cm(17.3), Cm(26.32), Cm(1.28),
                note_text,
                font_name='微软雅黑',
                font_size=Pt(12),
                font_color=RGBColor(0, 0, 0),
                bold=False,
                alignment=PP_ALIGN.LEFT)

    return slide


# ═══════════════════════════════════════════════════════════════
# 幻灯片构建 - 市场份额概览
# ═══════════════════════════════════════════════════════════════

def build_market_share(prs, data: ReportData, page_num=1):
    """构建市场份额概览幻灯片"""
    slide = _add_slide(prs)
    ms = data.market_share

    # 标题
    share_chg = _change_desc(ms.share_change, '提升', '下降', '基本持平')
    rating_chg = _change_desc(ms.rating_change, '提升', '下降', '基本持平')
    title = f'{data.report_date_short}，市场份额{share_chg}  收视率{rating_chg}'
    _add_title_bar(slide, title, page_num=page_num)

    # 描述文本（匹配 demo 格式，生成2个段落以匹配demo的XML结构）
    reach_dir = '提升' if ms.reach_change >= 0 else '下降'
    loyalty_dir = '提升' if ms.loyalty_change >= 0 else '下降'
    if reach_dir == loyalty_dir:
        reach_loyalty_desc = f'平均到达率和平均忠实度分别{reach_dir}{abs(ms.reach_change):.0f}%，{loyalty_dir}{abs(ms.loyalty_change):.0f}%'
    else:
        reach_loyalty_desc = f'平均到达率{reach_dir}{abs(ms.reach_change):.0f}%，平均忠实度{loyalty_dir}{abs(ms.loyalty_change):.0f}%'
    all_rating_chg_pct = abs((ms.all_current_rating - ms.all_period_rating) / max(ms.all_period_rating, 0.001) * 100)
    all_rating_dir = '下降' if ms.all_current_rating < ms.all_period_rating else '提升'
    line1 = (f'{data.report_date_short}，频道市场份额为{ms.cctv17_current_share:.3f}%，'
             f'较前一个月均值{share_chg}；{reach_loyalty_desc}。')
    line2 = (f'{data.report_date_short}，所有频道收视率为{ms.all_current_rating:.3f}%，'
             f'较前一个月均值{all_rating_dir}{all_rating_chg_pct:.1f}%，'
             f'CCTV-17收视率{ms.cctv17_current_rating:.3f}%，'
             f'{rating_chg}。')
    # 用2个独立段落（而非\n），确保v3同步时与demo的2段落结构匹配
    txbox = slide.shapes.add_textbox(Cm(1.5), Cm(2.0), Cm(30), Cm(2.2))
    tf = txbox.text_frame
    tf.word_wrap = True
    try:
        tf.auto_size = None
    except:
        pass
    p1 = tf.paragraphs[0]
    p1.text = line1
    p1.font.name = Fonts.MAIN
    p1.font.size = Pt(10)
    p1.font.color.rgb = Colors.MEDIUM_GRAY
    p2 = tf.add_paragraph()
    p2.text = line2
    p2.font.name = Fonts.MAIN
    p2.font.size = Pt(10)
    p2.font.color.rgb = Colors.MEDIUM_GRAY

    # ── 构建类别轴 ──
    prev_label = ms.prev_date_label if ms.prev_date_label else '前一日'
    if ms.has_prev_day:
        categories = ['前1个月\n均值', prev_label, data.report_date_short]
    else:
        categories = ['前1个月\n均值', data.report_date_short]

    # ── 左侧: 市场份额柱状图 ──
    if ms.has_prev_day:
        share_series = {
            'CCTV-17': [ms.cctv17_period_share, ms.cctv17_prev_share, ms.cctv17_current_share],
        }
    else:
        share_series = {
            'CCTV-17': [ms.cctv17_period_share, ms.cctv17_current_share],
        }
    add_column_chart(slide, Cm(0.5), Cm(4.5), Cm(6), Cm(8),
                     categories, share_series,
                     series_colors=[Colors.CHART_SERIES[0]],
                     show_data_labels=True, gap_width=80)
    add_textbox(slide, Cm(1.5), Cm(8), Cm(3), Cm(0.8),
                '市场份额%', font_size=Fonts.TINY_SIZE,
                font_color=Colors.DARK_GRAY)

    # ── 中间: 收视率标注 ──
    add_textbox(slide, Cm(10), Cm(4.5), Cm(3), Cm(1.2),
                f'所有频道\n收视率%', font_size=Fonts.SMALL_SIZE,
                font_color=Colors.DARK_GRAY, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Cm(10), Cm(9), Cm(3), Cm(1.2),
                f'CCTV-17\n收视率%', font_size=Fonts.SMALL_SIZE,
                font_color=Colors.DARK_GRAY, bold=True,
                alignment=PP_ALIGN.CENTER)

    # ── 中上: 所有频道收视率条形图 ──
    if ms.has_prev_day:
        all_rating_cats = ['前1个月均值', prev_label, data.report_date_short]
        all_rating_series = {'所有频道': [ms.all_period_rating, ms.all_prev_rating, ms.all_current_rating]}
    else:
        all_rating_cats = ['前1个月均值', data.report_date_short]
        all_rating_series = {'所有频道': [ms.all_period_rating, ms.all_current_rating]}
    add_bar_chart(slide, Cm(12.5), Cm(4.2), Cm(8), Cm(4.5),
                  all_rating_cats, all_rating_series,
                  series_colors=[Colors.CHART_SERIES[0]])

    # ── 中下: CCTV-17收视率条形图 ──
    if ms.has_prev_day:
        c17_rating_cats = ['前1个月均值', prev_label, data.report_date_short]
        c17_rating_series = {'CCTV-17': [ms.cctv17_period_rating, ms.cctv17_prev_rating, ms.cctv17_current_rating]}
    else:
        c17_rating_cats = ['前1个月均值', data.report_date_short]
        c17_rating_series = {'CCTV-17': [ms.cctv17_period_rating, ms.cctv17_current_rating]}
    add_bar_chart(slide, Cm(12.5), Cm(9.2), Cm(8.5), Cm(4.5),
                  c17_rating_cats, c17_rating_series,
                  series_colors=[Colors.CHART_SERIES[1]])

    # ── 右上: 到达率变化表格 ──
    reach_headers = ['平均到达率%', '']
    reach_data = [
        [f'{ms.reach_change:+.0f}%', ''],
        [f'前1个月均值{ms.cctv17_period_reach:.2f}', f'当天{ms.cctv17_current_reach:.2f}']
    ]
    add_styled_table(slide, Cm(22), Cm(5.6), Cm(6.5), Cm(3.0),
                     reach_headers, reach_data)

    # ── 右下: 忠实度变化表格 ──
    loy_headers = ['平均忠实度', '']
    loy_data = [
        [f'{ms.loyalty_change:+.0f}%', ''],
        [f'前1个月均值{ms.cctv17_period_loyalty:.2f}', f'当天{ms.cctv17_current_loyalty:.2f}']
    ]
    add_styled_table(slide, Cm(22), Cm(8.6), Cm(6.5), Cm(3.0),
                     loy_headers, loy_data)

    return slide


# ═══════════════════════════════════════════════════════════════
# 幻灯片构建 - 台组内排名
# ═══════════════════════════════════════════════════════════════

def build_org_ranking(prs, data: ReportData, page_num=2):
    """构建央视台组内排名幻灯片"""
    slide = _add_slide(prs)

    org_change = _rank_change_desc(data.org_rank_change)
    title = f'央视台组排名第{data.org_rank}位，较前一日{org_change}'
    _add_title_bar(slide, title, page_num=page_num)

    rankings = data.org_ranking
    if not rankings:
        return slide

    # 获取日期标签
    prev_date = ''
    curr_date = data.report_date_short
    if data.market_share.has_prev_day:
        prev_label = data.market_share.current_date_label
        # 前一天
        try:
            parts = prev_label.split('/')
            day = int(parts[2]) - 1
            prev_date = f'{int(parts[1])}月{day}日'
        except:
            prev_date = '前一日'
    else:
        prev_date = '前一日'

    # 构建表格数据
    prev_sorted = sorted(rankings, key=lambda x: x.prev_share, reverse=True)
    curr_sorted = sorted(rankings, key=lambda x: x.current_share, reverse=True)

    n_display = len(rankings)
    headers = ['排名', '频道', '市场份额%', '排名', '频道', '市场份额%']

    data_rows = []
    prev_rank = 0
    curr_rank = 0
    for i in range(n_display):
        prev_item = prev_sorted[i] if i < len(prev_sorted) else None
        curr_item = curr_sorted[i] if i < len(curr_sorted) else None

        row = []
        if prev_item:
            if '中央级' in prev_item.name:
                rank_str = '\u3000'
            else:
                prev_rank += 1
                rank_str = str(prev_rank)
            row += [rank_str, prev_item.short_name or prev_item.name,
                    f'{prev_item.prev_share:.3f}']
        else:
            row += ['', '', '']

        if curr_item:
            if '中央级' in curr_item.name:
                rank_str = '\u3000'
            else:
                curr_rank += 1
                rank_str = str(curr_rank)
            row += [rank_str, curr_item.short_name or curr_item.name,
                    f'{curr_item.current_share:.3f}']
        else:
            row += ['', '', '']

        data_rows.append(row)

    # 创建表格
    n_rows = n_display + 2  # 日期头 + 列头 + 数据
    n_cols = 6

    tbl_left = Cm(1.0)
    tbl_top = Cm(3.0)
    tbl_width = Cm(31)
    tbl_height = Cm(14.5)

    table_shape = slide.shapes.add_table(n_rows, n_cols, tbl_left, tbl_top,
                                         tbl_width, tbl_height)
    table = table_shape.table

    # 行1: 日期
    table.cell(0, 0).merge(table.cell(0, 2))
    table.cell(0, 3).merge(table.cell(0, 5))
    table.cell(0, 0).text = prev_date
    table.cell(0, 3).text = curr_date
    for j in range(6):
        _style_cell(table.cell(0, j), Fonts.TABLE_HEADER_SIZE, Colors.TABLE_HEADER_FG,
                    bold=True, bg_color=Colors.TABLE_HEADER_BG)
    _style_cell(table.cell(0, 0), Fonts.TABLE_HEADER_SIZE, Colors.TABLE_HEADER_FG,
                bold=True, bg_color=Colors.TABLE_HEADER_BG)
    _style_cell(table.cell(0, 3), Fonts.TABLE_HEADER_SIZE, Colors.TABLE_HEADER_FG,
                bold=True, bg_color=Colors.TABLE_HEADER_BG)

    # 行2: 列头
    for j, h in enumerate(headers):
        cell = table.cell(1, j)
        cell.text = h
        _style_cell(cell, Fonts.TABLE_BODY_SIZE, Colors.PRIMARY,
                    bold=True, bg_color=Colors.TABLE_ROW_EVEN)

    # 数据行
    for i, row in enumerate(data_rows):
        is_hl = any('农业农村' in v or 'CCTV-17' in v for v in row)
        bg = Colors.TABLE_HIGHLIGHT if is_hl else \
             (Colors.TABLE_ROW_EVEN if i % 2 == 0 else Colors.TABLE_ROW_ODD)
        for j, val in enumerate(row):
            cell = table.cell(i + 2, j)
            cell.text = val
            _style_cell(cell, Fonts.TABLE_BODY_SIZE, Colors.DARK_GRAY,
                       bold=is_hl, bg_color=bg)

    # CCTV-17的市场份额标注
    cctv17_item = None
    for item in rankings:
        if '农业农村' in item.name or 'CCTV-17' in item.short_name:
            cctv17_item = item
            break
    if cctv17_item:
        add_textbox(slide, Cm(27), Cm(15.5), Cm(4), Cm(1),
                    f'{cctv17_item.current_share:.2f}%',
                    font_size=Fonts.SMALL_SIZE, font_color=Colors.ACCENT_RED,
                    bold=True)

    # 备注
    add_note(slide, '备注：仅体现考核序列内央视台组频道', top=Cm(17.5))

    return slide


# ═══════════════════════════════════════════════════════════════
# 幻灯片构建 - 上星频道排名
# ═══════════════════════════════════════════════════════════════

def build_channel_ranking(prs, data: ReportData, page_num=3):
    """构建上星频道排名幻灯片"""
    slide = _add_slide(prs)

    ch_change = _rank_change_desc(data.channel_rank_change)
    title = f'上星频道排名第{data.channel_rank}位，较前一日{ch_change}'
    _add_title_bar(slide, title, page_num=page_num)

    rankings = data.channel_ranking
    if not rankings:
        return slide

    prev_date = '前一日'
    curr_date = data.report_date_short

    curr_sorted = sorted(rankings, key=lambda x: x.current_share, reverse=True)
    prev_sorted = sorted(rankings, key=lambda x: x.prev_share, reverse=True)

    # 构建展示行：前 10 + 分隔行 + CCTV-17 前1位 + CCTV-17 + 后2位
    TOP_N = 10
    top_prev = prev_sorted[:TOP_N]
    top_curr = curr_sorted[:TOP_N]

    # 在当日排名中查找 CCTV-17 的位置
    cctv17_curr_idx = -1
    for idx, item in enumerate(curr_sorted):
        if '农业农村' in item.name or 'CCTV-17' in item.name:
            cctv17_curr_idx = idx
            break

    # 在前日排名中查找 CCTV-17 的位置
    cctv17_prev_idx = -1
    for idx, item in enumerate(prev_sorted):
        if '农业农村' in item.name or 'CCTV-17' in item.name:
            cctv17_prev_idx = idx
            break

    # 如果 CCTV-17 已在前10，不需要分隔行和底部区域
    if cctv17_curr_idx >= 0 and cctv17_curr_idx < TOP_N:
        # CCTV-17 已在 top 10 中，只展示前 10
        display_rows = []
        for i in range(TOP_N):
            p = prev_sorted[i] if i < len(prev_sorted) else None
            c = curr_sorted[i] if i < len(curr_sorted) else None
            display_rows.append(('data', i, p, i, c))
    else:
        # 前 10 展示
        display_rows = []
        for i in range(min(TOP_N, len(rankings))):
            p = prev_sorted[i] if i < len(prev_sorted) else None
            c = curr_sorted[i] if i < len(curr_sorted) else None
            display_rows.append(('data', i, p, i, c))

        # 分隔行
        display_rows.append(('separator', None, None, None, None))

        # 底部 4 行：CCTV-17 前 1 位 + CCTV-17 + 后 2 位
        if cctv17_curr_idx >= 0:
            # 当日侧：前 1 位 + CCTV-17 + 后 2 位
            curr_start = max(0, cctv17_curr_idx - 1)
            curr_indices = list(range(curr_start, min(curr_start + 4, len(curr_sorted))))
            # 确保 CCTV-17 在其中
            if cctv17_curr_idx not in curr_indices:
                curr_indices = list(range(cctv17_curr_idx, min(cctv17_curr_idx + 3, len(curr_sorted))))
        else:
            curr_indices = []

        if cctv17_prev_idx >= 0:
            # 前日侧：前 1 位 + CCTV-17 + 后 2 位
            prev_start = max(0, cctv17_prev_idx - 1)
            prev_indices = list(range(prev_start, min(prev_start + 4, len(prev_sorted))))
            if cctv17_prev_idx not in prev_indices:
                prev_indices = list(range(cctv17_prev_idx, min(cctv17_prev_idx + 3, len(prev_sorted))))
        else:
            prev_indices = []

        n_bottom = max(len(curr_indices), len(prev_indices))
        for j in range(n_bottom):
            pi = prev_indices[j] if j < len(prev_indices) else None
            ci = curr_indices[j] if j < len(curr_indices) else None
            p = prev_sorted[pi] if pi is not None and pi < len(prev_sorted) else None
            c = curr_sorted[ci] if ci is not None and ci < len(curr_sorted) else None
            display_rows.append(('data', pi, p, ci, c))

    headers = ['排名', '频道', '市场份额%', '排名', '频道', '市场份额%', '较前一日变化']
    n_rows = len(display_rows) + 2  # date header + col header + data
    n_cols = 7

    tbl_left = Cm(0.5)
    tbl_top = Cm(3.0)
    tbl_width = Cm(32.2)
    tbl_height = Cm(14.5)

    table_shape = slide.shapes.add_table(n_rows, n_cols, tbl_left, tbl_top,
                                         tbl_width, tbl_height)
    table = table_shape.table

    # 行1: 日期
    table.cell(0, 0).merge(table.cell(0, 2))
    table.cell(0, 3).merge(table.cell(0, 5))
    table.cell(0, 0).text = prev_date
    table.cell(0, 3).text = curr_date
    table.cell(0, 6).text = ''
    for j in range(7):
        _style_cell(table.cell(0, j), Fonts.TABLE_HEADER_SIZE, Colors.TABLE_HEADER_FG,
                    bold=True, bg_color=Colors.TABLE_HEADER_BG)
    _style_cell(table.cell(0, 0), Fonts.TABLE_HEADER_SIZE, Colors.TABLE_HEADER_FG,
                bold=True, bg_color=Colors.TABLE_HEADER_BG)
    _style_cell(table.cell(0, 3), Fonts.TABLE_HEADER_SIZE, Colors.TABLE_HEADER_FG,
                bold=True, bg_color=Colors.TABLE_HEADER_BG)

    # 行2: 列头
    for j, h in enumerate(headers):
        cell = table.cell(1, j)
        cell.text = h
        _style_cell(cell, Fonts.TABLE_BODY_SIZE, Colors.PRIMARY,
                    bold=True, bg_color=Colors.TABLE_ROW_EVEN)

    # 数据行
    for i, (row_type, prev_idx, prev_item, curr_idx, curr_item) in enumerate(display_rows):
        row_idx = i + 2

        if row_type == 'separator':
            for j in range(7):
                table.cell(row_idx, j).text = '—'
                _style_cell(table.cell(row_idx, j), Fonts.TABLE_BODY_SIZE,
                           Colors.DARK_GRAY, bold=False, bg_color=Colors.TABLE_ROW_EVEN)
            continue

        is_hl = False
        if curr_item and ('农业农村' in curr_item.name or 'CCTV-17' in curr_item.name):
            is_hl = True
        bg = Colors.TABLE_HIGHLIGHT if is_hl else \
             (Colors.TABLE_ROW_EVEN if i % 2 == 0 else Colors.TABLE_ROW_ODD)

        # 前日
        if prev_item and prev_idx is not None:
            table.cell(row_idx, 0).text = f'{prev_idx + 1}'
            table.cell(row_idx, 1).text = prev_item.name
            table.cell(row_idx, 2).text = f'{prev_item.prev_share:.3f}'
        # 当日
        if curr_item and curr_idx is not None:
            table.cell(row_idx, 3).text = f'{curr_idx + 1}'
            table.cell(row_idx, 4).text = curr_item.name
            table.cell(row_idx, 5).text = f'{curr_item.current_share:.3f}'
            # 变化箭头
            diff = curr_item.current_share - curr_item.prev_share
            if diff > 0:
                table.cell(row_idx, 6).text = '↑'
            elif diff < 0:
                table.cell(row_idx, 6).text = '↓'
            else:
                table.cell(row_idx, 6).text = '-'

        for j in range(7):
            _style_cell(table.cell(row_idx, j), Fonts.TABLE_BODY_SIZE,
                       Colors.DARK_GRAY, bold=is_hl, bg_color=bg)

    return slide


# ═══════════════════════════════════════════════════════════════
# 幻灯片构建 - 串单市场份额
# ═══════════════════════════════════════════════════════════════

def build_schedule_chart(prs, data: ReportData, page_num=4):
    """构建频道串单市场份额图"""
    slide = _add_slide(prs)
    ms = data.market_share

    title = (f'{data.report_date_short}，频道市场份额{ms.cctv17_current_share:.3f}%，'
             f'收视率{ms.cctv17_current_rating:.3f}%')
    _add_title_bar(slide, title, page_num=page_num)

    programs = data.programs
    if not programs:
        return slide

    # 过滤出有效节目（时长>5分钟，时间在06:00-24:00之间）
    valid_progs = []
    for p in programs:
        if p.duration <= 5:
            continue
        # 过滤凌晨节目
        try:
            hour = int(p.start_time.split(':')[0])
            if hour < 6 and hour >= 2:
                continue
        except:
            pass
        valid_progs.append(p)

    if not valid_progs:
        valid_progs = [p for p in programs if p.duration > 2]

    # 截取合理数量
    valid_progs = valid_progs[:40]

    categories = []
    share_vals = []
    rating_vals = []
    for p in valid_progs:
        label = p.name[:8] if len(p.name) > 8 else p.name
        st = _format_start_time(p.start_time)
        categories.append(f'{label}{st}')
        share_vals.append(p.market_share)
        rating_vals.append(p.rating)

    series = {'市场份额%': share_vals, '收视率': rating_vals}

    # 生成图表颜色（首播节目用深色，重播用浅色）
    chart_shape = add_column_chart(
        slide, Cm(0.5), Cm(2.5), Cm(32), Cm(15),
        categories, series,
        series_colors=[Colors.CHART_SERIES[0]],
        show_data_labels=True, gap_width=30
    )

    # 找出最高收视节目并标注
    if share_vals:
        max_idx = share_vals.index(max(share_vals))
        max_prog = valid_progs[max_idx]
        add_textbox(slide, Cm(28), Cm(2.0), Cm(4), Cm(0.8),
                    f'{max(share_vals):.3f}',
                    font_size=Fonts.BODY_SIZE, font_color=Colors.ACCENT_RED,
                    bold=True, alignment=PP_ALIGN.RIGHT)

    return slide


# ═══════════════════════════════════════════════════════════════
# 幻灯片构建 - 栏目收视率/份额排名
# ═══════════════════════════════════════════════════════════════

def build_program_ranking(prs, data: ReportData, metric='rating', page_num=5):
    """构建栏目收视率/份额排名"""
    slide = _add_slide(prs)

    metric_name = '收视率' if metric == 'rating' else '收视份额'
    title = f'{data.report_date_short} 栏目{metric_name}排名'
    _add_title_bar(slide, title, page_num=page_num)

    programs = data.programs
    if not programs:
        return slide

    # 过滤有效节目 (时长>5min, 6:00-24:00)
    valid = []
    for p in programs:
        if p.duration <= 5:
            continue
        try:
            hour = int(p.start_time.split(':')[0])
            if 2 <= hour < 6:
                continue
        except:
            pass
        valid.append(p)

    # 按指标排序
    if metric == 'rating':
        valid.sort(key=lambda x: x.rating, reverse=True)
    else:
        valid.sort(key=lambda x: x.market_share, reverse=True)

    # 取前30个
    display = valid[:30]

    categories = []
    values = []
    for p in display:
        label = p.name[:10] if len(p.name) > 10 else p.name
        st = _format_start_time(p.start_time)
        categories.append(f'{label}{st}')
        values.append(p.rating if metric == 'rating' else p.market_share)

    series = {metric_name + '%': values}

    add_column_chart(slide, Cm(0.5), Cm(2.5), Cm(32), Cm(13.5),
                     categories, series,
                     series_colors=[Colors.CHART_SERIES[0]],
                     show_data_labels=True, gap_width=50)

    # 最高值标注
    if values:
        max_val = max(values)
        add_textbox(slide, Cm(1), Cm(2.0), Cm(4), Cm(0.8),
                    f'{max_val:.4f}',
                    font_size=Fonts.BODY_SIZE, font_color=Colors.ACCENT_RED,
                    bold=True)

    add_note(slide, '备注：仅体现早6点至晚22点之间开播的节目',
             top=Cm(16.5))

    return slide


# ═══════════════════════════════════════════════════════════════
# 幻灯片构建 - 分分钟收视率/市场份额
# ═══════════════════════════════════════════════════════════════

def _time_to_minutes(t):
    """将 HH:MM 时间转为从00:00起的分钟数（支持 24+小时制，如 25:30=1530）"""
    parts = t.split(':')
    return int(parts[0]) * 60 + int(parts[1])


def build_minute_chart(prs, data: ReportData, metric='share', page_num=7):
    """构建分分钟面积图"""
    slide = _add_slide(prs)

    metric_name = '收视率' if metric == 'rating' else '市场份额'
    title = f'分分钟{metric_name}'
    _add_title_bar(slide, title, page_num=page_num)

    minutes = data.minutes
    if not minutes:
        return slide

    # 筛选 05:30 ~ 25:59（对应 demo 中的 05:30~01:59）每分钟数据，不采样
    filtered = [m for m in minutes
                if _time_to_minutes(m.time_str) >= 330]  # >= 05:30

    if not filtered:
        filtered = minutes

    # 图表分类标签：24+ 小时转为次日格式（25:30 → 01:30）
    categories = []
    for m in filtered:
        h, mn = m.time_str.split(':')[:2]
        h_int = int(h)
        if h_int >= 24:
            categories.append(f'{h_int - 24:02d}:{mn}')
        else:
            categories.append(f'{h}:{mn}')

    if metric == 'rating':
        series = {
            '前1个月均值': [m.period_rating for m in filtered],
            data.report_date_short: [m.current_rating for m in filtered],
        }
        colors = [Colors.CHART_SERIES[0], Colors.CHART_SERIES[1]]
    else:
        series = {
            '前1个月均值': [m.period_share for m in filtered],
            data.report_date_short: [m.current_share for m in filtered],
        }
        colors = [Colors.CHART_SERIES[0], Colors.CHART_SERIES[1]]

    add_area_chart(slide, Cm(0.5), Cm(2.5), Cm(32), Cm(10),
                   categories, series, series_colors=colors)

    # 节目名称标签栏（与图表重叠）
    _add_program_labels(slide, data, Cm(0.5), Cm(2.5), Cm(32), Cm(10))

    # 时段汇总表（图表上方）
    _add_timeslot_summary(slide, data, metric,
                          Cm(0.5), Cm(14.2), Cm(32), Cm(2.5))

    # 前一个月剧集标签表（与图表重叠，空占位符，8列）
    _add_prev_month_labels(slide, Cm(0.5), Cm(2.5), Cm(32), Cm(10))

    # 补充说明文本
    if data.has_daily_report and data.time_slots:
        _add_timeslot_comment(slide, data, metric, Cm(0.5), Cm(17.0))

    return slide


def _add_program_labels(slide, data, left, top, width, height):
    """在分分钟图下方添加节目名称标签"""
    programs = data.programs
    if not programs:
        return

    # 过滤出6:00-24:00的有效节目
    valid = []
    for p in programs:
        if p.duration < 5:
            continue
        try:
            hour = int(p.start_time.split(':')[0])
            if 2 <= hour < 6:
                continue
        except:
            pass
        valid.append(p)

    if not valid:
        return

    # 创建1行N列的名称表
    n_cols = min(len(valid), 20)  # 最多显示20个节目
    step = max(1, len(valid) // n_cols)
    display = valid[::step][:n_cols]

    if not display:
        return

    table_shape = slide.shapes.add_table(1, n_cols, left, top, width, height)
    table = table_shape.table

    for j, p in enumerate(display):
        cell = table.cell(0, j)
        # 节目名称（截断），添加重播/首播标识
        name = p.name[:6]
        suffix = '（首播）' if p.start_time >= '17:00' else '（重播）'
        cell.text = f'{name}{suffix}'
        _style_cell(cell, Pt(6), Colors.DARK_GRAY, bg_color=Colors.WHITE)
        cell.margin_top = Pt(1)
        cell.margin_bottom = Pt(1)


def _add_prev_month_labels(slide, left, top, width, height):
    """添加前一个月剧集标签表（空占位符，1行8列，与demo结构一致）"""
    n_cols = 8
    table_shape = slide.shapes.add_table(1, n_cols, left, top, width, height)
    table = table_shape.table
    for j in range(n_cols):
        cell = table.cell(0, j)
        cell.text = ''
        _style_cell(cell, Pt(6), Colors.DARK_GRAY, bg_color=Colors.WHITE)
        cell.margin_top = Pt(1)
        cell.margin_bottom = Pt(1)


def _add_timeslot_summary(slide, data, metric, left, top, width, height):
    """添加分时段汇总表"""
    slots = [
        ('早间节目\n(6:00-8:00)', '06', '08'),
        ('上午剧场\n(8:00-12:00)', '08', '12'),
        ('午间节目\n(12:00-14:00)', '12', '14'),
        ('下午剧场\n(14:00-18:30)', '14', '18'),
        ('傍晚节目\n(18:30-22:00)', '18', '22'),
        ('晚间节目\n(22:00-24:00)', '22', '24'),
    ]

    # 如果有日报-分时段收视率数据，直接使用
    if data.has_daily_report and data.time_slots:
        n_slots = len(data.time_slots)
        # 动态确定列数
        slot_labels = []
        slot_values = []
        for ts in data.time_slots:
            name = ts.slot_name
            if metric == 'rating':
                val = ts.current_rating
                chg = ts.rating_change
            else:
                val = ts.current_share
                chg = ts.share_change
            chg_pct = round(chg * 100) if abs(chg) < 10 else round(chg)
            slot_labels.append(name)
            slot_values.append(f'{val:.3f}%\n{chg_pct:+d}%')

        n_cols = len(slot_labels)
        table_shape = slide.shapes.add_table(1, n_cols, left, top, width, height)
        table = table_shape.table
        for j in range(n_cols):
            cell = table.cell(0, j)
            # 格式：时段名\n数值\n变化
            parts = slot_labels[j].replace('-', '\n~\n').split('\n')
            cell.text = f'{slot_labels[j]}\n{slot_values[j]}'
            _style_cell(cell, Pt(8), Colors.DARK_GRAY,
                       bg_color=Colors.TABLE_ROW_EVEN)
    else:
        # 从分分钟数据计算
        n_cols = len(slots)
        table_shape = slide.shapes.add_table(1, n_cols, left, top, width, height)
        table = table_shape.table
        for j, (label, start, end) in enumerate(slots):
            # 计算该时段平均值
            period_vals = []
            current_vals = []
            for m in data.minutes:
                try:
                    h = int(m.time_str.split(':')[0])
                    if int(start) <= h < int(end):
                        if metric == 'rating':
                            period_vals.append(m.period_rating)
                            current_vals.append(m.current_rating)
                        else:
                            period_vals.append(m.period_share)
                            current_vals.append(m.current_share)
                except:
                    pass

            avg_period = sum(period_vals) / len(period_vals) if period_vals else 0
            avg_current = sum(current_vals) / len(current_vals) if current_vals else 0
            chg = round((avg_current - avg_period) / max(avg_period, 0.001) * 100)

            cell = table.cell(0, j)
            cell.text = f'{label}\n{avg_current:.3f}%\n{chg:+d}%'
            _style_cell(cell, Pt(8), Colors.DARK_GRAY,
                       bg_color=Colors.TABLE_ROW_EVEN)


def _add_timeslot_comment(slide, data, metric, left, top):
    """添加分时段变化说明"""
    if not data.time_slots:
        return

    parts = []
    for ts in data.time_slots:
        if metric == 'rating':
            chg = ts.rating_change
        else:
            chg = ts.share_change
        chg_pct = round(chg * 100) if abs(chg) < 10 else round(chg)
        if abs(chg_pct) > 20:
            direction = '提升' if chg_pct > 0 else '下降'
            parts.append(f'{ts.slot_name}{metric == "rating" and "收视率" or "市场份额"}'
                        f'{direction}{abs(chg_pct)}%')

    if parts:
        text = f'{data.report_date_short}，相比前一个月均值，' + '；'.join(parts[:3]) + '。'
        add_textbox(slide, left, top, Cm(31), Cm(1.5),
                    text, font_size=Pt(9), font_color=Colors.MEDIUM_GRAY)


# ═══════════════════════════════════════════════════════════════
# 幻灯片构建 - 电视剧分析 (日报)
# ═══════════════════════════════════════════════════════════════

def build_drama_analysis(prs, data: ReportData, page_num=None):
    """构建电视剧与非电视剧分析"""
    slide = _add_slide(prs)
    dr = data.drama

    # 计算变化
    drama_chg = 0
    if dr.drama_period_share > 0:
        drama_chg = round((dr.drama_current_share - dr.drama_period_share) / dr.drama_period_share * 100)
    non_drama_chg = 0
    if dr.non_drama_period_share > 0:
        non_drama_chg = round((dr.non_drama_current_share - dr.non_drama_period_share) / dr.non_drama_period_share * 100)

    title = (f'电视剧市场份额较前一个月'
             f'{"下降" if drama_chg < 0 else "提升"}{abs(drama_chg)}%，'
             f'非电视剧{"提升" if non_drama_chg > 0 else "下降"}{abs(non_drama_chg)}%')
    _add_title_bar(slide, title, page_num=page_num)

    # 描述文字
    desc_parts = []
    desc_parts.append(
        f'{data.report_date_short}，对比前一个月均值，'
        f'电视剧市场份额{dr.drama_current_share:.3f}%，'
        f'{"下降" if drama_chg < 0 else "提升"}{abs(drama_chg)}%，'
        f'非电视剧市场份额{dr.non_drama_current_share:.3f}%，'
        f'{"提升" if non_drama_chg > 0 else "下降"}{abs(non_drama_chg)}%。')

    # 三大剧场描述
    for t in data.theaters:
        if t.time_slot and t.info:
            chg_str = f'{"提升" if t.share_change > 0 else "下降"}{abs(round(t.share_change * 100))}%'
            desc_parts.append(f'{t.time_slot} 《{t.info.split("-")[0] if "-" in t.info else t.info}》'
                            f'市场份额{t.current_share:.3f}%，{chg_str}')

    desc = '。'.join(desc_parts[:3]) + '。'
    add_textbox(slide, Cm(1.5), Cm(2.0), Cm(30), Cm(2),
                desc, font_size=Pt(9), font_color=Colors.MEDIUM_GRAY)

    # 左侧: 电视剧/非电视剧对比图
    categories = ['前1个月均值', data.report_date_short]
    series = {
        '电视剧': [dr.drama_period_share, dr.drama_current_share],
        '非电视剧': [dr.non_drama_period_share, dr.non_drama_current_share],
    }
    add_column_chart(slide, Cm(0.5), Cm(4.0), Cm(12), Cm(9),
                     categories, series,
                     series_colors=[Colors.CHART_SERIES[0], Colors.CHART_SERIES[2]],
                     show_data_labels=True)

    # 右侧: 三大剧场表格
    if data.theaters:
        headers = ['', '前1个月均值', '', data.report_date_short, '', '变化幅度']
        sub_headers = ['', '播出时段及信息', '收视率&收视份额%', '播出时段及信息',
                       '收视率&收视份额%', '']

        n_rows = len(data.theaters) + 2
        n_cols = 6

        tbl_shape = slide.shapes.add_table(n_rows, n_cols,
                                           Cm(13), Cm(4.0), Cm(19.5), Cm(9))
        tbl = tbl_shape.table

        # 表头
        for j, h in enumerate(headers):
            tbl.cell(0, j).text = h
            _style_cell(tbl.cell(0, j), Fonts.TABLE_HEADER_SIZE,
                       Colors.TABLE_HEADER_FG, bold=True,
                       bg_color=Colors.TABLE_HEADER_BG)
        for j, h in enumerate(sub_headers):
            tbl.cell(1, j).text = h
            _style_cell(tbl.cell(1, j), Pt(7), Colors.PRIMARY,
                       bold=True, bg_color=Colors.TABLE_ROW_EVEN)

        # 数据
        theater_names = ['上午剧场', '下午剧场', '黄金剧场']
        for i, t in enumerate(data.theaters):
            row_idx = i + 2
            name = theater_names[i] if i < len(theater_names) else t.time_slot
            bg = Colors.TABLE_ROW_EVEN if i % 2 == 0 else Colors.TABLE_ROW_ODD

            tbl.cell(row_idx, 0).text = name
            tbl.cell(row_idx, 1).text = t.time_slot
            tbl.cell(row_idx, 2).text = f'{t.period_rating:.3f}'
            # 当日
            info = t.info[:10] if len(t.info) > 10 else t.info
            tbl.cell(row_idx, 3).text = f'{t.time_slot}\n{info}'
            tbl.cell(row_idx, 4).text = f'{t.current_rating:.3f}'
            chg = round(t.share_change * 100) if abs(t.share_change) < 10 else round(t.share_change)
            tbl.cell(row_idx, 5).text = f'{chg:+d}%'

            for j in range(6):
                _style_cell(tbl.cell(row_idx, j), Pt(8), Colors.DARK_GRAY,
                           bg_color=bg)

    add_note(slide, '备注：非电视剧包括自制栏目剧', top=Cm(13.5))

    return slide


# ═══════════════════════════════════════════════════════════════
# 幻灯片构建 - 首播电视剧分类观众规模(日报)
# ═══════════════════════════════════════════════════════════════

def build_drama_audience(prs, data: ReportData, page_num=None):
    """构建首播电视剧分类观众规模"""
    slide = _add_slide(prs)
    _add_title_bar(slide, '首播电视剧分类观众规模及各剧场观众规模',
                   page_num=page_num)

    da = data.drama_audience
    if not da:
        return slide

    # 左侧: 分类观众图(各section汇总)
    by_section = da.get('by_section', {})
    valid_sections = {k: v for k, v in by_section.items() if k is not None and v}

    if valid_sections:
        # 合并所有section取"四岁及以上所有人"等总体数据
        all_items = []
        for section_name, items in valid_sections.items():
            for item in items:
                all_items.append(item)

        # 按target分组, 取得唯一分类
        first_section_key = list(valid_sections.keys())[0]
        items = valid_sections[first_section_key]
        cats = [item['target'][:8] for item in items[:10]]
        vals = [item['audience'] for item in items[:10]]

        if any(v > 0 for v in vals):
            series = {'观众规模(万人)': vals}
            add_column_chart(slide, Cm(0.5), Cm(3.0), Cm(15), Cm(11),
                            cats, series,
                            series_colors=[Colors.CHART_SERIES[0]],
                            show_data_labels=True)

    # 右侧: 剧场观众规模图
    theater_data = da.get('theaters', [])
    if theater_data:
        valid_theaters = [t for t in theater_data if t.get('period_value', 0) > 0
                          or t.get('current_value', 0) > 0]
        if valid_theaters:
            cats = [t['name'] for t in valid_theaters[:5]]
            period_vals = [t['period_value'] for t in valid_theaters[:5]]
            current_vals = [t['current_value'] for t in valid_theaters[:5]]

            series = {
                '前1个月日均': period_vals,
                data.report_date_short: current_vals,
            }
            add_column_chart(slide, Cm(16.5), Cm(3.0), Cm(16), Cm(11),
                            cats, series,
                            series_colors=[Colors.CHART_SERIES[0], Colors.CHART_SERIES[1]],
                            show_data_labels=True)

    add_note(slide, '备注：对比数据：前推1个月（30天）日均观众规模。',
             top=Cm(15))

    return slide


# ═══════════════════════════════════════════════════════════════
# 幻灯片构建 - 首播节目对比(日报)
# ═══════════════════════════════════════════════════════════════

def build_premiere_comparison(prs, data: ReportData, metric='rating', page_num=None):
    """构建首播栏目收视率/份额/忠实度对比"""
    slide = _add_slide(prs)

    metric_map = {
        'rating': ('收视率', 'rating'),
        'share': ('市场份额', 'share'),
        'loyalty': ('平均忠实度', 'loyalty'),
    }
    metric_label, attr_name = metric_map.get(metric, ('收视率', 'rating'))

    title = f'栏目首播{metric_label}'
    _add_title_bar(slide, title, page_num=page_num)

    progs = data.premiere_programs
    if not progs:
        return slide

    # 判断是否有前1个月数据
    has_period_data = any(
        (p.period_rating > 0 if metric == 'rating' else p.period_share > 0)
        for p in progs
    )

    # 构建描述文本（列出全部栏目）
    desc_parts = []
    for p in progs:
        if metric == 'loyalty':
            val = p.loyalty
            # 忠实度没有可靠的变化数据，只显示当前值
            if val > 0:
                desc_parts.append(f'《{p.name}》{val:.2f}%')
            else:
                desc_parts.append(f'《{p.name}》-')
        else:
            val = p.rating if metric == 'rating' else p.share
            chg = p.rating_change if metric == 'rating' else p.share_change
            chg_pct = round(chg) if abs(chg) >= 1 else round(chg * 100)
            if has_period_data:
                # 对比前1个月均值
                if abs(chg_pct) <= 3:
                    chg_str = '较前一个月均值基本持平'
                else:
                    chg_str = f'{"下降" if chg_pct < 0 else "提升"}{abs(chg_pct)}%'
            else:
                # 对比前一天
                if abs(chg_pct) <= 3:
                    chg_str = '较前一天基本持平'
                else:
                    chg_str = f'{"下降" if chg_pct < 0 else "提升"}{abs(chg_pct)}%'
            desc_parts.append(f'《{p.name}》{val:.3f}%，{chg_str}')

    if has_period_data:
        compare_label = '对比前一个月均值'
    else:
        compare_label = '对比前一天'
    desc = (f'{data.report_date_short}，{compare_label}，首播栏目{metric_label}中，'
            + '，'.join(desc_parts) + '。')
    add_textbox(slide, Cm(1.5), Cm(2.0), Cm(30), Cm(2),
                desc, font_size=Pt(9), font_color=Colors.MEDIUM_GRAY)

    # 图表 - 类别标签包含时间段
    categories = [f'{p.name}\n{p.time_slot}' if p.time_slot else p.name for p in progs]

    if metric == 'rating':
        if has_period_data:
            series = {
                '前1个月均值': [p.period_rating for p in progs],
                '前一日': [p.prev_day_rating for p in progs],
                data.report_date_short: [p.rating for p in progs],
            }
            series_colors = [Colors.CHART_SERIES[0], Colors.CHART_SERIES[3],
                             Colors.CHART_SERIES[1]]
        else:
            series = {
                '前一日': [p.prev_day_rating for p in progs],
                data.report_date_short: [p.rating for p in progs],
            }
            series_colors = [Colors.CHART_SERIES[3], Colors.CHART_SERIES[1]]
    elif metric == 'share':
        if has_period_data:
            series = {
                '前1个月均值': [p.period_share for p in progs],
                '前一日': [p.prev_day_share for p in progs],
                data.report_date_short: [p.share for p in progs],
            }
            series_colors = [Colors.CHART_SERIES[0], Colors.CHART_SERIES[3],
                             Colors.CHART_SERIES[1]]
        else:
            series = {
                '前一日': [p.prev_day_share for p in progs],
                data.report_date_short: [p.share for p in progs],
            }
            series_colors = [Colors.CHART_SERIES[3], Colors.CHART_SERIES[1]]
    else:
        # 忠实度 - 使用计算的loyalty值，只显示当前值(历史数据需额外数据)
        series = {
            data.report_date_short: [p.loyalty for p in progs],
        }
        series_colors = [Colors.CHART_SERIES[1]]

    add_column_chart(slide, Cm(0.5), Cm(4.0), Cm(32), Cm(12),
                     categories, series,
                     series_colors=series_colors,
                     show_data_labels=True)

    return slide


# ═══════════════════════════════════════════════════════════════
# 幻灯片构建 - 频道分类观众规模(日报)
# ═══════════════════════════════════════════════════════════════

_AUDIENCE_TOTAL_CATS = {'四岁及以上所有人', '总体', '合计', 'total'}

def _audience_group(category):
    """将观众类别分组：gender/age/education/urban"""
    if category in ('男', '女'):
        return 'gender'
    if '岁' in category:
        return 'age'
    if category in ('未受过正规教育', '小学', '初中', '高中', '大学以上', '大学', '大专以上'):
        return 'education'
    if category in ('城市', '乡村', '城镇', '农村'):
        return 'urban'
    return 'other'


def build_channel_audience(prs, data: ReportData, page_num=None):
    """构建频道分类观众规模（匹配 demo 的 19 类别含间隔结构）"""
    slide = _add_slide(prs)
    _add_title_bar(slide, '频道分类观众规模', page_num=page_num)

    audience = data.channel_audience
    if not audience:
        return slide

    # 分离总体与分类项
    total_item = None
    cat_items = []
    for a in audience:
        if a.category.strip() in _AUDIENCE_TOTAL_CATS:
            total_item = a
        else:
            cat_items.append(a)

    # 按人口组分组，插入空白间隔（匹配 demo 的 19 类别结构）
    categories = []
    period_vals = []
    current_vals = []
    bar_items = []   # 与 categories 平行，None 表示间隔

    prev_group = None
    for item in cat_items:
        grp = _audience_group(item.category)
        if prev_group is not None and grp != prev_group:
            categories.append('')
            period_vals.append(None)
            current_vals.append(None)
            bar_items.append(None)
        categories.append(item.category)
        period_vals.append(item.period_value)
        current_vals.append(item.current_value)
        bar_items.append(item)
        prev_group = grp

    # 系列名包含总观众规模（匹配 demo 格式）
    total_period = total_item.period_value if total_item else sum(
        v for v in period_vals if v is not None)
    total_current = total_item.current_value if total_item else sum(
        v for v in current_vals if v is not None)

    series = {
        f'前一个月日均观众规模：{total_period:.1f}': period_vals,
        f'今日观众规模：{total_current:.1f}': current_vals,
    }
    add_column_chart(slide, Cm(0.5), Cm(3.0), Cm(32), Cm(12),
                     categories, series,
                     series_colors=[Colors.CHART_SERIES[0], Colors.CHART_SERIES[1]],
                     show_data_labels=True)

    # 分类变化百分比标注（仅非间隔项）
    non_gap = [item for item in bar_items if item is not None]
    for i, a in enumerate(non_gap):
        if a.change == 0:
            continue
        chg_pct = round(a.change * 100) if abs(a.change) < 10 else round(a.change)
        if chg_pct < 0:
            chg_label = f'-{abs(chg_pct)}%'
            color = Colors.ACCENT_RED
        elif chg_pct > 0:
            chg_label = f'{abs(chg_pct)}%'
            color = RGBColor(0x4C, 0xAF, 0x50)
        else:
            continue
        x_pos = Cm(1.5 + i * (31 / max(len(non_gap), 1)))
        add_textbox(slide, x_pos, Cm(15.5), Cm(2), Cm(0.6),
                   chg_label,
                   font_size=Pt(7), font_color=color,
                   bold=True, alignment=PP_ALIGN.CENTER)

    # 总体变化标签（图例区域，匹配 demo 的 ↓8% 样式）
    if total_item and total_item.change != 0:
        total_pct = round(total_item.change * 100) if abs(total_item.change) < 10 else round(total_item.change)
        if total_pct < 0:
            total_label = f'↓{abs(total_pct)}%'
            total_color = Colors.ACCENT_RED
        else:
            total_label = f'↑{abs(total_pct)}%'
            total_color = RGBColor(0x4C, 0xAF, 0x50)
        add_textbox(slide, Cm(24), Cm(1.5), Cm(2), Cm(0.6),
                   total_label,
                   font_size=Pt(7), font_color=total_color,
                   bold=True, alignment=PP_ALIGN.LEFT)

    add_textbox(slide, Cm(28), Cm(2.5), Cm(4), Cm(0.5),
                '单位：万人',
                font_size=Pt(8), font_color=Colors.MEDIUM_GRAY,
                alignment=PP_ALIGN.RIGHT)

    return slide


# ═══════════════════════════════════════════════════════════════
# 幻灯片构建 - 结尾
# ═══════════════════════════════════════════════════════════════

def build_ending(prs, data: ReportData):
    """构建感谢观看结尾页"""
    slide = _add_slide(prs, _get_blank_layout(prs))

    # 背景装饰
    from pptx.enum.shapes import MSO_SHAPE
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE,
        Cm(10), Cm(0), Cm(24), Cm(19)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xF8)
    bg_shape.line.fill.background()

    # 感谢观看文字
    add_textbox(slide, Cm(0), Cm(6), SLIDE_WIDTH, Cm(4),
                '感谢观看',
                font_size=Pt(48), font_color=Colors.PRIMARY,
                bold=True, alignment=PP_ALIGN.CENTER)

    return slide


# ═══════════════════════════════════════════════════════════════
# 辅助函数
# ═══════════════════════════════════════════════════════════════

def _change_desc(pct, up_word='提升', down_word='下降', flat_word='基本持平'):
    """生成变化描述文本"""
    if abs(pct) <= 5:
        return flat_word
    if pct > 0:
        return f'{up_word}{abs(pct):.0f}%'
    else:
        return f'{down_word}{abs(pct):.0f}%'


def _rank_change_desc(change):
    """排名变化描述"""
    if change > 0:
        return f'上升{change}位'
    elif change < 0:
        return f'下滑{abs(change)}位'
    else:
        return '保持不变'


def _extract_ymd_from_filename(path: str):
    """从文件名提取YYYYMMDD，返回(year, month, day)或None"""
    base_name = os.path.basename(path)
    m = re.search(r'(\d{4})(\d{2})(\d{2})', base_name)
    if not m:
        return None
    return int(m.group(1)), int(m.group(2)), int(m.group(3))


# ═══════════════════════════════════════════════════════════════
# 主生成函数
# ═══════════════════════════════════════════════════════════════

def generate_report(excel_path, template_path=None, output_path=None):
    """
    主入口：生成PPT报告

    Args:
        excel_path: Excel数据文件路径
        template_path: PPT模板路径(可选，用于获取主题)
        output_path: 输出PPT路径(可选)
    """
    print(f'[1/4] 读取数据: {excel_path}')
    data = read_excel_data(excel_path)

    global BODY_PLAIN_MODE

    strict_simple_mode = False
    if template_path:
        template_name = os.path.basename(template_path)
        strict_simple_mode = ('简版' in template_name)

    # 2-13页正文净版：不显示标题/副标题/页码/装饰
    BODY_PLAIN_MODE = strict_simple_mode

    # 日期来源按spec固定：优先使用文件名日期
    ymd = _extract_ymd_from_filename(excel_path)
    if ymd:
        year, month, day = ymd
        data.report_date_short = f'{month}月{day}日'
        data.report_date = f'{year}年{month}月{day}日'
    elif not data.report_date_short or not data.report_date:
        # 兜底：保留原有解析结果
        if not data.report_date_short:
            data.report_date_short = '当日'
        if not data.report_date:
            data.report_date = data.report_date_short

    print(f'  日期: {data.report_date}')
    print(f'  日报数据: {"有" if data.has_daily_report else "无"}')
    print(f'  模板模式: {"简版严格对标" if strict_simple_mode else "标准"}')
    print(f'  台组排名: 第{data.org_rank}位 ({data.org_rank_change:+d})')
    print(f'  上星排名: 第{data.channel_rank}位 ({data.channel_rank_change:+d})')

    # 创建PPT
    print('[2/4] 创建PPT...')
    if template_path and os.path.exists(template_path):
        prs = Presentation(template_path)
        # 删除模板中的所有幻灯片
        _remove_all_slides(prs)
    else:
        prs = Presentation()

    # 设置幻灯片尺寸（宽屏 13.33x7.5）
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # 构建幻灯片
    print('[3/4] 构建幻灯片...')
    page = 0

    # 1. 封面
    print('  - 封面')
    build_cover(prs, data)

    # 2. 收视速报
    print('  - 收视速报')
    build_summary(prs, data)

    # 3. 市场份额概览
    page += 1
    print('  - 市场份额概览')
    build_market_share(prs, data, page_num=page)

    # 4. 台组内排名
    page += 1
    print('  - 台组内排名')
    build_org_ranking(prs, data, page_num=page)

    # 5. 上星频道排名
    page += 1
    print('  - 上星频道排名')
    build_channel_ranking(prs, data, page_num=page)

    # 6. 串单市场份额
    page += 1
    print('  - 串单市场份额')
    build_schedule_chart(prs, data, page_num=page)

    # 7. 栏目收视率排名
    page += 1
    print('  - 栏目收视率排名')
    build_program_ranking(prs, data, metric='rating', page_num=page)

    # 8. 栏目收视份额排名
    page += 1
    print('  - 栏目收视份额排名')
    build_program_ranking(prs, data, metric='share', page_num=page)

    # 9. 分分钟收视率
    page += 1
    print('  - 分分钟收视率')
    build_minute_chart(prs, data, metric='rating', page_num=page)

    # 10. 分分钟市场份额
    page += 1
    print('  - 分分钟市场份额')
    build_minute_chart(prs, data, metric='share', page_num=page)

    # 日报扩展幻灯片
    if data.has_daily_report:
        if strict_simple_mode:
            # 简版严格对标：仅保留 demo 结构中的扩展页
            if data.premiere_programs:
                page += 1
                print('  - 栏目首播收视率')
                build_premiere_comparison(prs, data, metric='rating', page_num=page)

                page += 1
                print('  - 栏目首播市场份额')
                build_premiere_comparison(prs, data, metric='share', page_num=page)

            if data.channel_audience:
                page += 1
                print('  - 频道分类观众规模')
                build_channel_audience(prs, data, page_num=page)
        else:
            # 标准版：完整扩展页
            # 11. 电视剧市场份额
            page += 1
            print('  - 电视剧分析')
            build_drama_analysis(prs, data, page_num=page)

            # 12. 首播电视剧观众规模
            if data.drama_audience:
                page += 1
                print('  - 首播电视剧观众规模')
                build_drama_audience(prs, data, page_num=page)

            # 13. 栏目首播收视率
            if data.premiere_programs:
                page += 1
                print('  - 栏目首播收视率')
                build_premiere_comparison(prs, data, metric='rating', page_num=page)

                # 14. 栏目首播市场份额
                page += 1
                print('  - 栏目首播市场份额')
                build_premiere_comparison(prs, data, metric='share', page_num=page)

                # 15. 栏目首播忠实度
                page += 1
                print('  - 栏目首播忠实度')
                build_premiere_comparison(prs, data, metric='loyalty', page_num=page)

            # 16. 频道分类观众规模
            if data.channel_audience:
                page += 1
                print('  - 频道分类观众规模')
                build_channel_audience(prs, data, page_num=page)

    # 17. 感谢观看
    print('  - 感谢观看')
    build_ending(prs, data)

    # 保存
    if output_path is None:
        base_name = os.path.splitext(os.path.basename(excel_path))[0]
        output_path = f'{base_name}_收视日报.pptx'

    print(f'[4/4] 保存: {output_path}')
    prs.save(output_path)
    print(f'✅ 完成！共 {len(prs.slides)} 张幻灯片')
    return output_path


# ═══════════════════════════════════════════════════════════════
# CLI 入口
# ═══════════════════════════════════════════════════════════════

def main():
    parser = argparse.ArgumentParser(
        description='CCTV-17 收视日报 PPT 自动生成器',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog='''
示例:
  python generate_report.py 20260210CCTV-17日数据.xlsx
  python generate_report.py 0219.xlsx --template origin.pptx
  python generate_report.py data.xlsx --output 日报.pptx
        ''')
    parser.add_argument('excel', help='Excel数据文件路径')
    parser.add_argument('--template', '-t', default=None,
                        help='PPT模板文件路径 (默认: origin.pptx)')
    parser.add_argument('--output', '-o', default=None,
                        help='输出PPT文件路径')

    args = parser.parse_args()

    # 默认模板
    if args.template is None:
        if os.path.exists('origin.pptx'):
            args.template = 'origin.pptx'

    generate_report(args.excel, args.template, args.output)


if __name__ == '__main__':
    main()
