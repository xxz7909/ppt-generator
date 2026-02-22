# -*- coding: utf-8 -*-
"""
全页配置驱动版生成器 v3 —— XML 级精确替换

核心思路：
1) 用 generate_report.py 从 Excel 生成"数据稿"PPT（含各页数值、图表、表格）
2) **直接复制 demo** 作为输出文件（保留所有形状、格式、样式、换行 <a:br>）
3) 逐页把数据稿内容**在 XML 级别**回填到 demo 副本的对应控件
   - 文本框：仅替换 <a:t> 文本值，保留 <a:rPr>（字体/字号/颜色/粗细）
   - 图表：用 replace_data 替换数据
   - 表格：逐单元格 XML 级替换文本
4) 第 1 页封面做特殊处理：精确替换日期 run，保留 <a:br> 换行和标题

确保输出 PPT 与 demo 格式完全一致（字号、字体、颜色、换行、空行、标点）。
"""

import argparse
import json
import os
import re
import shutil
import sys
import tempfile
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.chart.data import CategoryChartData

import generate_report as base

# ─── XML 命名空间 ───
A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
NSMAP = {'a': A_NS, 'p': P_NS}

# 控制字符清洗
_VTAB_PATTERN = re.compile(r'_x000[0-9A-Fa-f]_')


def _sanitize(text):
    """清除 _x000B_ 等控制字符转义"""
    if text is None:
        return ''
    s = str(text)
    s = _VTAB_PATTERN.sub('', s)
    s = s.replace('\x0b', '').replace('\x0c', '')
    return s


# ═══════════════════════════════════════════════════════════════
# XML 级文本操作
# ═══════════════════════════════════════════════════════════════

def _get_para_text(p_elem):
    """从 <a:p> 元素提取纯文本（<a:br> → \\n）"""
    parts = []
    for child in p_elem:
        tag = etree.QName(child.tag).localname
        if tag == 'r':
            t = child.find('a:t', NSMAP)
            if t is not None and t.text:
                parts.append(t.text)
        elif tag == 'br':
            parts.append('\n')
    return ''.join(parts)


def _get_runs(p_elem):
    """获取段落内所有 <a:r> 元素"""
    return p_elem.findall('a:r', NSMAP)


def _get_breaks(p_elem):
    """获取段落内所有 <a:br> 元素"""
    return p_elem.findall('a:br', NSMAP)


def _set_run_text(run_elem, text):
    """设置 run 的 <a:t> 文本值，保留 <a:rPr>"""
    t = run_elem.find('a:t', NSMAP)
    if t is not None:
        t.text = text
    else:
        t = etree.SubElement(run_elem, '{http://schemas.openxmlformats.org/drawingml/2006/main}t')
        t.text = text


def _replace_para_text(target_p, new_text):
    """
    替换段落文本，保留所有格式。

    策略：
    - 如果段落有 <a:br>：按 \\n 拆分文本，分配到 br 前后的 run 组
    - 否则：将全部文本写入第一个 run，清空其余 run
    - 始终保留每个 run 的 <a:rPr>（字体/字号/颜色/粗细）
    """
    runs = _get_runs(target_p)
    breaks = _get_breaks(target_p)

    if not runs:
        return

    if breaks and '\n' in new_text:
        # 有 <a:br>：按换行分段分配到 run 组
        segments = new_text.split('\n')
        # 按 XML 子元素顺序将 run 分组（以 br 为分隔）
        run_groups = []
        current_group = []
        for child in target_p:
            tag = etree.QName(child.tag).localname
            if tag == 'r':
                current_group.append(child)
            elif tag == 'br':
                run_groups.append(current_group)
                current_group = []
        run_groups.append(current_group)

        for gi, group in enumerate(run_groups):
            seg_text = segments[gi] if gi < len(segments) else ''
            if group:
                _set_run_text(group[0], seg_text)
                for r in group[1:]:
                    _set_run_text(r, '')
    else:
        # 无 br：全部文本写入第一个 run，清空其余
        # 若源文本含 \n（来自 <a:br>）但目标无 <a:br>，只取第一行，
        # 避免将后续行重复写入（其余目标段落保持原样）
        text_to_write = new_text.split('\n')[0] if '\n' in new_text else new_text
        _set_run_text(runs[0], text_to_write)
        for r in runs[1:]:
            _set_run_text(r, '')


def _clear_para_text(target_p):
    """清空段落所有 run 的文本"""
    for r in _get_runs(target_p):
        _set_run_text(r, '')


# ═══════════════════════════════════════════════════════════════
# Shape 级操作
# ═══════════════════════════════════════════════════════════════

def _shape_type(shape):
    if getattr(shape, 'has_chart', False):
        return 'chart'
    if getattr(shape, 'has_table', False):
        return 'table'
    if getattr(shape, 'has_text_frame', False):
        return 'text'
    return 'other'


def _shape_center(shape):
    return (int(shape.left + shape.width / 2), int(shape.top + shape.height / 2))


def _dist2(a, b):
    return (a[0] - b[0]) ** 2 + (a[1] - b[1]) ** 2


def _iter_shapes_by_type(slide, type_name):
    out = []
    for idx, shp in enumerate(slide.shapes):
        if _shape_type(shp) == type_name:
            out.append((idx, shp))
    return out


def _match_nearest(target_items, src_items):
    """
    两轮匹配：
    1. 短文本先做内容匹配（防止静态标签被交叉覆盖）
    2. 剩余的做最近邻位置匹配
    """
    mapping = []
    used_src = set()
    used_tgt = set()

    # 第一轮：短文本内容匹配
    for ti, tshp in target_items:
        if not getattr(tshp, 'has_text_frame', False):
            continue
        ttext = tshp.text.strip().replace('\x0b', '\n')
        if not ttext or len(ttext) > 30:
            continue
        for si, sshp in src_items:
            if si in used_src:
                continue
            if not getattr(sshp, 'has_text_frame', False):
                continue
            stext = sshp.text.strip().replace('\x0b', '\n')
            if stext == ttext:
                mapping.append((ti, si))
                used_src.add(si)
                used_tgt.add(ti)
                break

    # 第二轮：剩余的做最近邻位置匹配
    for ti, tshp in target_items:
        if ti in used_tgt:
            continue
        tc = _shape_center(tshp)
        best = None
        best_d = None
        for si, sshp in src_items:
            if si in used_src:
                continue
            d = _dist2(tc, _shape_center(sshp))
            if best is None or best_d is None or d < best_d:
                best = si
                best_d = d
        if best is not None:
            used_src.add(best)
            mapping.append((ti, best))
    return mapping


def _match_charts_by_y_order(target_items, src_items):
    """
    按 Y 坐标（top）排序配对图表。

    最近邻匹配在图表位置偏差较大时会交叉配对（如市场份额页的
    所有频道 vs CCTV-17 收视率条形图）。改用按 top 排序后一一
    配对，确保上方图表匹配上方图表、下方匹配下方。
    """
    t_sorted = sorted(target_items, key=lambda x: x[1].top)
    s_sorted = sorted(src_items, key=lambda x: x[1].top)
    n = min(len(t_sorted), len(s_sorted))
    return [(t_sorted[i][0], s_sorted[i][0]) for i in range(n)]


# ═══════════════════════════════════════════════════════════════
# 文本同步（XML 级）
# ═══════════════════════════════════════════════════════════════

def _sync_text_shape(target_shape, src_shape):
    """
    将 src_shape 的文本内容写入 target_shape，
    但完全保留 target_shape 的 XML 格式（rPr / br / pPr）。

    使用"非空段落对齐"策略：
    - 提取两边的非空段落索引
    - 按顺序一一配对
    - 空段落保持不变（保留空行分隔结构）
    """
    src_text = src_shape.text if hasattr(src_shape, 'text') else ''
    if not src_text or not src_text.strip():
        return

    # 内容相同则跳过（防止静态标签被错误交叉覆盖）
    # 统一 \x0b / \n 后比较
    target_text = target_shape.text if hasattr(target_shape, 'text') else ''
    _norm = lambda t: _sanitize(t).replace('\x0b', '').replace('\n', '').strip()
    if _norm(src_text) == _norm(target_text):
        return

    # 获取 <p:txBody> → <a:p> 列表
    t_body = target_shape._element.find(f'.//{{{P_NS}}}txBody')
    s_body = src_shape._element.find(f'.//{{{P_NS}}}txBody')
    if t_body is None:
        t_body = target_shape._element.find(f'.//{{{A_NS}}}txBody')
    if s_body is None:
        s_body = src_shape._element.find(f'.//{{{A_NS}}}txBody')
    if t_body is None or s_body is None:
        return

    t_paras = t_body.findall('a:p', NSMAP)
    s_paras = s_body.findall('a:p', NSMAP)

    # 提取非空段落索引
    t_content = [i for i, p in enumerate(t_paras) if _get_para_text(p).strip()]
    s_content = [i for i, p in enumerate(s_paras) if _get_para_text(p).strip()]

    # 按非空段落一一配对
    n = min(len(t_content), len(s_content))
    for j in range(n):
        ti = t_content[j]
        si = s_content[j]
        src_text_j = _sanitize(_get_para_text(s_paras[si]))
        _replace_para_text(t_paras[ti], src_text_j)

    # 如果源有更多非空段落，追加到最后一个已配对的目标段落
    if len(s_content) > len(t_content) and t_content:
        extra_texts = []
        for j in range(len(t_content), len(s_content)):
            txt = _sanitize(_get_para_text(s_paras[s_content[j]]))
            if txt:
                extra_texts.append(txt)
        if extra_texts:
            last_tp = t_paras[t_content[-1]]
            runs = _get_runs(last_tp)
            if runs:
                t_elem = runs[-1].find('a:t', NSMAP)
                if t_elem is not None:
                    t_elem.text = (t_elem.text or '') + '\n'.join([''] + extra_texts)


# ═══════════════════════════════════════════════════════════════
# 图表数据同步
# ═══════════════════════════════════════════════════════════════

def _extract_chart_data(src_chart):
    chart_data = CategoryChartData()
    categories = []
    try:
        categories = [_sanitize(str(c.label)) if c.label is not None else ''
                      for c in src_chart.plots[0].categories]
    except (AttributeError, TypeError, IndexError):
        pass
    chart_data.categories = categories

    for ser in src_chart.series:
        values = []
        try:
            values = list(ser.values)
        except (AttributeError, TypeError):
            pass
        name = _sanitize(ser.name) if ser.name else 'Series'
        chart_data.add_series(name, values)
    return chart_data


def _sync_chart(target_shape, src_shape):
    try:
        # 先保存 demo 模板中的数据标签数字格式（replace_data 会重置）
        C_NS = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
        cs = target_shape.chart._chartSpace
        saved_nf_info = []  # [(dLbls_element, formatCode, sourceLinked)]
        for dLbls in cs.findall(f'.//{{{C_NS}}}dLbls'):
            nf = dLbls.find(f'{{{C_NS}}}numFmt')
            if nf is not None:
                saved_nf_info.append((
                    dLbls,
                    nf.get('formatCode', ''),
                    nf.get('sourceLinked', '0'),
                ))

        target_shape.chart.replace_data(_extract_chart_data(src_shape.chart))

        # 恢复 demo 模板的数据标签数字格式
        if saved_nf_info:
            # replace_data 后重新获取 chartSpace（可能被替换）
            cs = target_shape.chart._chartSpace
            all_dLbls = cs.findall(f'.//{{{C_NS}}}dLbls')
            for i, (_, fmt_code, src_linked) in enumerate(saved_nf_info):
                if i >= len(all_dLbls) or not fmt_code:
                    continue
                dLbls = all_dLbls[i]
                # 移除旧 numFmt
                for old_nf in dLbls.findall(f'{{{C_NS}}}numFmt'):
                    dLbls.remove(old_nf)
                # 插入保存的 numFmt
                nf_elem = etree.SubElement(dLbls, f'{{{C_NS}}}numFmt')
                nf_elem.set('formatCode', fmt_code)
                nf_elem.set('sourceLinked', src_linked)
        else:
            # 模板无 numFmt，尝试从 draft 取
            src_plot = src_shape.chart.plots[0]
            tgt_plot = target_shape.chart.plots[0]
            if src_plot.has_data_labels:
                src_nf = src_plot.data_labels.number_format
                src_linked = src_plot.data_labels.number_format_is_linked
                if not src_linked and src_nf and src_nf != 'General':
                    cs = target_shape.chart._chartSpace
                    for dLbls in cs.findall(f'.//{{{C_NS}}}dLbls'):
                        for old_nf in dLbls.findall(f'{{{C_NS}}}numFmt'):
                            dLbls.remove(old_nf)
                    tgt_plot.has_data_labels = True
                    tgt_plot.data_labels.number_format = src_nf
                    tgt_plot.data_labels.number_format_is_linked = False
    except (AttributeError, TypeError):
        pass


def _sync_combo_chart(target_shape, src_shape):
    """同步组合图表（柱状+折线），XML 级别替换数据，保留双图表结构。

    target 是从 demo 拷贝来的组合图表（barChart + lineChart），
    src 是 draft 生成的纯柱状图（barChart 含两个 series：市场份额% 和 收视率）。
    本函数：
      1. 从 src 提取 categories、bar values (series 0)、line values (series 1)
      2. 在 target 的 barChart/ser 和 lineChart/ser 中替换 strCache 和 numCache
      3. 更新 ptCount 以匹配新数据长度
    """
    C_NS = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
    try:
        src_chart = src_shape.chart
        tgt_chart = target_shape.chart

        # --- 从 src 提取数据 ---
        src_categories = []
        src_bar_vals = []
        src_line_vals = []

        try:
            src_categories = [
                _sanitize(str(c.label)) if c.label is not None else ''
                for c in src_chart.plots[0].categories
            ]
        except (AttributeError, TypeError, IndexError):
            pass

        # Series 0 = 市场份额 (bar), Series 1 = 收视率 (line)
        for i, ser in enumerate(src_chart.series):
            vals = list(ser.values) if ser.values else []
            if i == 0:
                src_bar_vals = vals
            elif i == 1:
                src_line_vals = vals

        n = len(src_categories)
        if n == 0:
            return

        # --- XML 级替换 ---
        tgt_cs = tgt_chart._chartSpace
        pa = tgt_cs.find(f'.//{{{C_NS}}}plotArea')

        bar_chart = pa.find(f'{{{C_NS}}}barChart')
        line_chart = pa.find(f'{{{C_NS}}}lineChart')

        def _replace_cache(ser_elem, categories, values, is_num=False):
            """替换一个 <c:ser> 下的 strCache (cat) 和 numCache (val)"""
            # 替换分类 strCache
            cat_elem = ser_elem.find(f'{{{C_NS}}}cat')
            if cat_elem is not None:
                str_ref = cat_elem.find(f'{{{C_NS}}}strRef')
                if str_ref is not None:
                    old_cache = str_ref.find(f'{{{C_NS}}}strCache')
                    if old_cache is not None:
                        str_ref.remove(old_cache)
                    new_cache = etree.SubElement(str_ref, f'{{{C_NS}}}strCache')
                    pt_count = etree.SubElement(new_cache, f'{{{C_NS}}}ptCount')
                    pt_count.set('val', str(len(categories)))
                    for idx, cat_text in enumerate(categories):
                        pt = etree.SubElement(new_cache, f'{{{C_NS}}}pt')
                        pt.set('idx', str(idx))
                        v = etree.SubElement(pt, f'{{{C_NS}}}v')
                        v.text = cat_text

            # 替换数据 numCache
            val_elem = ser_elem.find(f'{{{C_NS}}}val')
            if val_elem is not None:
                num_ref = val_elem.find(f'{{{C_NS}}}numRef')
                if num_ref is not None:
                    old_cache = num_ref.find(f'{{{C_NS}}}numCache')
                    if old_cache is not None:
                        num_ref.remove(old_cache)
                    new_cache = etree.SubElement(num_ref, f'{{{C_NS}}}numCache')
                    fmt = etree.SubElement(new_cache, f'{{{C_NS}}}formatCode')
                    fmt.text = '0.000'
                    pt_count = etree.SubElement(new_cache, f'{{{C_NS}}}ptCount')
                    pt_count.set('val', str(len(values)))
                    for idx, val in enumerate(values):
                        pt = etree.SubElement(new_cache, f'{{{C_NS}}}pt')
                        pt.set('idx', str(idx))
                        v_elem = etree.SubElement(pt, f'{{{C_NS}}}v')
                        v_elem.text = str(val) if val is not None else '0'

        # 替换 barChart 数据
        if bar_chart is not None:
            bar_ser = bar_chart.find(f'{{{C_NS}}}ser')
            if bar_ser is not None:
                _replace_cache(bar_ser, src_categories, src_bar_vals)

        # 替换 lineChart 数据
        if line_chart is not None:
            line_ser = line_chart.find(f'{{{C_NS}}}ser')
            if line_ser is not None:
                _replace_cache(line_ser, src_categories, src_line_vals)

    except (AttributeError, TypeError) as e:
        print(f'[WARN] _sync_combo_chart error: {e}')


# ═══════════════════════════════════════════════════════════════
# 表格数据同步
# ═══════════════════════════════════════════════════════════════

def _sync_table(target_shape, src_shape):
    """替换表格单元格文本，保留目标格式（XML 级）"""
    tt = target_shape.table
    st = src_shape.table
    rows = min(len(tt.rows), len(st.rows))
    cols = min(len(tt.columns), len(st.columns))

    a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    for r in range(rows):
        for c in range(cols):
            new_text = _sanitize(st.cell(r, c).text)
            tc_elem = tt.cell(r, c)._tc
            tc_body = tc_elem.find(f'.//{{{a_ns}}}txBody')
            if tc_body is not None:
                paras = tc_body.findall('a:p', NSMAP)
                if paras:
                    _replace_para_text(paras[0], new_text)
                    for p in paras[1:]:
                        _clear_para_text(p)
            else:
                tt.cell(r, c).text = new_text


# ═══════════════════════════════════════════════════════════════
# 标题修复（简版模式下 data draft 无标题，需内联构建）
# ═══════════════════════════════════════════════════════════════

def _change_desc(pct, up_word='提升', down_word='下降', flat_word='基本持平'):
    if abs(pct) <= 5:
        return flat_word
    return f'{up_word}{abs(pct):.0f}%' if pct > 0 else f'{down_word}{abs(pct):.0f}%'


def _rank_change_desc(change):
    if change > 0:
        return f'上升{change}位'
    elif change < 0:
        return f'下滑{abs(change)}位'
    return '保持不变'


def _build_slide_titles(data):
    """
    根据 ReportData 构建各页标题（与 generate_report.py 逻辑一致）。
    返回 dict: {slide_index: title_text}，仅含数据相关标题。
    """
    ms = data.market_share
    titles = {}

    # Page 3 (idx=2): 市场份额概览
    share_chg = _change_desc(ms.share_change)
    rating_chg = _change_desc(ms.rating_change)
    titles[2] = f'{data.report_date_short}，市场份额{share_chg}，收视率{rating_chg}'

    # Page 4 (idx=3): 台组内排名
    org_chg = _rank_change_desc(data.org_rank_change)
    titles[3] = f'央视台组排名第{data.org_rank}位，较前一日{org_chg}'

    # Page 5 (idx=4): 上星频道排名
    ch_chg = _rank_change_desc(data.channel_rank_change)
    titles[4] = f'上星频道排名第{data.channel_rank}位，较前一日{ch_chg}'

    # Page 6 (idx=5): 串单市场份额
    titles[5] = (f'{data.report_date_short}，频道市场份额{ms.cctv17_current_share:.3f}%，'
                 f'收视率{ms.cctv17_current_rating:.3f}%')

    # Page 7 (idx=6): 栏目收视率排名
    titles[6] = f'{data.report_date_short}栏目收视率排名'

    # Page 8 (idx=7): 栏目收视份额排名
    titles[7] = f'{data.report_date_short}栏目收视份额排名'

    # Page 9-13: 静态标题，无需更新
    return titles


def _find_title_shape(slide):
    """查找幻灯片中的标题 shape（按名称识别）"""
    for shp in slide.shapes:
        if not getattr(shp, 'has_text_frame', False):
            continue
        name = shp.name or ''
        if '标题' in name:
            return shp
    return None


def _fix_titles(target_prs, data):
    """修复简版模板中未被数据稿同步的标题文本"""
    titles = _build_slide_titles(data)
    for slide_idx, title_text in titles.items():
        if slide_idx >= len(target_prs.slides):
            continue
        shp = _find_title_shape(target_prs.slides[slide_idx])
        if shp is None:
            continue
        # 找到标题段落并用 XML 级替换
        body = shp._element.find(f'.//{{{P_NS}}}txBody')
        if body is None:
            body = shp._element.find(f'.//{{{A_NS}}}txBody')
        if body is None:
            continue
        paras = body.findall('a:p', NSMAP)
        if paras:
            # 标题通常只有一个段落，有前导空格的 run（如"  收视速报"）需保留空格
            # 但 data-dependent titles 不含前导空格
            _replace_para_text(paras[0], title_text)
            for p in paras[1:]:
                _clear_para_text(p)
        print(f'    ✧ 标题已更新: 第{slide_idx + 1}页 → {title_text[:40]}...')


# ═══════════════════════════════════════════════════════════════
# 第 1 页封面特殊处理
# ═══════════════════════════════════════════════════════════════

def _update_cover(target_slide, src_slide):
    """
    封面页精确更新：
    - 标题 shape：只替换日期相关 run 的 <a:t>，保留 <a:br> 换行和标题文字
    - 署名 shape：只替换年份和期号 run
    """
    # ── 从源幻灯片提取日期和期号 ──
    src_title_text = ''
    src_subtitle_text = ''
    for shp in src_slide.shapes:
        if not getattr(shp, 'has_text_frame', False):
            continue
        text = shp.text
        if '频道收视日报' in text or '农业' in text:
            src_title_text = text
        elif '统筹策划部' in text or '策划部' in text:
            src_subtitle_text = text

    # 提取年月日
    date_match = re.search(r'(\d{4})年(\d{1,2})月(\d{1,2})日', src_title_text)
    new_year = date_match.group(1) if date_match else None
    new_month = date_match.group(2) if date_match else None
    new_day = date_match.group(3) if date_match else None

    # 提取期号
    period_match = re.search(r'第(\d+)期', src_subtitle_text)
    new_period = period_match.group(1) if period_match else None

    # ── 更新目标标题 shape 的日期 run ──
    for shp in target_slide.shapes:
        if not getattr(shp, 'has_text_frame', False):
            continue
        full_text = shp.text

        if '频道收视日报' in full_text or '农业' in full_text:
            # 标题 shape：逐 run 精确替换日期数字
            for p in shp.text_frame._element.findall('a:p', NSMAP):
                runs = p.findall('a:r', NSMAP)
                for i, r in enumerate(runs):
                    t = r.find('a:t', NSMAP)
                    if t is None or t.text is None:
                        continue
                    val = t.text.strip()

                    # 年份 run：4 位数字
                    if val.isdigit() and len(val) == 4 and new_year:
                        t.text = new_year
                        continue

                    # 月份 / 日期 run：1-2 位数字，根据相邻 run 判断
                    if val.isdigit() and len(val) <= 2:
                        # 看后一个 run
                        if i + 1 < len(runs):
                            next_t = runs[i + 1].find('a:t', NSMAP)
                            if next_t is not None and next_t.text:
                                if '月' in next_t.text and new_month:
                                    t.text = new_month
                                    continue
                        # 看前一个 run
                        if i > 0:
                            prev_t = runs[i - 1].find('a:t', NSMAP)
                            if prev_t is not None and prev_t.text:
                                if '月' in prev_t.text and new_day:
                                    t.text = new_day
                                    continue

        elif '统筹策划部' in full_text or '策划部' in full_text:
            # 署名 shape：更新年份和期号
            for p in shp.text_frame._element.findall('a:p', NSMAP):
                runs = p.findall('a:r', NSMAP)
                for i, r in enumerate(runs):
                    t = r.find('a:t', NSMAP)
                    if t is None or t.text is None:
                        continue

                    # 年份 run：含空格+4位数字如 " 2026"
                    year_m = re.match(r'^(\s*)(\d{4})$', t.text)
                    if year_m and new_year:
                        t.text = year_m.group(1) + new_year
                        continue

                    # 期号 run：纯数字且前面 run 含"第"
                    val = t.text.strip()
                    if val.isdigit() and new_period and i > 0:
                        prev_t = runs[i - 1].find('a:t', NSMAP)
                        if prev_t is not None and prev_t.text and '第' in prev_t.text:
                            t.text = new_period


# ═══════════════════════════════════════════════════════════════
# 逐页同步
# ═══════════════════════════════════════════════════════════════

def _is_change_label(text):
    """判断是否为变化百分比标签（如 ↑14%、-15%、↓8%）"""
    t = text.strip()
    if not t:
        return False
    return bool(re.match(r'^[↑↓▲▼+-]?\d+\.?\d*%$', t))


def _sync_audience_labels(target_slide, src_slide):
    """
    页面13（频道分类观众规模）专用：按 X 坐标排序匹配变化标签。
    demo 的标签散布在图表区各位置，draft 的标签排成一行。
    按 X 排序后一一对应。
    """
    # 收集 demo/draft 中的变化标签
    t_labels = []
    s_labels = []
    for idx, shp in enumerate(target_slide.shapes):
        if getattr(shp, 'has_text_frame', False) and _is_change_label(shp.text):
            t_labels.append((shp.left, idx, shp))
    for idx, shp in enumerate(src_slide.shapes):
        if getattr(shp, 'has_text_frame', False) and _is_change_label(shp.text):
            s_labels.append((shp.left, idx, shp))

    # 按 X 排序
    t_labels.sort(key=lambda x: x[0])
    s_labels.sort(key=lambda x: x[0])

    n = min(len(t_labels), len(s_labels))
    for i in range(n):
        t_shp = t_labels[i][2]
        s_shp = s_labels[i][2]
        _sync_text_shape(t_shp, s_shp)


def _compute_nice_axis(max_val):
    """计算 Y 轴最大值（模拟 PowerPoint 自动缩放）。"""
    import math
    if max_val <= 0:
        return 100
    raw_interval = max_val / 6
    magnitude = 10 ** math.floor(math.log10(raw_interval))
    candidates = [magnitude, 2 * magnitude, 5 * magnitude, 10 * magnitude]
    best = None
    for interval in candidates:
        n = math.ceil(max_val / interval) + 1          # +1 留余量
        axis_max = n * interval
        n_ticks = n + 1
        if 5 <= n_ticks <= 12:
            if best is None or axis_max < best:
                best = axis_max
    if best is None:
        interval = 2 * magnitude
        best = (math.ceil(max_val / interval) + 1) * interval
    return best


def _set_label_color(shape, hex_color):
    """
    将文本框中所有 run 和 endParaRPr 的字体颜色设置为指定的 srgbClr。

    替换现有 solidFill（无论是 schemeClr 还是 srgbClr）。
    """
    A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    P = 'http://schemas.openxmlformats.org/presentationml/2006/main'

    txBody = shape._element.find(f'.//{{{A}}}txBody')
    if txBody is None:
        txBody = shape._element.find(f'.//{{{P}}}txBody')
    if txBody is None:
        return

    # 处理所有 rPr 和 endParaRPr
    rpr_tags = [f'{{{A}}}rPr', f'{{{A}}}endParaRPr']
    for tag in rpr_tags:
        for rPr in txBody.findall(f'.//{tag}'):
            # 移除旧 solidFill
            for old_fill in rPr.findall(f'{{{A}}}solidFill'):
                rPr.remove(old_fill)
            # 添加新 solidFill
            sf = etree.SubElement(rPr, f'{{{A}}}solidFill')
            srgb = etree.SubElement(sf, f'{{{A}}}srgbClr')
            srgb.set('val', hex_color)
            # 放到 rPr 的第一个子元素位置（solidFill 通常在前面）
            rPr.insert(0, sf)


def _sync_and_position_audience(target_slide, src_slide):
    """
    Page 13（频道分类观众规模）：同步图表数据 + 动态定位变化标签。

    1. 同步图表数据（replace_data）
    2. 根据图表数据 **计算** 变化百分比文本
    3. 将标签 X 居中于对应类别柱组，Y 紧贴最高柱上方
    4. 图例区总体变化标签单独处理（仅替换文字）
    """
    C_NS = 'http://schemas.openxmlformats.org/drawingml/2006/chart'

    # ── 1. 图表数据同步 ──
    target_charts = _iter_shapes_by_type(target_slide, 'chart')
    src_charts = _iter_shapes_by_type(src_slide, 'chart')
    src_chart_map = {idx: shp for idx, shp in src_charts}
    chart_shape = None
    for ti, si in _match_nearest(target_charts, src_charts):
        t_shp = dict(target_charts)[ti]
        s_shp = src_chart_map[si]
        _sync_chart(t_shp, s_shp)
        chart_shape = t_shp

    if chart_shape is None:
        return

    # ── 1b. 清除 replace_data 产生的 "None" 间隔类别 ──
    cs_tmp = chart_shape.chart._chartSpace
    for str_cache in cs_tmp.findall(f'.//{{{C_NS}}}strCache'):
        for pt in list(str_cache.findall(f'{{{C_NS}}}pt')):
            v = pt.find(f'{{{C_NS}}}v')
            if v is not None and v.text in ('None', ''):
                str_cache.remove(pt)

    # ── 2. 图表几何计算 ──
    chart = chart_shape.chart
    cs = chart._chartSpace
    pa = cs.find(f'.//{{{C_NS}}}plotArea')
    man = pa.find(f'{{{C_NS}}}layout/{{{C_NS}}}manualLayout')

    if man is not None:
        px = float(man.find(f'{{{C_NS}}}x').get('val'))
        py = float(man.find(f'{{{C_NS}}}y').get('val'))
        pw = float(man.find(f'{{{C_NS}}}w').get('val'))
        ph = float(man.find(f'{{{C_NS}}}h').get('val'))
    else:
        px, py, pw, ph = 0.05, 0.15, 0.9, 0.6

    plot_left = chart_shape.left + int(px * chart_shape.width)
    plot_top = chart_shape.top + int(py * chart_shape.height)
    plot_width = int(pw * chart_shape.width)
    plot_height = int(ph * chart_shape.height)
    plot_bottom = plot_top + plot_height

    # ── 3. 读取图表数据 ──
    s0 = list(chart.series[0].values)
    s1 = list(chart.series[1].values)
    num_cats = len(s0)

    # 每个类别的 max 值与变化百分比
    cat_data = []          # (index, max_value, change_text)
    for i in range(num_cats):
        v0 = s0[i] if s0[i] is not None else 0
        v1 = s1[i] if s1[i] is not None else 0
        max_v = max(v0, v1)
        if max_v <= 0:
            continue          # 间隔类别
        change_pct = ((v1 - v0) / v0 * 100) if v0 > 0 else 0
        # 格式：|val|<1% 保留 1 位小数，否则取整；负值用 - 前缀，正值无前缀
        abs_chg = abs(change_pct)
        if abs_chg < 1 and abs_chg > 0:
            val_str = f'{abs_chg:.1f}'
        else:
            val_str = str(round(abs_chg))
        if change_pct < 0:
            text = f'-{val_str}%'
        elif change_pct > 0:
            text = f'{val_str}%'
        else:
            text = '0%'
        cat_data.append((i, max_v, text, change_pct))

    if not cat_data:
        return

    # ── 4. 计算坐标轴范围 & 设置显式最大值 ──
    overall_max = max(d[1] for d in cat_data)
    axis_max = _compute_nice_axis(overall_max)

    val_ax = pa.find(f'{{{C_NS}}}valAx')
    if val_ax is not None:
        scaling = val_ax.find(f'{{{C_NS}}}scaling')
        if scaling is not None:
            max_elem = scaling.find(f'{{{C_NS}}}max')
            if max_elem is None:
                max_elem = etree.SubElement(scaling, f'{{{C_NS}}}max')
            max_elem.set('val', str(axis_max))

    # ── 5. 分离柱上标签 vs 图例标签 ──
    # 图例标签用 ↑↓▲▼ 箭头前缀；柱上标签用 -/数字前缀
    bar_labels = []
    legend_label = None
    for sh in target_slide.shapes:
        if getattr(sh, 'has_text_frame', False) and _is_change_label(sh.text):
            t = sh.text.strip()
            if t and t[0] in '↑↓▲▼':
                legend_label = sh          # 图例区（如 ↓8%）
            else:
                bar_labels.append(sh)
    bar_labels.sort(key=lambda s: s.left)

    # 颜色常量：红升绿降
    COLOR_UP = 'FF0000'
    COLOR_DOWN = '00B050'

    # ── 6. 定位每个柱上标签 ──
    GAP_EMU = 55000          # 标签底边与柱顶的间距（~0.15 cm）
    n = min(len(bar_labels), len(cat_data))
    for i in range(n):
        lbl = bar_labels[i]
        cat_idx, max_v, new_text, chg_pct = cat_data[i]

        # X: 居中于类别柱组
        cat_center_x = plot_left + int(plot_width * (cat_idx + 0.5) / num_cats)
        lbl.left = cat_center_x - lbl.width // 2

        # Y: 紧贴最高柱上方
        bar_top = plot_bottom - int((max_v / axis_max) * plot_height)
        lbl.top = bar_top - lbl.height - GAP_EMU

        # 文字：XML 级替换（保留字体/颜色/粗细格式）
        p_ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        txBody = lbl._element.find(f'.//{{{A_NS}}}txBody')
        if txBody is None:
            txBody = lbl._element.find(f'.//{{{p_ns}}}txBody')
        if txBody is not None:
            paras = txBody.findall(f'{{{A_NS}}}p')
            if paras:
                _replace_para_text(paras[0], new_text)

        # 颜色：红升绿降
        if chg_pct != 0:
            _set_label_color(lbl, COLOR_UP if chg_pct > 0 else COLOR_DOWN)

    # ── 7. 图例变化标签（仅替换文字 + 颜色） ──
    if legend_label:
        total_s0 = sum(v for v in s0 if v is not None)
        total_s1 = sum(v for v in s1 if v is not None)
        total_chg = round((total_s1 - total_s0) / total_s0 * 100) if total_s0 > 0 else 0
        if total_chg > 0:
            leg_text = f'↑{total_chg}%'
        elif total_chg < 0:
            leg_text = f'↓{abs(total_chg)}%'
        else:
            leg_text = '0%'
        txBody = legend_label._element.find(f'.//{{{A_NS}}}txBody')
        if txBody is None:
            p_ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
            txBody = legend_label._element.find(f'.//{{{p_ns}}}txBody')
        if txBody is not None:
            paras = txBody.findall(f'{{{A_NS}}}p')
            if paras:
                _replace_para_text(paras[0], leg_text)
        # 颜色：红升绿降
        if total_chg != 0:
            _set_label_color(legend_label, COLOR_UP if total_chg > 0 else COLOR_DOWN)


def _fix_change_pct_colors(slide):
    """
    修正市场份额页（第3页）表格中变化百分比的字体颜色。

    规则（匹配 demo 风格）：
    - 正值（提升）→ 红色 FF0000
    - 负值（下降）→ 绿色 00B050
    - 零值 → 不改
    """
    COLOR_UP = 'FF0000'
    COLOR_DOWN = '00B050'

    for shp in slide.shapes:
        if not getattr(shp, 'has_table', False):
            continue
        tbl = shp.table
        for r in range(len(tbl.rows)):
            for c in range(len(tbl.columns)):
                cell = tbl.cell(r, c)
                text = cell.text.strip()
                if not text:
                    continue
                # 匹配变化百分比单元格：如 "+7%", "-8%", "6%"
                m = re.match(r'^([+-]?)(\d+\.?\d*)%$', text)
                if not m:
                    continue
                sign = m.group(1)
                val = float(m.group(2))
                if val == 0:
                    continue
                # 无符号正数或 + 号 → 提升（红），- 号 → 下降（绿）
                color = COLOR_DOWN if sign == '-' else COLOR_UP
                # XML 级别修改 run 字体颜色
                tc = cell._tc
                for run in tc.findall(f'.//{{{A_NS}}}r'):
                    rpr = run.find(f'{{{A_NS}}}rPr')
                    if rpr is None:
                        rpr = etree.SubElement(run, f'{{{A_NS}}}rPr')
                        # 插入到 <a:t> 之前
                        t_elem = run.find(f'{{{A_NS}}}t')
                        if t_elem is not None:
                            run.remove(rpr)
                            run.insert(list(run).index(t_elem), rpr)
                    # 移除旧的 solidFill
                    old_fill = rpr.find(f'{{{A_NS}}}solidFill')
                    if old_fill is not None:
                        rpr.remove(old_fill)
                    # 添加新颜色
                    fill = etree.SubElement(rpr, f'{{{A_NS}}}solidFill')
                    clr = etree.SubElement(fill, f'{{{A_NS}}}srgbClr')
                    clr.set('val', color)


def _fix_page3_numfmt(slide):
    """
    强制将第 3 页所有图表的数据标签数字格式统一为 3 位小数。

    demo 模板中 "图表 7"（所有频道收视率%）的 numFmt 是 #,##0.00（2 位），
    需要改为 #,##0.000 保持与其他图表一致。
    """
    C_NS = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
    TARGET_FMT = '#,##0.000_);[Red]\\(#,##0.000\\)'

    for shp in slide.shapes:
        if not shp.has_chart:
            continue
        cs = shp.chart._chartSpace
        for dLbls in cs.findall(f'.//{{{C_NS}}}dLbls'):
            nf = dLbls.find(f'{{{C_NS}}}numFmt')
            if nf is not None:
                fmt = nf.get('formatCode', '')
                # 只修正 2 位小数的格式（0.00）为 3 位（0.000）
                if '0.00' in fmt and '0.000' not in fmt:
                    new_fmt = fmt.replace('0.00', '0.000')
                    nf.set('formatCode', new_fmt)


def _load_threshold_config():
    """
    从模板配置文件 demo_layout_config.json 读取台组排名页（第 4 页）
    橘色阈值线的颜色和阈值文本。

    颜色来源：slideNo=4 中名为 "文本框 10" 的文本框字体颜色（与线同色）。
    阈值来源：同一文本框的文本内容（如 "0.83%"）。

    返回 (color_hex, threshold_text, threshold_value)
    """
    # 兼容 PyInstaller 打包：优先使用 _MEIPASS（打包后的临时目录）
    _base = Path(getattr(sys, '_MEIPASS', Path(__file__).resolve().parent))
    config_path = _base / 'demo_layout_config.json'
    default_color = 'ED7D31'
    default_text = '0.83%'
    default_value = 0.83
    try:
        with open(config_path, 'r', encoding='utf-8') as f:
            config = json.load(f)
        color = default_color
        text = default_text
        for slide_cfg in config.get('slides', []):
            if slide_cfg.get('slideNo') != 4:
                continue
            for ts in slide_cfg.get('textShapes', []):
                if ts.get('name') == '文本框 10':
                    # 读取阈值文本
                    raw_text = ts.get('text', '').strip()
                    if raw_text:
                        text = raw_text
                    for para in ts.get('paragraphs', []):
                        for run in para.get('runs', []):
                            rgb = run.get('colorRGB')
                            if rgb and len(rgb) == 3:
                                color = f'{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}'
        # 解析阈值数值
        m = re.match(r'([\d.]+)', text)
        value = float(m.group(1)) if m else default_value
        return color, text, value
    except (FileNotFoundError, json.JSONDecodeError, KeyError):
        return default_color, default_text, default_value


def _fix_org_ranking_page(slide, threshold_color='ED7D31',
                          threshold_text='0.83%', threshold_value=0.83):
    """
    修正台组内排名页（第 4 页，slide_index=3）：

    1. CCTV-17（农业农村）行设置绿色字体（5B7F5B）
    2. 橘色阈值线（直接连接符 10）定位到右侧排名中市场份额
       跨越阈值的行边界，标签（文本框 10）恢复阈值文本

    :param slide: 目标幻灯片
    :param threshold_color: 阈值线/标签颜色（从模板配置文件读取），如 'ED7D31'
    :param threshold_text: 阈值标签文本，如 '0.83%'
    :param threshold_value: 阈值数值，如 0.83
    """
    GREEN_FONT = '5B7F5B'

    table_shape = None
    connector = None
    label_box = None

    for shp in slide.shapes:
        if getattr(shp, 'has_table', False):
            table_shape = shp
        elif shp.name == '直接连接符 10':
            connector = shp
        elif shp.name == '文本框 10':
            label_box = shp

    if not table_shape:
        return

    tbl = table_shape.table
    n_rows = len(tbl.rows)
    n_cols = len(tbl.columns)

    # ── 查找含 CCTV-17 的行 和 阈值跨越行 ──
    cctv17_all_rows = set()     # 任意列含 CCTV-17 的行
    threshold_row = -1          # 右侧排名中第一个 < 阈值的行

    for r in range(2, n_rows):  # 跳过日期 + 列头
        for c in range(n_cols):
            cell_text = tbl.cell(r, c).text
            if '农业农村' in cell_text or 'CCTV-17' in cell_text:
                cctv17_all_rows.add(r)
        # 右侧份额列（col 5）：找第一个 < 阈值的行
        if threshold_row < 0:
            try:
                share_val = float(tbl.cell(r, 5).text)
                if share_val < threshold_value:
                    threshold_row = r
            except (ValueError, TypeError):
                pass

    # ── 0. 清除非 CCTV-17 数据行的残留绿色/加粗（demo 模板中旧位置） ──
    for r in range(2, n_rows):
        if r in cctv17_all_rows:
            continue
        for c in range(n_cols):
            tc = tbl.cell(r, c)._tc
            for run in tc.findall(f'.//{{{A_NS}}}r'):
                rpr = run.find(f'{{{A_NS}}}rPr')
                if rpr is None:
                    continue
                # 移除绿色 solidFill
                for sf in rpr.findall(f'{{{A_NS}}}solidFill'):
                    clr = sf.find(f'{{{A_NS}}}srgbClr')
                    if clr is not None and clr.get('val') == GREEN_FONT:
                        rpr.remove(sf)
                # 移除加粗（恢复为非粗体）
                if rpr.get('b') == '1':
                    rpr.attrib.pop('b', None)
            for eprpr in tc.findall(f'.//{{{A_NS}}}endParaRPr'):
                for sf in eprpr.findall(f'{{{A_NS}}}solidFill'):
                    clr = sf.find(f'{{{A_NS}}}srgbClr')
                    if clr is not None and clr.get('val') == GREEN_FONT:
                        eprpr.remove(sf)
                if eprpr.get('b') == '1':
                    eprpr.attrib.pop('b', None)

    # ── 1. 绿色字体 ──
    for r in cctv17_all_rows:
        for c in range(n_cols):
            tc = tbl.cell(r, c)._tc
            # 处理所有 <a:r> 的 rPr：设置绿色 + 加粗
            for run in tc.findall(f'.//{{{A_NS}}}r'):
                rpr = run.find(f'{{{A_NS}}}rPr')
                if rpr is None:
                    rpr = etree.SubElement(run, f'{{{A_NS}}}rPr')
                    t_elem = run.find(f'{{{A_NS}}}t')
                    if t_elem is not None:
                        run.remove(rpr)
                        run.insert(list(run).index(t_elem), rpr)
                # 移除旧 solidFill，添加绿色
                for old in rpr.findall(f'{{{A_NS}}}solidFill'):
                    rpr.remove(old)
                fill = etree.SubElement(rpr, f'{{{A_NS}}}solidFill')
                clr = etree.SubElement(fill, f'{{{A_NS}}}srgbClr')
                clr.set('val', GREEN_FONT)
                # 加粗
                rpr.set('b', '1')

            # 处理 defRPr：移除冲突颜色
            for defrpr in tc.findall(f'.//{{{A_NS}}}defRPr'):
                for old in defrpr.findall(f'{{{A_NS}}}solidFill'):
                    defrpr.remove(old)

            # 处理 endParaRPr：设置绿色 + 加粗
            for eprpr in tc.findall(f'.//{{{A_NS}}}endParaRPr'):
                for old in eprpr.findall(f'{{{A_NS}}}solidFill'):
                    eprpr.remove(old)
                fill = etree.SubElement(eprpr, f'{{{A_NS}}}solidFill')
                clr = etree.SubElement(fill, f'{{{A_NS}}}srgbClr')
                clr.set('val', GREEN_FONT)
                eprpr.set('b', '1')

    # ── 2. 重新定位橘色阈值线和标签 ──
    if threshold_row < 0 or not connector:
        return

    # 从表格 XML 累加行高计算阈值跨越行的上沿 Y 坐标
    tbl_xml = table_shape._element.find(f'.//{{{A_NS}}}tbl')
    if tbl_xml is None:
        return

    trs = tbl_xml.findall(f'{{{A_NS}}}tr')
    cum_y = table_shape.top
    for idx in range(threshold_row):
        cum_y += int(trs[idx].get('h', 0))

    target_y = cum_y  # 阈值跨越行上沿（线放在此处）

    # 更新连接符位置（仅调整 Y，X / 宽度不变）
    connector.top = target_y

    # 更新标签文本框位置（保持在线上方的固定偏移）+ 恢复阈值文本
    if label_box:
        LABEL_ABOVE_LINE = 137093   # demo 中标签比线高约 137 kEMU
        label_box.top = target_y - LABEL_ABOVE_LINE
        # 恢复阈值标签文本（text sync 可能把 demo 的值覆盖了）
        if label_box.has_text_frame:
            for para in label_box.text_frame.paragraphs:
                for run in para.runs:
                    run.text = threshold_text
                    break
                break


def _fix_channel_ranking_page(slide):
    """
    修正上星频道排名页（第 5 页，slide_index=4）：

    1. 箭头列（col 6）：↑ 标红（FF0000），↓ 保持黑色
    2. CCTV-17（农业农村）行：全行字体设为红色（FF0000）+ 加粗
    """
    RED = 'FF0000'

    for shp in slide.shapes:
        if not getattr(shp, 'has_table', False):
            continue
        tbl = shp.table
        n_rows = len(tbl.rows)
        n_cols = len(tbl.columns)

        # ── 查找 CCTV-17 所在行 ──
        cctv17_rows = set()
        for r in range(2, n_rows):
            for c in range(n_cols):
                t = tbl.cell(r, c).text
                if '农业农村' in t or 'CCTV-17' in t:
                    cctv17_rows.add(r)

        for r in range(2, n_rows):
            # ── 箭头列上色 ──
            if n_cols >= 7:
                arrow_text = tbl.cell(r, 6).text.strip()
                if arrow_text == '↑':
                    _set_cell_font_color(tbl.cell(r, 6), RED)
                elif arrow_text == '↓':
                    _set_cell_font_color(tbl.cell(r, 6), '000000')

            # ── CCTV-17 行：清除继承自 demo 的红色，恢复为黑色 ──
            is_cctv17 = any('农业农村' in tbl.cell(r, c).text or 'CCTV-17' in tbl.cell(r, c).text
                            for c in range(n_cols))
            if is_cctv17:
                for c in range(min(6, n_cols)):  # col 0-5，箭头列已单独处理
                    _set_cell_font_color(tbl.cell(r, c), '000000')


def _set_cell_font_color(cell, color_hex, bold=None):
    """设置单元格所有 run / endParaRPr 的字体颜色，可选加粗"""
    tc = cell._tc
    for run in tc.findall(f'.//{{{A_NS}}}r'):
        rpr = run.find(f'{{{A_NS}}}rPr')
        if rpr is None:
            rpr = etree.SubElement(run, f'{{{A_NS}}}rPr')
            t_elem = run.find(f'{{{A_NS}}}t')
            if t_elem is not None:
                run.remove(rpr)
                run.insert(list(run).index(t_elem), rpr)
        for old in rpr.findall(f'{{{A_NS}}}solidFill'):
            rpr.remove(old)
        fill = etree.SubElement(rpr, f'{{{A_NS}}}solidFill')
        clr = etree.SubElement(fill, f'{{{A_NS}}}srgbClr')
        clr.set('val', color_hex)
        if bold is not None:
            rpr.set('b', '1' if bold else '0')

    # defRPr：移除冲突颜色
    for defrpr in tc.findall(f'.//{{{A_NS}}}defRPr'):
        for old in defrpr.findall(f'{{{A_NS}}}solidFill'):
            defrpr.remove(old)

    # endParaRPr
    for eprpr in tc.findall(f'.//{{{A_NS}}}endParaRPr'):
        for old in eprpr.findall(f'{{{A_NS}}}solidFill'):
            eprpr.remove(old)
        fill = etree.SubElement(eprpr, f'{{{A_NS}}}solidFill')
        clr = etree.SubElement(fill, f'{{{A_NS}}}srgbClr')
        clr.set('val', color_hex)
        if bold is not None:
            eprpr.set('b', '1' if bold else '0')


def _fix_chart_max_annotation(target_slide, annotation_name=None, num_format=':.3f'):
    """
    通用：更新图表超轴标注文本框。

    找到柱状图（barChart）的所有超出坐标轴最大值的数据点，
    按从左到右顺序匹配幻灯片上已有的标注文本框，写入对应值。
    多余的文本框清空。

    annotation_name: 首选标注文本框名称（如 '文本框 2'），用于仅处理一个超轴值。
                     若为 None，则自动搜索所有候选文本框。
    num_format: 数字格式字符串，如 ':.3f' 或 ':.4f'。
    """
    C_NS = 'http://schemas.openxmlformats.org/drawingml/2006/chart'

    chart_shape = None
    for shp in target_slide.shapes:
        if shp.has_chart:
            chart_shape = shp
            break
    if not chart_shape:
        return

    cs = chart_shape.chart._chartSpace
    pa = cs.find(f'.//{{{C_NS}}}plotArea')
    bar_chart = pa.find(f'{{{C_NS}}}barChart')
    if bar_chart is None:
        return

    # 获取 barChart 绑定的值轴最大值
    axis_max = None
    bar_ax_ids = [ax.get('val') for ax in bar_chart.findall(f'{{{C_NS}}}axId')]
    for vax in pa.findall(f'{{{C_NS}}}valAx'):
        ax_id = vax.find(f'{{{C_NS}}}axId').get('val')
        if ax_id in bar_ax_ids:
            scl = vax.find(f'{{{C_NS}}}scaling')
            mx = scl.find(f'{{{C_NS}}}max') if scl is not None else None
            if mx is not None:
                axis_max = float(mx.get('val'))
            break

    # 收集所有超出 axis_max 的数据点 (index, value)
    overflow_bars = []
    bar_ser = bar_chart.find(f'{{{C_NS}}}ser')
    if bar_ser is not None and axis_max is not None:
        val_elem = bar_ser.find(f'{{{C_NS}}}val')
        if val_elem is not None:
            num_ref = val_elem.find(f'{{{C_NS}}}numRef')
            if num_ref is not None:
                nc = num_ref.find(f'{{{C_NS}}}numCache')
                if nc is not None:
                    for pt in nc.findall(f'{{{C_NS}}}pt'):
                        idx = int(pt.get('idx', '-1'))
                        v = pt.find(f'{{{C_NS}}}v')
                        if v is not None and v.text:
                            try:
                                fv = float(v.text)
                                if fv > axis_max:
                                    overflow_bars.append((idx, fv))
                            except ValueError:
                                pass

    # 按 index 排序（对应从左到右的柱子）
    overflow_bars.sort(key=lambda x: x[0])

    # 收集候选标注文本框（排除标题、备注等），按 left 排序
    EXCLUDE_KW = ['频道', '收视', '市场', '，', '备注', '栏目', '排名']
    candidate_boxes = []
    for shp in target_slide.shapes:
        if not (shp.has_text_frame and not shp.has_chart and not shp.has_table):
            continue
        name = shp.name or ''
        if '标题' in name:
            continue
        if '文本框' not in name:
            continue
        # 通过备注等文字内容排除
        txt = shp.text_frame.text.strip()
        if any(kw in txt for kw in EXCLUDE_KW):
            continue
        candidate_boxes.append(shp)

    # 按 left 位置排序
    candidate_boxes.sort(key=lambda s: s.left)

    fmt_str = f'{{{num_format}}}'

    # ── 计算每个超轴柱子对应的标注位置 ──
    # 需要图表 shape 的位置和数据点数
    n_pts = 0
    if bar_ser is not None:
        val_elem = bar_ser.find(f'{{{C_NS}}}val')
        if val_elem is not None:
            nr = val_elem.find(f'{{{C_NS}}}numRef')
            if nr is not None:
                nc = nr.find(f'{{{C_NS}}}numCache')
                if nc is not None:
                    ptc = nc.find(f'{{{C_NS}}}ptCount')
                    if ptc is not None:
                        n_pts = int(ptc.get('val', '0'))

    # 参考第一个文本框（已有正确位置）来推算绘图区参数
    ref_box = candidate_boxes[0] if candidate_boxes else None

    # 依次将超轴值写入对应文本框，并调整位置
    for i, box in enumerate(candidate_boxes):
        if i < len(overflow_bars):
            bar_idx, val = overflow_bars[i]
            text = fmt_str.format(val)

            # 动态调整文本框位置 —— 对齐到对应柱子上方
            if ref_box is not None and n_pts > 0 and chart_shape is not None:
                if i == 0:
                    # 第一个文本框保持原位（模板已对齐好）
                    pass
                else:
                    ref_idx = overflow_bars[0][0]
                    ref_val = overflow_bars[0][1]
                    # 用 ref_box 的中心点反推柱宽：ref_box 对齐 bar[ref_idx]
                    ref_center_x = ref_box.left + ref_box.width // 2
                    chart_w = chart_shape.width
                    # 柱槽宽 ≈ 绘图区宽 / 柱数，绘图区约占 chart 宽 93%
                    bar_slot = int(chart_w * 0.93) / n_pts

                    # bar[bar_idx] 中心 = ref_center + (bar_idx - ref_idx) * bar_slot
                    bar_center_x = ref_center_x + int((bar_idx - ref_idx) * bar_slot)
                    box.left = bar_center_x - box.width // 2

                    # 垂直位置：根据溢出比例计算
                    ref_bottom = ref_box.top + ref_box.height
                    overflow_ratio = (val - axis_max) / (ref_val - axis_max) if ref_val > axis_max else 0.5
                    annot_height = int(overflow_ratio * ref_box.height)
                    annot_height = max(annot_height, box.height)
                    box.top = ref_bottom - annot_height
                    box.height = annot_height
        else:
            text = ''
        tf = box.text_frame
        if tf.paragraphs:
            para = tf.paragraphs[0]
            if para.runs:
                para.runs[0].text = text
            elif text:
                # 段落里没有 run —— 从第一个有 run 的文本框复制格式新建
                from pptx.oxml.ns import qn as _qn
                from copy import deepcopy
                # 找一个同页有 run 的标注文本框作为样板
                ref_para_elem = None
                ref_run_elem = None
                for cb in candidate_boxes:
                    for rp in cb.text_frame.paragraphs:
                        if rp.runs:
                            ref_para_elem = rp._p
                            ref_run_elem = rp.runs[0]._r
                            break
                    if ref_run_elem is not None:
                        break
                if ref_run_elem is not None:
                    # 复制段落属性 (pPr) 以继承 defRPr（字号、粗体等）
                    ref_pPr = ref_para_elem.find(_qn('a:pPr'))
                    if ref_pPr is not None and para._p.find(_qn('a:pPr')) is None:
                        para._p.insert(0, deepcopy(ref_pPr))
                    new_r = deepcopy(ref_run_elem)
                    new_r.find(_qn('a:t')).text = text
                    para._p.append(new_r)
                else:
                    para.text = text


def _fix_schedule_chart_page(target_slide, src_slide):
    """
    修正串单市场份额页（第 6 页，slide_index=5）：

    1. 组合图表（barChart + lineChart）的数据同步：不使用 replace_data（会破坏双图结构），
       而是 XML 级别逐 series 替换 strCache / numCache。
    2. 超轴最大值标注文本框。
    """
    # ── 1. 找到 target 和 src 的图表 shape ──
    target_chart_shape = None
    src_chart_shape = None
    for shp in target_slide.shapes:
        if shp.has_chart:
            target_chart_shape = shp
            break
    for shp in src_slide.shapes:
        if shp.has_chart:
            src_chart_shape = shp
            break

    if target_chart_shape and src_chart_shape:
        _sync_combo_chart(target_chart_shape, src_chart_shape)

    # ── 2. 超轴最大值标注 ──
    _fix_chart_max_annotation(target_slide, annotation_name='文本框 23', num_format=':.3f')


# ═══════════════════════════════════════════════════════════════
# 分分钟页覆盖表格修正（第 9、10 页，slide_index=8,9）
# ═══════════════════════════════════════════════════════════════

def _parse_time_minutes(t):
    """HH:MM → 分钟数。支持 '1900-01-01 HH:MM:SS' 格式和 24+ 小时制"""
    t = str(t).strip()
    if ' ' in t:
        t = t.split(' ')[1]  # 取时间部分
    parts = t.replace('.', ':').split(':')
    h, m = int(parts[0]), int(parts[1]) if len(parts) > 1 else 0
    if h < 5:
        h += 24  # 次日凌晨归入 24+
    return h * 60 + m


def _group_programs(programs, max_cols=20):
    """
    将串单节目合并为 ≤ max_cols 个显示组。

    每个节目按"开始时间"和"结束时间"确定其在时间轴上的位置与宽度。
    首播/重播判定：下午剧场~晚间节目时段 (13:30-22:30) 的节目标记为（首播），
    其余时段（早间、上午、午间、夜间）标记为（重播）。
    连续播出的同名节目（如电视剧连续几集）自动合并为一个单元格。

    返回: list of dict {'name': str, 'start': int, 'end': int, 'duration': int, 'premiere': bool}
    """
    # 不显示的节目名称
    _HIDDEN_NAMES = {'国歌', '歌曲', '再见'}

    valid = []
    for p in programs:
        if p.duration < 1:
            continue
        if p.name.strip() in _HIDDEN_NAMES:
            continue
        start_m = _parse_time_minutes(p.start_time)
        # 排除凌晨 02:00~05:29
        if 2 * 60 <= start_m < 5 * 60 + 30:
            continue
        end_m = _parse_time_minutes(p.end_time)

        # 首播/重播判定：下午剧场(13:30)起到晚间节目(22:30)止为首播，其余为重播
        is_premiere = (13 * 60 + 30 <= start_m < 22 * 60 + 30)
        suffix = '（首播）' if is_premiere else '（重播）'

        # 不加"电视剧："前缀，只保留节目名
        display_name = f'{p.name}{suffix}'

        valid.append({
            'name': display_name,
            'raw_name': p.name,
            'start': start_m,
            'end': end_m,
            'duration': max(end_m - start_m, 1),
            'premiere': is_premiere,
        })

    if not valid:
        return valid

    # 合并连续播出的同名节目（如电视剧连续几集）
    consecutive = [dict(valid[0])]
    for v in valid[1:]:
        prev = consecutive[-1]
        if v['raw_name'] == prev['raw_name'] and v['premiere'] == prev['premiere']:
            prev['end'] = v['end']
            prev['duration'] = prev['end'] - prev['start']
        else:
            consecutive.append(dict(v))
    valid = consecutive

    # 合并短节目 (< 5 min) 到相邻节目
    _MAX_MERGE_LEN = 20   # 合并后名称上限（防止垂直文本溢出重叠）
    merged = []
    i = 0
    while i < len(valid):
        grp = dict(valid[i])
        while grp['duration'] < 5 and i + 1 < len(valid):
            i += 1
            nxt = valid[i]
            if len(grp['name']) + len(nxt['name']) < _MAX_MERGE_LEN:
                grp['name'] = grp['name'] + '+' + nxt['name']
            grp['end'] = nxt['end']
            grp['duration'] = grp['end'] - grp['start']
        merged.append(grp)
        i += 1

    # 如果还是太多列，反复合并最短的
    while len(merged) > max_cols:
        min_idx = min(range(len(merged)), key=lambda k: merged[k]['duration'])
        if min_idx < len(merged) - 1:
            nb = min_idx + 1
        else:
            nb = min_idx - 1
        a, b = sorted([min_idx, nb])
        combined_name = merged[a]['name']
        if len(combined_name) + len(merged[b]['name']) < _MAX_MERGE_LEN:
            combined_name += '+' + merged[b]['name']
        combined = {
            'name': combined_name,
            'start': merged[a]['start'],
            'end': merged[b]['end'],
            'duration': merged[b]['end'] - merged[a]['start'],
        }
        merged = merged[:a] + [combined] + merged[b + 1:]

    # 截断过长名称：合并后仍超长时只保留第一个节目名
    for g in merged:
        if len(g['name']) > _MAX_MERGE_LEN and '+' in g['name']:
            g['name'] = g['name'].split('+')[0]

    return merged


def _proportional_widths(durations, total_width):
    """按时长占比分配列宽（EMU），确保每列至少有最小宽度"""
    total_dur = sum(durations)
    if total_dur <= 0:
        n = len(durations)
        return [total_width // n] * n
    widths = [max(int(d / total_dur * total_width), 50000) for d in durations]
    # 消除舍入误差
    diff = total_width - sum(widths)
    if widths:
        widths[-1] += diff
    return widths


# 固定时段划分（分钟数）— 与频道编排一致
_FIXED_TIME_SLOTS = [
    ('早间节目',  6 * 60,       10 * 60),       # 06:00-10:00  240min
    ('上午剧场', 10 * 60,       12 * 60),       # 10:00-12:00  120min
    ('午间节目', 12 * 60,       13 * 60 + 30),  # 12:00-13:30   90min
    ('下午剧场', 13 * 60 + 30,  17 * 60),       # 13:30-17:00  210min
    ('傍晚节目', 17 * 60,       18 * 60 + 30),  # 17:00-18:30   90min
    ('黄金剧场', 18 * 60 + 30,  20 * 60 + 30),  # 18:30-20:30  120min
    ('晚间节目', 20 * 60 + 30,  22 * 60 + 30),  # 20:30-22:30  120min
    ('夜间节目', 22 * 60 + 30,  25 * 60),       # 22:30-25:00  150min
]


def _fix_minute_chart_page(target_slide, data, metric='rating'):
    """
    修正分分钟收视率/市场份额页的 3 个覆盖表格：

    1. 表头表格 (8 cols): 固定时段名称（早间节目…夜间节目）+ 按固定时段比例列宽
    2. 分隔线表格 (8 cols, 空): 同上比例列宽（绿色竖线对齐时段边界）
    3. 节目名称表格 (20/21 cols): 串单节目名 + 按实际开始/结束时间比例列宽
    """
    programs = data.programs
    if not programs:
        return

    # ── 固定时段列宽 ──
    slot_durations = [end - start for _, start, end in _FIXED_TIME_SLOTS]

    # ── 节目分组 ──
    prog_groups = _group_programs(programs, max_cols=20)
    prog_durations = [g['duration'] for g in prog_groups]

    # ── 找到 3 个表格并分类 ──
    header_tbl_shape = None   # 8 cols, 有文字（表头）
    divider_tbl_shape = None  # 8 cols, 空（分隔线）
    prog_tbl_shape = None     # 20/21 cols（节目名）

    for shp in target_slide.shapes:
        if not getattr(shp, 'has_table', False):
            continue
        tbl = shp.table
        ncols = len(tbl.columns)
        if ncols in (20, 21):
            prog_tbl_shape = shp
        elif ncols == 8:
            cell_text = tbl.rows[0].cells[0].text.strip()
            if cell_text:
                header_tbl_shape = shp
            else:
                divider_tbl_shape = shp

    a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    # ── 更新表头文字（固定时段名称）──
    if header_tbl_shape:
        tbl = header_tbl_shape.table
        n = min(len(tbl.columns), len(_FIXED_TIME_SLOTS))
        for ci in range(n):
            new_text = _FIXED_TIME_SLOTS[ci][0]  # 早间节目 / 上午剧场 / …
            tc_elem = tbl.cell(0, ci)._tc
            tc_body = tc_elem.find(f'.//{{{a_ns}}}txBody')
            if tc_body is not None:
                paras = tc_body.findall(f'{{{a_ns}}}p')
                if paras:
                    _replace_para_text(paras[0], new_text)
                    for p in paras[1:]:
                        tc_body.remove(p)
            else:
                tbl.cell(0, ci).text = new_text

    # ── 获取图表绘图区域用于对齐 ──
    chart_shape = None
    for shp in target_slide.shapes:
        if getattr(shp, 'has_chart', False):
            chart_shape = shp
            break

    chart_start_min = 330   # 05:30
    chart_end_min = 1560    # 26:00
    plot_left = None
    plot_width = None

    if chart_shape:
        chart = chart_shape.chart
        cats = list(chart.plots[0].categories)
        if cats:
            chart_start_min = _parse_time_minutes(cats[0])
            chart_end_min = _parse_time_minutes(cats[-1]) + 1
        # 从图表 XML 解析 plotArea 的 x/w
        chart_xml = etree.tostring(chart._chartSpace).decode()
        layout_match = re.search(r'<c:plotArea>(.*?)</c:plotArea>', chart_xml, re.DOTALL)
        if layout_match:
            x_m = re.search(r'<c:x val="([^"]+)"', layout_match.group(1))
            w_m = re.search(r'<c:w val="([^"]+)"', layout_match.group(1))
            if x_m and w_m:
                pa_x = float(x_m.group(1))
                pa_w = float(w_m.group(1))
                plot_left = chart_shape.left + int(pa_x * chart_shape.width)
                plot_width = int(pa_w * chart_shape.width)

    chart_duration = chart_end_min - chart_start_min

    # 表头/分隔线表格：保留模板原始位置和列宽，不做重新定位
    # （模板已手动调好对齐，重新计算会导致列太窄、文字换行）

    # ── 更新节目名称表格 ──
    if prog_tbl_shape and prog_groups:
        tbl = prog_tbl_shape.table
        ncols = len(tbl.columns)

        # 使节目段连续覆盖整个图表时间范围
        for i in range(len(prog_groups) - 1):
            prog_groups[i]['end'] = prog_groups[i + 1]['start']
        prog_groups[0]['start'] = chart_start_min
        prog_groups[-1]['end'] = chart_end_min
        for g in prog_groups:
            g['duration'] = g['end'] - g['start']
        prog_durations = [g['duration'] for g in prog_groups]

        # 定位表格以对齐图表绘图区
        if plot_left is not None:
            prog_tbl_shape.left = plot_left
            prog_total_w = plot_width
        else:
            prog_total_w = prog_tbl_shape.width

        # 判断是否有前导空列（21 列 = 1 spacer + 20 program）
        has_spacer = (ncols == 21)
        data_start = 1 if has_spacer else 0
        data_cols = ncols - data_start

        if has_spacer:
            tbl.columns[0].width = 0

        n = min(data_cols, len(prog_groups))
        widths = _proportional_widths(prog_durations[:n], prog_total_w)

        # 首播节目颜色：收视率页=橙色, 市场份额页=绿色
        _PREMIERE_COLOR = 'F39C12' if metric == 'rating' else '4A7C31'

        for ci in range(n):
            col_idx = data_start + ci
            tbl.columns[col_idx].width = widths[ci]
            grp = prog_groups[ci]
            new_text = grp['name']
            is_premiere = grp.get('premiere', False)
            tc_elem = tbl.cell(0, col_idx)._tc
            tc_body = tc_elem.find(f'.//{{{a_ns}}}txBody')
            if tc_body is not None:
                paras = tc_body.findall(f'{{{a_ns}}}p')
                if paras:
                    _replace_para_text(paras[0], new_text)
                    # 居中对齐
                    pPr = paras[0].find(f'{{{a_ns}}}pPr')
                    if pPr is None:
                        pPr = etree.SubElement(paras[0], f'{{{a_ns}}}pPr')
                        paras[0].insert(0, pPr)
                    pPr.set('algn', 'ctr')
                    # 首播节目上色
                    if is_premiere:
                        for rr in paras[0].findall(f'{{{a_ns}}}r'):
                            rPr = rr.find(f'{{{a_ns}}}rPr')
                            if rPr is not None:
                                for old_f in rPr.findall(f'{{{a_ns}}}solidFill'):
                                    rPr.remove(old_f)
                                sf = etree.SubElement(rPr, f'{{{a_ns}}}solidFill')
                                sc = etree.SubElement(sf, f'{{{a_ns}}}srgbClr')
                                sc.set('val', _PREMIERE_COLOR)
                                rPr.insert(0, sf)
                    for p in paras[1:]:
                        tc_body.remove(p)
            else:
                tbl.cell(0, col_idx).text = new_text

        # 多余列清空
        for ci in range(n, data_cols):
            col_idx = data_start + ci
            tbl.columns[col_idx].width = 0
            tbl.cell(0, col_idx).text = ''


def _sync_slide(target_slide, src_slide, slide_index, total_slides):
    """同步单页内容"""
    # 第 1 页封面：精确替换日期和期号
    if slide_index == 0:
        _update_cover(target_slide, src_slide)
        return

    # 最后一页感谢观看：不做任何替换
    if slide_index == total_slides - 1:
        return

    # 第 13 页（idx=12）频道分类观众规模：图表同步 + 变化标签动态定位
    if slide_index == 12:
        _sync_and_position_audience(target_slide, src_slide)
        return

    # 其他页：按类型匹配并同步
    target_texts = _iter_shapes_by_type(target_slide, 'text')
    src_texts = _iter_shapes_by_type(src_slide, 'text')

    target_charts = _iter_shapes_by_type(target_slide, 'chart')
    src_charts = _iter_shapes_by_type(src_slide, 'chart')

    target_tables = _iter_shapes_by_type(target_slide, 'table')
    src_tables = _iter_shapes_by_type(src_slide, 'table')

    # 文本框：最近邻匹配 + XML 级文本替换
    src_text_map = {idx: shp for idx, shp in src_texts}
    for ti, si in _match_nearest(target_texts, src_texts):
        t_shp = dict(target_texts)[ti]
        s_shp = src_text_map[si]
        _sync_text_shape(t_shp, s_shp)

    # 第 6 页（idx=5）串单市场份额：组合图表 + 超轴标注
    if slide_index == 5:
        _fix_schedule_chart_page(target_slide, src_slide)

    # 图表：Y 排序匹配 + 数据替换（第 6 页走 combo chart 专用逻辑，跳过通用同步）
    if slide_index != 5:
        src_chart_map = {idx: shp for idx, shp in src_charts}
        for ti, si in _match_charts_by_y_order(target_charts, src_charts):
            t_shp = dict(target_charts)[ti]
            s_shp = src_chart_map[si]
            _sync_chart(t_shp, s_shp)

    # 表格：最近邻匹配 + 单元格文本替换
    # 第 9/10 页（idx=8,9）分分钟页覆盖表格由专用函数处理，跳过通用表格同步
    if slide_index not in (8, 9):
        src_table_map = {idx: shp for idx, shp in src_tables}
        for ti, si in _match_nearest(target_tables, src_tables):
            t_shp = dict(target_tables)[ti]
            s_shp = src_table_map[si]
            _sync_table(t_shp, s_shp)

    # 第 3 页（idx=2）市场份额：修正变化百分比颜色（提升=红，下降=绿）+ 数据标签统一3位小数
    if slide_index == 2:
        _fix_change_pct_colors(target_slide)
        _fix_page3_numfmt(target_slide)

    # 第 4 页（idx=3）台组排名：CCTV-17 刷绿 + 橘色阈值线定位
    if slide_index == 3:
        t_color, t_text, t_value = _load_threshold_config()
        _fix_org_ranking_page(target_slide, threshold_color=t_color,
                              threshold_text=t_text, threshold_value=t_value)

    # 第 5 页（idx=4）上星频道排名：箭头上色 + CCTV-17 行标红
    if slide_index == 4:
        _fix_channel_ranking_page(target_slide)

    # 第 7 页（idx=6）栏目收视率排名：超轴最大值标注
    if slide_index == 6:
        _fix_chart_max_annotation(target_slide, annotation_name='文本框 2', num_format=':.4f')

    # 第 8 页（idx=7）栏目收视份额排名：超轴最大值标注
    if slide_index == 7:
        _fix_chart_max_annotation(target_slide, annotation_name='文本框 2', num_format=':.3f')

    # 第 11/12 页（idx=10,11）栏目首播收视率/市场份额：修正当日系列颜色为橙色
    if slide_index in (10, 11):
        _fix_premiere_chart_colors(target_slide)


def _fix_premiere_chart_colors(target_slide):
    """
    修正栏目首播页图表的系列颜色。

    问题：demo 模板有 3 个系列(前1个月均值/前一日/当日)，当只有 2 个系列时
    （无月均数据），replace_data 后 series[1](当日) 继承了 demo 的 series[1]
    的浅灰色，应改为橙色(FE9B1C)以匹配 demo 的当日系列。
    """
    C_NS = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
    A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    for shp in target_slide.shapes:
        if not getattr(shp, 'has_chart', False):
            continue
        chart = shp.chart
        n_series = len(list(chart.series))
        if n_series < 2:
            continue

        # 修正最后一个系列（当日数据）为橙色 FE9B1C
        last_ser = list(chart.series)[-1]._element
        sp = last_ser.find(f'{{{C_NS}}}spPr')
        if sp is None:
            sp = etree.SubElement(last_ser, f'{{{C_NS}}}spPr')
        # 清除旧 solidFill
        for old_fill in sp.findall(f'{{{A_NS}}}solidFill'):
            sp.remove(old_fill)
        # 添加橙色 solidFill
        solid = etree.SubElement(sp, f'{{{A_NS}}}solidFill')
        srgb = etree.SubElement(solid, f'{{{A_NS}}}srgbClr')
        srgb.set('val', 'FE9B1C')
        # 将 solidFill 放到 spPr 的第一个子元素位置
        sp.insert(0, solid)


# ═══════════════════════════════════════════════════════════════
# 主入口
# ═══════════════════════════════════════════════════════════════

def _make_data_draft(excel_path, template_path):
    """用现有生成器产出数据稿（临时文件）"""
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
    tmp.close()
    base.generate_report(excel_path, template_path, tmp.name)
    return tmp.name


def _strip_comments(pptx_path):
    """
    从 pptx 文件中移除所有批注（comments）和批注作者信息。

    操作：ZIP 级别删除 comment*.xml / commentAuthors.xml，
    并清理 [Content_Types].xml 和 .rels 中的引用。
    """
    import zipfile
    tmp_path = pptx_path + '.tmp'
    with zipfile.ZipFile(pptx_path, 'r') as zin, \
         zipfile.ZipFile(tmp_path, 'w', zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            name_lower = item.filename.lower()
            # 跳过批注文件
            if 'commentauthors' in name_lower or '/comments/' in name_lower:
                continue
            data = zin.read(item.filename)
            # 清理 Content_Types.xml 中的批注引用
            if item.filename == '[Content_Types].xml':
                text = data.decode('utf-8')
                text = re.sub(
                    r'<Override[^>]*PartName="[^"]*comment[^"]*"[^>]*/>\s*',
                    '', text, flags=re.IGNORECASE)
                data = text.encode('utf-8')
            # 清理 .rels 文件中的批注关系引用
            if item.filename.endswith('.rels'):
                text = data.decode('utf-8')
                if 'comment' in text.lower():
                    text = re.sub(
                        r'<Relationship[^>]*Target="[^"]*comment[^"]*"[^>]*/>\s*',
                        '', text, flags=re.IGNORECASE)
                    data = text.encode('utf-8')
            zout.writestr(item, data)
    # 替换原文件
    os.replace(tmp_path, pptx_path)


def generate_report_config_driven(excel_path, template_path, output_path):
    """
    配置驱动生成：

    1. 生成数据稿（临时）
    2. 复制 demo → 输出文件
    3. XML 级逐页回填数据
    """
    if not template_path or not Path(template_path).exists():
        raise ValueError(f'模板不存在: {template_path}')

    print(f'[1/4] 生成数据稿: {excel_path}')
    draft_path = _make_data_draft(excel_path, template_path)

    print(f'[2/4] 复制模板: {template_path} → {output_path}')
    shutil.copy2(template_path, output_path)

    print('[3/5] XML 级精确回填...')
    target = Presentation(output_path)
    source = Presentation(draft_path)

    n = min(len(target.slides), len(source.slides))
    for i in range(n):
        print(f'  - 第 {i + 1}/{n} 页')
        _sync_slide(target.slides[i], source.slides[i], i, n)

    # 修复标题（简版模式下 data draft 不生成标题）
    print('[4/5] 修复标题...')
    from data_reader import read_excel_data
    data = read_excel_data(excel_path)
    # 确保日期正确（与 generate_report.py 同逻辑）
    ymd_match = re.search(r'(\d{4})(\d{2})(\d{2})', os.path.basename(excel_path))
    if ymd_match:
        year, month, day = int(ymd_match.group(1)), int(ymd_match.group(2)), int(ymd_match.group(3))
        data.report_date_short = f'{month}月{day}日'
        data.report_date = f'{year}年{month}月{day}日'
    _fix_titles(target, data)

    # 修复分分钟页覆盖表格（第 9/10 页）
    n = len(target.slides)
    if n > 8:
        print('  - 修复第 9 页分分钟收视率覆盖表格')
        _fix_minute_chart_page(target.slides[8], data, metric='rating')
    if n > 9:
        print('  - 修复第 10 页分分钟市场份额覆盖表格')
        _fix_minute_chart_page(target.slides[9], data, metric='share')

    print(f'[5/5] 保存: {output_path}')
    target.save(output_path)

    # 清除从 demo 模板继承的批注（commentAuthors / comments）
    _strip_comments(output_path)

    print(f'✅ 完成！共 {len(target.slides)} 张幻灯片')


def main():
    parser = argparse.ArgumentParser(description='全页配置驱动 PPT 生成器 v3')
    parser.add_argument('excel', help='Excel 数据文件路径')
    parser.add_argument('--template', '-t', required=True, help='demo 模板 pptx 路径')
    parser.add_argument('--output', '-o', required=True, help='输出 PPT 路径')
    args = parser.parse_args()
    generate_report_config_driven(args.excel, args.template, args.output)


if __name__ == '__main__':
    main()
