import sys

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
import xlrd
import pandas as pd
import numpy as np

date = '0000'


def str_date(date):
    return f'{date[:2]}月{date[2:4]}日'


def str_p_date(date):
    day = int(date[2:])
    if 1 < day < 11:
        return f'{date[:2]}月0{day - 1}日'
    elif day >= 11:
        return f'{date[:2]}月{day - 1}日'
    elif day == 1:
        if date[:2] in ['02', '04', '06', '07', '09', '11']:
            return f'0{int(date[:2]) - 1}月31日'
        elif date[:2] in ['12', '05', '08', '10']:
            return f'0{int(date[:2]) - 1}月30日'
        if date[:2] == '01':
            return '12月31日'
        elif date[:2] == '03':
            return '02月28日'


def increase_rank(p, n):
    if n < p:
        return f'提升{p - n}位'
    elif n > p:
        return f'下滑{n - p}位'
    else:
        return '保持不变'


def increase_percent(p, n):
    if n < p:
        return f'下滑{round(100 * (p - n) / p, 1)}%'
    elif n > p:
        return f'提升{round(100 * (n - p) / p, 1)}%'
    else:
        return '保持不变'


def chart_title_text(shapes, shape, p, n):
    if n < p:
        shapes[shape].chart.chart_title.text_frame.paragraphs[0].runs[0].text = '▼'
        shapes[shape].chart.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(
            0, 255, 0)
        shapes[shape].chart.chart_title.text_frame.paragraphs[0].runs[1].text = f'{round(100 * (p - n) / p,1)}%'
    elif n > p:
        shapes[shape].chart.chart_title.text_frame.paragraphs[0].runs[0].text = '▲'
        shapes[shape].chart.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(
            255, 0, 0)
        shapes[shape].chart.chart_title.text_frame.paragraphs[0].runs[1].text = f'{round(100 * (n - p) / p,1)}%'
    else:
        shapes[shape].chart.chart_title.text_frame.text = '持平'
        shapes[shape].chart.chart_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(
            0, 0, 0)
        shapes[shape].chart.chart_title.text_frame.paragraphs[0].font.size = Pt(
            16)
    chart_data = CategoryChartData()
    chart_data.categories = ['前1月均值', str_date(date)]
    chart_data.add_series('Series 2', (p, n))
    shapes[shape].chart.replace_data(chart_data)


def table_data(shapes, shape, p, n):
    if n < p:
        shapes[shape].table.cell(
            1, 0).text_frame.paragraphs[0].runs[0].text = '▼'
        shapes[shape].table.cell(
            1, 0).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(
            0, 255, 0)
        run = shapes[shape].table.cell(1, 0).text_frame.paragraphs[0].add_run()
        run.text = f' {round(100 * (p - n) / p, 1)}%'
        shapes[shape].table.cell(
            1, 0).text_frame.paragraphs[0].runs[1].font.name = '微软雅黑'
    elif n > p:
        shapes[shape].table.cell(
            1, 0).text_frame.paragraphs[0].runs[0].text = '▲'
        shapes[shape].table.cell(
            1, 0).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(
            255, 0, 0)
        run = shapes[shape].table.cell(1, 0).text_frame.paragraphs[0].add_run()
        run.text = f' {round(100 * (n - p) / p, 1)}%'
        shapes[shape].table.cell(
            1, 0).text_frame.paragraphs[0].runs[1].font.name = '微软雅黑'
    else:
        shapes[shape].table.cell(1, 0).text_frame.text = '持平'
        shapes[shape].table.cell(
            1, 0).text_frame.paragraphs[0].font.color.rgb = RGBColor(
            0, 0, 0)
        shapes[shape].table.cell(1, 0).text_frame.font.size = Pt(16)
    shapes[shape].table.cell(2, 0).text_frame.paragraphs[0].runs[0].text = f'前1月均值'
    shapes[shape].table.cell(2, 0).text_frame.paragraphs[1].runs[0].text = f'{p}'
    shapes[shape].table.cell(2, 1).text_frame.paragraphs[0].runs[0].text = f'{str_date(date)}'
    shapes[shape].table.cell(2, 1).text_frame.paragraphs[1].runs[0].text = f'{n}'


def rating_increase(shapes, increase, col):
    if increase < 0:
        run = shapes[3].table.cell(0, col).text_frame.paragraphs[0].add_run()
        run.text = f'↓{str(round(abs(increase * 100), 1))}%'
        shapes[3].table.cell(
            0, col).text_frame.paragraphs[0].runs[1].font.color.rgb = RGBColor(
            0, 255, 0)
    elif increase > 0:
        run = shapes[3].table.cell(0, col).text_frame.paragraphs[0].add_run()
        run.text = f'↑{str(round(abs(increase * 100), 1))}%'
        shapes[3].table.cell(
            0, col).text_frame.paragraphs[0].runs[1].font.color.rgb = RGBColor(
            255, 0, 0)
    else:
        run = shapes[3].table.cell(0, col).text_frame.paragraphs[0].add_run()
        run.text = '持平'
        shapes[3].table.cell(
            0,
            col).text_frame.paragraphs[0].runs[1].font.color.rgb = RGBColor(
            0,
            0,
            0)
    shapes[3].table.cell(
        0, col).text_frame.paragraphs[0].runs[1].font.size = Pt(10)
    shapes[3].table.cell(
        0, col).text_frame.paragraphs[0].runs[1].font.name = '微软雅黑'
    shapes[3].table.cell(
        0, col).text_frame.paragraphs[0].runs[1].font.bold = True


def main(argv):
    global date
    date = argv[0]

    ppt = Presentation('./origin.pptx')
    excel = xlrd.open_workbook(f'./{date}.xlsx')

    # 台组内排名
    df = pd.read_excel(f'./{date}.xlsx', sheet_name='台组内排名', usecols=[0, 1, 2, 3])[1:]
    df.sort_values(
        by='市场份额%',
        inplace=True,
        ascending=False,
        ignore_index=True)
    p_market_sort = df.loc[df['单位 >>'] == '中央电视台农业农村频道'].index[0]
    df.sort_values(
        by='市场份额%.1',
        inplace=True,
        ascending=False,
        ignore_index=True)
    market_sort = df.loc[df['单位 >>'] == '中央电视台农业农村频道'].index[0]
    in_market_sort = increase_rank(p_market_sort, market_sort)

    # 上星频道排名
    df = pd.read_excel(f'./{date}.xlsx', sheet_name='上星频道排名', usecols=[0, 1, 2, 3])[1:]
    df.columns = ['单位 >>', '市场份额%', '市场份额%.1', '频道调整']
    df.sort_values(
        by='市场份额%',
        inplace=True,
        ascending=False,
        ignore_index=True)
    p_star_sort = df.loc[df['单位 >>'] == '中央电视台农业农村频道'].index[0] + 1
    df.sort_values(
        by='市场份额%.1',
        inplace=True,
        ascending=False,
        ignore_index=True)
    star_sort = df.loc[df['单位 >>'] == '中央电视台农业农村频道'].index[0] + 1
    in_star_sort = increase_rank(p_star_sort, star_sort)

    # 第一页ppt
    slide1 = ppt.slides[0]
    shapes = slide1.shapes  # 获取所有的shape
    shapes[2].text_frame.paragraphs[1].runs[0].text = str_date(date)
    month = date[1] if date[0] == '0' else date[:2]
    shapes[3].text_frame.paragraphs[0].runs[2].text = month

    # 第二页ppt
    sheet1 = excel.sheet_by_index(0)
    market = sheet1.cell_value(2, 3)
    rating = sheet1.cell_value(2, 2)
    all_rating = sheet1.cell_value(4, 2)
    p_market: object = sheet1.cell_value(1, 3)
    p_rating = sheet1.cell_value(1, 2)
    p_all_rating = sheet1.cell_value(3, 2)
    sheet2 = excel.sheet_by_index(1)
    p_play = sheet2.cell_value(4, 5)
    play = sheet2.cell_value(4, 7)
    p_nplay = sheet2.cell_value(5, 5)
    nplay = sheet2.cell_value(5, 7)
    p_arrive = sheet1.cell_value(1, 4)
    arrive = sheet1.cell_value(2, 4)
    p_fidelity = sheet1.cell_value(1, 5)
    fidelity = sheet1.cell_value(2, 5)
    all_rating = sheet1.cell_value(4, 2)
    p_all_rating = sheet1.cell_value(3, 2)
    in_all_rating = increase_percent(p_all_rating, all_rating)
    in_arrive = increase_percent(p_arrive, arrive)
    in_fidelity = increase_percent(p_fidelity, fidelity)
    in_market = increase_percent(p_market, market)
    in_rating = increase_percent(p_rating, rating)
    in_play = increase_percent(p_play, play)
    in_nplay = increase_percent(p_nplay, nplay)

    slide2 = ppt.slides[1]
    shapes = slide2.shapes  # 获取所有的shape

    shapes[4].text_frame.paragraphs[0].text = f'{str_date(date)}，频道市场份额{market}%，较前一月均值{in_market}；'
    shapes[4].text_frame.paragraphs[1].text = f'                 收视率{rating}%，较前一月均值{in_rating}。'
    shapes[4].text_frame.paragraphs[3].text = f'{str_date(date)}，央视台组排名第{market_sort}位，较前一日{in_market_sort}；'
    shapes[4].text_frame.paragraphs[4].text = f'                 上星频道排名第{star_sort}位，较前一日{in_star_sort}。'
    shapes[4].text_frame.paragraphs[6].text = f'{str_date(date)}，电视剧市场份额{play}% ，较前一月均值{in_play}；'
    shapes[4].text_frame.paragraphs[7].text = f'                 非电视剧市场份额{nplay}% ，较前一月均值{in_nplay}。'
    for i in range(8):
        shapes[4].text_frame.paragraphs[i].font.name = '微软雅黑'
        shapes[4].text_frame.paragraphs[i].font.size = Pt(24)
        shapes[4].text_frame.paragraphs[i].font.bold = True

    # 第三页ppt

    slide3 = ppt.slides[2]
    shapes = slide3.shapes  # 获取所有的shape
    shapes[2].text_frame.paragraphs[0].text = f'{str_date(date)}，频道市场份额{in_market}，收视率{in_rating}。'
    shapes[2].text_frame.paragraphs[0].font.size = Pt(34)
    shapes[2].text_frame.paragraphs[0].font.name = '微软雅黑'
    shapes[17].text_frame.paragraphs[
        0].text = f'{str_date(date)}，频道市场份额{market}%，较前一月均值{in_market}，观众到达率{in_arrive}，观众忠实度{in_fidelity}。'
    shapes[17].text_frame.paragraphs[0].font.size = Pt(16)
    shapes[17].text_frame.paragraphs[0].font.bold = True
    shapes[17].text_frame.paragraphs[0].font.name = '微软雅黑'
    shapes[17].text_frame.paragraphs[
        1].text = f'{str_date(date)}，所有频道收视率为{all_rating}%，较前一月均值{in_all_rating}，CCTV-17收视率{rating}%，{in_rating}。'
    shapes[17].text_frame.paragraphs[1].font.size = Pt(16)
    shapes[17].text_frame.paragraphs[1].font.bold = True
    shapes[17].text_frame.paragraphs[1].font.name = '微软雅黑'

    chart_title_text(shapes, 18, p_market, market)
    chart_title_text(shapes, 20, p_rating, rating)
    chart_title_text(shapes, 19, p_all_rating, all_rating)

    table_data(shapes, 21, p_arrive, arrive)
    table_data(shapes, 22, p_fidelity, fidelity)

    # 第4页ppt
    slide4 = ppt.slides[3]
    shapes = slide4.shapes  # 获取所有的shape
    shapes[0].text_frame.paragraphs[0].text = f'{str_date(date)}，央视台组排名第{market_sort}位，较前一日{in_market_sort}。'
    shapes[0].text_frame.paragraphs[0].font.size = Pt(32)

    df = pd.read_excel(f'./{date}.xlsx', sheet_name='台组内排名', usecols=[0, 1, 2, 3])[1:]
    df.sort_values(
        by='市场份额%',
        inplace=True,
        ascending=False,
        ignore_index=True)

    shapes[2].table.cell(
    0, 0).text_frame.paragraphs[0].runs[0].text = str_p_date(date)
    shapes[2].table.cell(
    0, 3).text_frame.paragraphs[0].runs[0].text = str_date(date)
    
    shapes[2].table.cell(
        2, 2).text_frame.paragraphs[0].runs[0].text = str(
        df.iloc[0]['市场份额%'])
    for i in range(1, 21):
        shapes[2].table.cell(
            i + 2, 0).text_frame.paragraphs[0].runs[0].text = str(i)
        shapes[2].table.cell(
            i + 2,
            1).text_frame.paragraphs[0].runs[0].text = str(
            df.iloc[i]['Unnamed: 3'])
        shapes[2].table.cell(
            i + 2,
            2).text_frame.paragraphs[0].runs[0].text = str(
            df.iloc[i]['市场份额%'])
    tv_index = df.loc[df['单位 >>'] == '中央电视台农业农村频道'].index[0] + 2
    for i in range(3):
        shapes[2].table.cell(
            tv_index,
            i).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(
            184,
            1,
            0)
        shapes[2].table.cell(
            tv_index, i).text_frame.paragraphs[0].runs[0].font.bold = True

    df.sort_values(
        by='市场份额%.1',
        inplace=True,
        ascending=False,
        ignore_index=True)
    shapes[2].table.cell(
        2, 5).text_frame.paragraphs[0].runs[0].text = str(
        df.iloc[0]['市场份额%.1'])
    border = 0
    for i in range(1, 21):
        shapes[2].table.cell(
            i + 2, 3).text_frame.paragraphs[0].runs[0].text = str(i)
        shapes[2].table.cell(
            i + 2,
            4).text_frame.paragraphs[0].runs[0].text = str(
            df.iloc[i]['Unnamed: 3'])
        shapes[2].table.cell(
            i + 2,
            5).text_frame.paragraphs[0].runs[0].text = str(
            df.iloc[i]['市场份额%.1'])
        if df.iloc[i]['市场份额%.1'] < 0.83 and border == 0:
            border = i
    tv_index = df.loc[df['单位 >>'] == '中央电视台农业农村频道'].index[0] + 2

    for i in range(3):
        shapes[2].table.cell(
            tv_index,
            i +
            3).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(
            184,
            1,
            0)
        shapes[2].table.cell(
            tv_index,
            i +
            3).text_frame.paragraphs[0].runs[0].font.bold = True

    shapes[4].top = 1950000 + 215000 * (border - 2)
    shapes[3].top = 2140000 + 215000 * (border - 2)

    # 第5页ppt
    slide5 = ppt.slides[4]
    shapes = slide5.shapes  # 获取所有的shape
    shapes[0].text_frame.paragraphs[0].text = f'{str_date(date)}，上星频道排名第{star_sort}位，较前一日{in_star_sort}。'
    shapes[0].text_frame.paragraphs[0].font.size = Pt(32)

    # 上星频道排名
    df = pd.read_excel(f'./{date}.xlsx', sheet_name='上星频道排名', usecols=[0, 1, 2, 3])[1:]
    df.columns = ['单位 >>', '市场份额%', '市场份额%.1', '频道调整']
    df.sort_values(
        by='市场份额%',
        inplace=True,
        ascending=False,
        ignore_index=True)
    tv_index = df.loc[df['单位 >>'] == '中央电视台农业农村频道'].index[0] - 1
    shapes[2].table.cell(
        0, 0).text_frame.paragraphs[0].runs[0].text = str_p_date(date)
    shapes[2].table.cell(
        0, 3).text_frame.paragraphs[0].runs[0].text = str_date(date)

    for i in range(1, 11):
        shapes[2].table.cell(
            i + 1, 0).text_frame.paragraphs[0].runs[0].text = str(i)
        shapes[2].table.cell(
            i + 1, 1).text_frame.paragraphs[0].runs[0].text = str(df.iloc[i - 1]['频道调整'])
        shapes[2].table.cell(
            i + 1, 2).text_frame.paragraphs[0].runs[0].text = str(df.iloc[i - 1]['市场份额%'])

    for i in range(4):
        shapes[2].table.cell(
            13 + i,
            0).text_frame.paragraphs[0].runs[0].text = str(
            tv_index + 1 + i)
        shapes[2].table.cell(
            13 + i, 1).text_frame.paragraphs[0].runs[0].text = str(df.iloc[tv_index + i]['频道调整'])
        shapes[2].table.cell(
            13 + i, 2).text_frame.paragraphs[0].runs[0].text = str(df.iloc[tv_index + i]['市场份额%'])

    for i in range(3):
        shapes[2].table.cell(
            14,
            i).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(
            184,
            1,
            0)
        shapes[2].table.cell(
            14, i).text_frame.paragraphs[0].runs[0].font.bold = True

    df.sort_values(
        by='市场份额%.1',
        inplace=True,
        ascending=False,
        ignore_index=True)
    tv_index = df.loc[df['单位 >>'] == '中央电视台农业农村频道'].index[0] - 1
    for i in range(1, 11):
        shapes[2].table.cell(
            i + 1, 3).text_frame.paragraphs[0].runs[0].text = str(i)
        shapes[2].table.cell(
            i + 1, 4).text_frame.paragraphs[0].runs[0].text = str(df.iloc[i - 1]['频道调整'])
        shapes[2].table.cell(
            i + 1, 5).text_frame.paragraphs[0].runs[0].text = str(df.iloc[i - 1]['市场份额%.1'])

    for i in range(4):
        shapes[2].table.cell(
            13 + i,
            3).text_frame.paragraphs[0].runs[0].text = str(
            tv_index + 1 + i)
        shapes[2].table.cell(
            13 + i, 4).text_frame.paragraphs[0].runs[0].text = str(df.iloc[tv_index + i]['频道调整'])
        shapes[2].table.cell(
            13 + i, 5).text_frame.paragraphs[0].runs[0].text = str(df.iloc[tv_index + i]['市场份额%.1'])

    for i in range(3):
        shapes[2].table.cell(
            14,
            i +
            3).text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(
            184,
            1,
            0)
        shapes[2].table.cell(
            14, i + 3).text_frame.paragraphs[0].runs[0].font.bold = True

    # 第6页ppt
    slide6 = ppt.slides[5]
    shapes = slide6.shapes  # 获取所有的shape
    shapes[0].text_frame.paragraphs[0].text = f'{str_date(date)}，频道市场份额{market}%，收视率{rating}%。'
    shapes[0].text_frame.paragraphs[0].font.size = Pt(32)

    df = pd.read_excel(f'./{date}.xlsx', sheet_name='串单', usecols=[0, 5, 11, 12])

    categories = []
    for i in range(len(df)):
        start_time = str(df.iloc[i]['开始时间'])[:-3]
        if ' ' in start_time:
            start_time = start_time.split(' ')[1]
        categories.append(df.iloc[i]['名称'] + start_time)
    rating_values = [df.iloc[i]['收视率%'] for i in range(len(df))]
    market_values = [df.iloc[i]['市场份额%'] for i in range(len(df))]

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('市场份额%', market_values)
    chart_data.add_series('收视率%', rating_values)
    shapes[1].chart.replace_data(chart_data)

    if max(rating_values) > 0.24:
        shapes[1].chart.value_axis.maximum_scale = max(rating_values) + 0.04

    goodBye = df.loc[df['名称'] == '再见']
    sing = df.loc[df['名称'] == '国歌']
    df = df.drop(df[df['名称'] == '再见'].index)
    df = df.drop(df[df['名称'] == '国歌'].index)
    df = df.append(goodBye)
    df2 = pd.DataFrame(np.insert(df.values, 0, values=sing, axis=0))
    df2.columns = df.columns
    df = df2

    # 第7页ppt

    df = df.drop(df[df['名称'] == '再见'].index)
    df = df.drop(df[df['名称'] == '国歌'].index)
    slide7 = ppt.slides[6]
    shapes = slide7.shapes  # 获取所有的shape
    shapes[1].text_frame.paragraphs[0].text = f'{str_date(date)}栏目市场份额排名'
    shapes[1].text_frame.paragraphs[0].font.size = Pt(32)

    df.sort_values(
        by='市场份额%',
        inplace=True,
        ascending=False,
        ignore_index=True)
    categories = []
    for i in range(len(df)):
        start_time = str(df.iloc[i]['开始时间'])[:-3]
        if ' ' in start_time:
            start_time = start_time.split(' ')[1]
        categories.append(df.iloc[i]['名称'] + start_time)
    market_values = [df.iloc[i]['市场份额%'] for i in range(len(df))]

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('市场份额%', market_values)
    shapes[0].chart.replace_data(chart_data)

    # 第8页ppt
    slide8 = ppt.slides[7]
    shapes = slide8.shapes  # 获取所有的shape
    shapes[0].text_frame.paragraphs[0].text = f'{str_date(date)}栏目收视率排名'
    shapes[0].text_frame.paragraphs[0].font.size = Pt(32)
    df.sort_values(by='收视率%', inplace=True, ascending=False, ignore_index=True)
    categories = []
    for i in range(len(df)):
        start_time = str(df.iloc[i]['开始时间'])[:-3]
        if ' ' in start_time:
            start_time = start_time.split(' ')[1]
        categories.append(df.iloc[i]['名称'] + start_time)
    rating_values = [df.iloc[i]['收视率%'] for i in range(len(df))]

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('收视率%', rating_values)
    shapes[1].chart.replace_data(chart_data)

    # 第9页ppt
    slide9 = ppt.slides[8]
    shapes = slide9.shapes  # 获取所有的shape
    shapes[0].text_frame.paragraphs[0].text = f'{str_date(date)}分分钟市场份额%'
    shapes[0].text_frame.paragraphs[0].font.size = Pt(32)

    df = pd.read_excel(f'./{date}.xlsx', sheet_name='分分钟')[1:]
    df['单位 >>'] = df['单位 >>'].astype(str)
    hour_six = df.loc[df['单位 >>'] == '06:00:00'].index[0]
    df_copy = df.copy()
    df = df[hour_six - 1:]
    categories = []
    for i in range(len(df)):
        start_time = str(df.iloc[i]['单位 >>'])[:-3]
        if ' ' in start_time:
            start_time = start_time.split(' ')[1]
        categories.append(start_time)
    p_market_values = [df.iloc[i]['市场份额%'] for i in range(len(df))]
    p_rating_values = [df.iloc[i]['收视率%'] for i in range(len(df))]
    market_values = [df.iloc[i]['市场份额%.1'] for i in range(len(df))]
    rating_values = [df.iloc[i]['收视率%.1'] for i in range(len(df))]

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series(f'前1月均值%', p_market_values)
    chart_data.add_series(f'{str_date(date)}市场份额%', market_values)
    shapes[1].chart.replace_data(chart_data)

    # 第10页ppt
    slide10 = ppt.slides[9]
    shapes = slide10.shapes  # 获取所有的shape
    shapes[0].text_frame.paragraphs[0].text = f'{str_date(date)}分分钟收视率%'
    shapes[0].text_frame.paragraphs[0].font.size = Pt(32)

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series(f'前1月均值%', p_rating_values)
    chart_data.add_series(f'{str_date(date)}收视率%', rating_values)
    shapes[2].chart.replace_data(chart_data)

    hour_six = df_copy.loc[df_copy['单位 >>'] == '06:00:00'].index[0]
    hour_eight = df_copy.loc[df_copy['单位 >>'] == '08:00:00'].index[0]
    hour_twelve = df_copy.loc[df_copy['单位 >>'] == '12:00:00'].index[0]
    hour_fourteen = df_copy.loc[df_copy['单位 >>'] == '14:00:00'].index[0]
    hour_185 = df_copy.loc[df_copy['单位 >>'] == '18:30:00'].index[0]
    hour_22 = df_copy.loc[df_copy['单位 >>'] == '22:00:00'].index[0]
    hour_00 = df_copy.loc[df_copy['单位 >>'] == '1900-01-01 00:00:00'].index[0]
    df_zao = df_copy[hour_six - 1:hour_eight]
    df_shang = df_copy[hour_eight - 1:hour_twelve]
    df_zhong = df_copy[hour_twelve - 1:hour_fourteen]
    df_xia = df_copy[hour_fourteen - 1:hour_185]
    df_huang = df_copy[hour_185 - 1:hour_22]
    df_wan = df_copy[hour_22 - 1:hour_00]

    p_zao_rating_values = sum(
        df_zao.iloc[i]['收视率%'] for i in range(
            len(df_zao)))
    zao_rating_values = sum(df_zao.iloc[i]['收视率%.1']
                            for i in range(len(df_zao)))
    increase_zao_rating_values = (
        zao_rating_values - p_zao_rating_values) / p_zao_rating_values

    p_shang_rating_values = sum(
        df_shang.iloc[i]['收视率%'] for i in range(
            len(df_shang)))
    shang_rating_values = sum(
        df_shang.iloc[i]['收视率%.1'] for i in range(
            len(df_shang)))
    increase_shang_rating_values = (
        shang_rating_values - p_shang_rating_values) / p_shang_rating_values

    p_zhong_rating_values = sum(
        df_zhong.iloc[i]['收视率%'] for i in range(
            len(df_zhong)))
    zhong_rating_values = sum(
        df_zhong.iloc[i]['收视率%.1'] for i in range(
            len(df_zhong)))
    increase_zhong_rating_values = (
        zhong_rating_values - p_zhong_rating_values) / p_zhong_rating_values

    p_xia_rating_values = sum(
        df_xia.iloc[i]['收视率%'] for i in range(
            len(df_xia)))
    xia_rating_values = sum(df_xia.iloc[i]['收视率%.1']
                            for i in range(len(df_xia)))
    increase_xia_rating_values = (
        xia_rating_values - p_xia_rating_values) / p_xia_rating_values

    p_huang_rating_values = sum(
        df_huang.iloc[i]['收视率%'] for i in range(
            len(df_huang)))
    huang_rating_values = sum(
        df_huang.iloc[i]['收视率%.1'] for i in range(
            len(df_huang)))
    increase_huang_rating_values = (
        huang_rating_values - p_huang_rating_values) / p_huang_rating_values

    p_wan_rating_values = sum(
        df_wan.iloc[i]['收视率%'] for i in range(
            len(df_wan)))
    wan_rating_values = sum(df_wan.iloc[i]['收视率%.1']
                            for i in range(len(df_wan)))
    increase_wan_rating_values = (
        wan_rating_values - p_wan_rating_values) / p_wan_rating_values

    rating_increase(shapes, increase_zao_rating_values, 0)
    rating_increase(shapes, increase_shang_rating_values, 1)
    rating_increase(shapes, increase_zhong_rating_values, 2)
    rating_increase(shapes, increase_xia_rating_values, 3)
    rating_increase(shapes, increase_huang_rating_values, 4)
    rating_increase(shapes, increase_wan_rating_values, 5)

    ppt.save(f'./{date}频道收视日报.pptx')


if __name__ == '__main__':
    main(sys.argv[1:])
