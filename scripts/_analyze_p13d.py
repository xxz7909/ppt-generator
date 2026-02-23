# _analyze_p13d.py
from pptx import Presentation
import math

prs = Presentation('output_ppt/CCTV-17频道收视日报 0210 demo（简版）.pptx')
slide = prs.slides[12]
C = 'http://schemas.openxmlformats.org/drawingml/2006/chart'

for sh in slide.shapes:
    if sh.has_chart:
        chart = sh.chart
        from lxml import etree
        cs = chart._chartSpace
        pa = cs.find(f'.//{{{C}}}plotArea')
        layout = pa.find(f'{{{C}}}layout')
        man = layout.find(f'{{{C}}}manualLayout')
        
        px = float(man.find(f'{{{C}}}x').get('val'))
        py = float(man.find(f'{{{C}}}y').get('val'))
        pw = float(man.find(f'{{{C}}}w').get('val'))
        ph = float(man.find(f'{{{C}}}h').get('val'))
        
        shape_top = sh.top
        shape_height = sh.height
        shape_left = sh.left
        shape_width = sh.width
        
        plot_top = shape_top + int(py * shape_height)
        plot_height = int(ph * shape_height)
        plot_bottom = plot_top + plot_height
        plot_left = shape_left + int(px * shape_width)
        plot_width = int(pw * shape_width)
        
        print(f"Plot: top={plot_top}, bottom={plot_bottom}, height={plot_height}")
        print(f"Plot: left={plot_left}, width={plot_width}")
        
        # Get data
        s0 = list(chart.series[0].values)
        s1 = list(chart.series[1].values)
        num_cats = len(s0)
        
        # Max per category
        cat_max = []
        for i in range(num_cats):
            v0 = s0[i] if s0[i] is not None else 0
            v1 = s1[i] if s1[i] is not None else 0
            cat_max.append(max(v0, v1))
        
        overall_max = max(cat_max)
        print(f"Overall max: {overall_max}")
        
        # Try different axis max values
        for axis_max in [1400, 1500, 1600, 1800]:
            print(f"\n--- axis_max = {axis_max} ---")
            non_gap = [(i, cat_max[i]) for i in range(num_cats) if cat_max[i] > 0]
            
            # Get labels sorted by X
            labels = []
            for s in slide.shapes:
                if getattr(s, 'has_text_frame', False):
                    t = s.text_frame.text.strip()
                    if '%' in t and len(t) < 10:
                        labels.append((s.left, s.top, s.width, s.height, t))
            labels.sort()
            
            # Filter out legend label (the one at very top, near chart top)
            cat_labels = [l for l in labels if l[1] > shape_top + shape_height * 0.3]
            
            for j, (ci, mv) in enumerate(non_gap):
                bar_top = plot_bottom - int((mv / axis_max) * plot_height)
                if j < len(cat_labels):
                    lbl = cat_labels[j]
                    lbl_bottom = lbl[1] + lbl[3]
                    gap = bar_top - lbl_bottom
                    print(f"  cat {ci:2d} max={mv:7.1f} bar_top={bar_top:7d}  lbl_top={lbl[1]:7d} lbl_h={lbl[3]:5d} lbl_bottom={lbl_bottom:7d} gap_to_bar={gap:7d} ({gap/914400*2.54:.2f}cm)  text={lbl[4]}")
