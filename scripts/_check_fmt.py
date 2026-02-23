from pptx import Presentation
from lxml import etree

A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'

prs = Presentation('output_ppt/result_0215.pptx')
slide = prs.slides[3]  # page 4

for shp in slide.shapes:
    if not getattr(shp, 'has_table', False):
        continue
    t = shp.table
    n_cols = len(t.columns)
    for r in range(len(t.rows)):
        cell_texts = [t.cell(r, c).text for c in range(n_cols)]
        # check col 1 formatting
        tc = t.cell(r, 1)._tc
        runs = tc.findall(f'.//{{{A_NS}}}r')
        has_green = False
        has_bold = False
        for run in runs:
            rpr = run.find(f'{{{A_NS}}}rPr')
            if rpr is not None:
                if rpr.get('b') == '1':
                    has_bold = True
                fill = rpr.find(f'{{{A_NS}}}solidFill')
                if fill is not None:
                    clr = fill.find(f'{{{A_NS}}}srgbClr')
                    if clr is not None and clr.get('val') == '5B7F5B':
                        has_green = True
        if has_green or has_bold:
            print(f'  row{r}: bold={has_bold}, green={has_green}, texts={cell_texts}')
