"""检查 demo 表头表格字体和单元格属性"""
from pptx import Presentation
from lxml import etree

a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
demo = Presentation('output_ppt/CCTV-17频道收视日报 0210 demo（简版）.pptx')

for pi in (8, 9):
    slide = demo.slides[pi]
    print(f'\n=== Page {pi+1} ===')
    for shp in slide.shapes:
        if not getattr(shp, 'has_table', False):
            continue
        tbl = shp.table
        if len(tbl.columns) != 8:
            continue
        if not tbl.rows[0].cells[0].text.strip():
            continue
        # header table
        print(f'Header table: left={shp.left}, width={shp.width}')
        for ci in range(8):
            tc = tbl.rows[0].cells[ci]._tc
            tc_body = tc.find(f'.//{{{a_ns}}}txBody')
            runs = tc_body.findall(f'.//{{{a_ns}}}r') if tc_body is not None else []
            for r in runs:
                rPr = r.find(f'{{{a_ns}}}rPr')
                sz = rPr.get('sz') if rPr is not None else None
                b = rPr.get('b') if rPr is not None else None
                txt_el = r.find(f'{{{a_ns}}}t')
                txt = txt_el.text if txt_el is not None else ''
                print(f'  col[{ci}] text="{txt}" sz={sz} b={b} w={tbl.columns[ci].width} w_cm={round(tbl.columns[ci].width/914400*2.54,2)}')
            # check bodyPr
            bodyPr = tc_body.find(f'{{{a_ns}}}bodyPr') if tc_body is not None else None
            if bodyPr is not None:
                attrs = dict(bodyPr.attrib)
                print(f'    bodyPr: {attrs}')
            # check pPr
            paras = tc_body.findall(f'{{{a_ns}}}p') if tc_body is not None else []
            for p in paras:
                pPr = p.find(f'{{{a_ns}}}pPr')
                if pPr is not None:
                    algn = pPr.get('algn')
                    defRPr = pPr.find(f'{{{a_ns}}}defRPr')
                    dsz = defRPr.get('sz') if defRPr is not None else None
                    print(f'    pPr: algn={algn} defRPr.sz={dsz}')
