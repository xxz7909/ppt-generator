"""检查 demo 模板 P6 图表 catAx 原始字号"""
from pptx import Presentation
from lxml import etree

C = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
A = 'http://schemas.openxmlformats.org/drawingml/2006/main'

prs = Presentation('output_ppt/CCTV-17频道收视日报 0210 demo（简版）.pptx')
slide = prs.slides[5]  # P6

for shp in slide.shapes:
    if not getattr(shp, 'has_chart', False):
        continue
    cs = shp.chart._chartSpace
    pa = cs.find(f'.//{{{C}}}plotArea')

    # catAx font sizes
    for cat_ax in pa.findall(f'{{{C}}}catAx'):
        txPr = cat_ax.find(f'{{{C}}}txPr')
        if txPr is not None:
            for defRPr in txPr.findall(f'.//{{{A}}}defRPr'):
                sz = defRPr.get('sz', 'unset')
                print(f'demo catAx defRPr sz={sz}')

    # Check labels
    max_len = 0
    for str_cache in pa.findall(f'.//{{{C}}}strCache'):
        for pt in str_cache.findall(f'{{{C}}}pt'):
            v = pt.find(f'{{{C}}}v')
            if v is not None and v.text:
                l = len(v.text)
                if l > max_len:
                    max_len = l
                    print(f'  label len={l}: {v.text[:40]}')
    print(f'max_label_len={max_len}')

# Check generated output
print('\n=== Generated ===')
prs2 = Presentation('output_ppt/test_axis_adaptive.pptx')
slide2 = prs2.slides[5]
for shp in slide2.shapes:
    if not getattr(shp, 'has_chart', False):
        continue
    cs = shp.chart._chartSpace
    pa = cs.find(f'.//{{{C}}}plotArea')
    for cat_ax in pa.findall(f'{{{C}}}catAx'):
        txPr = cat_ax.find(f'{{{C}}}txPr')
        if txPr is not None:
            for defRPr in txPr.findall(f'.//{{{A}}}defRPr'):
                sz = defRPr.get('sz', 'unset')
                print(f'generated catAx defRPr sz={sz}')
    max_len = 0
    for str_cache in pa.findall(f'.//{{{C}}}strCache'):
        for pt in str_cache.findall(f'{{{C}}}pt'):
            v = pt.find(f'{{{C}}}v')
            if v is not None and v.text:
                l = len(v.text)
                if l > max_len:
                    max_len = l
    print(f'max_label_len={max_len}')
