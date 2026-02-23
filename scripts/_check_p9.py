# -*- coding: utf-8 -*-
"""Inspect page 9 chart and table structure."""
from pptx import Presentation
from pptx.util import Emu, Cm
from lxml import etree
import re

prs = Presentation("output_ppt/result_0215_v4.pptx")
slide = prs.slides[8]

print("=== Page 9 shapes ===")
for i, shp in enumerate(slide.shapes):
    tp = "table" if getattr(shp, "has_table", False) else ("chart" if getattr(shp, "has_chart", False) else "other")
    print(f"  [{i}] {tp}  left={shp.left}  top={shp.top}  w={shp.width}  h={shp.height}")
    
    if tp == "chart":
        chart = shp.chart
        cats = list(chart.plots[0].categories)
        print(f"      time range: first={cats[0]}, last={cats[-1]}, count={len(cats)}")
        # Check plot area layout from chart XML
        chart_xml = etree.tostring(chart._chartSpace, pretty_print=True).decode()
        for m in re.finditer(r'<c:plotArea>(.*?)</c:plotArea>', chart_xml, re.DOTALL):
            for m2 in re.finditer(r'<c:(x|y|w|h) val="([^"]+)"', m.group(1)):
                print(f"      plotArea {m2.group(1)} = {m2.group(2)}")

# Also check demo template
print("\n=== Demo template page 9 ===")
demo = Presentation("output_ppt/CCTV-17频道收视日报 0210 demo（简版）.pptx")
dslide = demo.slides[8]
for i, shp in enumerate(dslide.shapes):
    tp = "table" if getattr(shp, "has_table", False) else ("chart" if getattr(shp, "has_chart", False) else "other")
    extra = ""
    if tp == "table":
        tbl = shp.table
        ncols = len(tbl.columns)
        cell0 = tbl.rows[0].cells[0].text[:30]
        extra = f"  cols={ncols} cell0={repr(cell0)}"
        # Print column widths
        for j in range(ncols):
            w = tbl.columns[j].width
            t = tbl.rows[0].cells[j].text[:40]
            extra += f"\n        col[{j}] w={w} text={repr(t)}"
    if tp == "chart":
        chart = shp.chart
        cats = list(chart.plots[0].categories)
        extra = f"  first={cats[0]} last={cats[-1]} count={len(cats)}"
        chart_xml = etree.tostring(chart._chartSpace, pretty_print=True).decode()
        for m2 in re.finditer(r'<c:(x|y|w|h) val="([^"]+)"', chart_xml):
            extra += f"\n      plotArea {m2.group(1)} = {m2.group(2)}"
    print(f"  [{i}] {tp}  left={shp.left}  top={shp.top}  w={shp.width}  h={shp.height}{extra}")
