"""对比 demo 和生成文件中 Page 9/10 的表格结构"""
from pptx import Presentation
from pptx.util import Emu, Pt
from lxml import etree
import sys

def dump_tables(prs_path, label):
    prs = Presentation(prs_path)
    for page_idx in (8, 9):
        slide = prs.slides[page_idx]
        print(f'\n{"="*60}')
        print(f'{label} - Page {page_idx+1}')
        print(f'{"="*60}')
        
        for si, shp in enumerate(slide.shapes):
            if getattr(shp, 'has_table', False):
                tbl = shp.table
                ncols = len(tbl.columns)
                nrows = len(tbl.rows)
                texts = []
                for ri in range(nrows):
                    row_texts = [tbl.rows[ri].cells[ci].text.strip()[:8] for ci in range(min(ncols, 10))]
                    texts.append(row_texts)
                
                # 检查表格边框
                tbl_xml = etree.tostring(shp._element, pretty_print=True).decode()
                has_borders = 'ln' in tbl_xml and 'solidFill' in tbl_xml
                
                print(f'\n  Shape[{si}] Table: {nrows}r x {ncols}c')
                print(f'    left={shp.left}, top={shp.top}, width={shp.width}, height={shp.height}')
                print(f'    left_cm={round(shp.left/914400*2.54,2)}, top_cm={round(shp.top/914400*2.54,2)}')
                print(f'    width_cm={round(shp.width/914400*2.54,2)}, height_cm={round(shp.height/914400*2.54,2)}')
                
                for ri, row_texts in enumerate(texts):
                    print(f'    row[{ri}]: {row_texts}')
                
                # 检查行高
                for ri in range(nrows):
                    print(f'    row[{ri}].height = {tbl.rows[ri].height}')
                
                # 检查是否有边框线
                # 查看 tblPr
                a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                tbl_elem = shp._element.find(f'.//{{{a_ns}}}tbl')
                if tbl_elem is not None:
                    tblPr = tbl_elem.find(f'{{{a_ns}}}tblPr')
                    if tblPr is not None:
                        print(f'    tblPr attribs: {dict(tblPr.attrib)}')
                        # bandRow, bandCol etc
                
                # 检查单元格边框 (第一个单元格)
                tc = tbl.rows[0].cells[0]._tc
                tc_xml_short = etree.tostring(tc, pretty_print=True).decode()[:500]
                # 只看有没有 tcPr 中的 border
                borders_found = []
                tcPr = tc.find(f'{{{a_ns}}}tcPr')
                if tcPr is not None:
                    for border_name in ['lnL', 'lnR', 'lnT', 'lnB']:
                        ln = tcPr.find(f'{{{a_ns}}}{border_name}')
                        if ln is not None:
                            fill = ln.find(f'{{{a_ns}}}solidFill')
                            nofill = ln.find(f'{{{a_ns}}}noFill')
                            if nofill is not None:
                                borders_found.append(f'{border_name}=noFill')
                            elif fill is not None:
                                clr = fill.find(f'{{{a_ns}}}srgbClr')
                                if clr is not None:
                                    borders_found.append(f'{border_name}={clr.get("val")}')
                                else:
                                    borders_found.append(f'{border_name}=solidFill')
                            else:
                                w = ln.get('w', '?')
                                borders_found.append(f'{border_name}=w:{w}')
                    print(f'    cell[0,0] borders: {borders_found}')
            
            elif getattr(shp, 'has_chart', False):
                print(f'\n  Shape[{si}] Chart')
                print(f'    left={shp.left}, top={shp.top}')
            else:
                txt = getattr(shp, 'text', '')[:20]
                print(f'\n  Shape[{si}] {shp.shape_type}: "{txt}"')
                print(f'    left={shp.left}, top={shp.top}')

# Demo
demo_path = 'output_ppt/CCTV-17频道收视日报 0210 demo（简版）.pptx'
dump_tables(demo_path, 'DEMO')

# Generated
gen_path = 'output_ppt/test_0222_aligned.pptx'
dump_tables(gen_path, 'GENERATED')
