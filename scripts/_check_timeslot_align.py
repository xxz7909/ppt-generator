"""检查分分钟页分时段标签是否与图表横坐标对齐"""
from pptx import Presentation
from pptx.util import Emu
from lxml import etree
import re

prs = Presentation('output_ppt/test_0222_aligned.pptx')

for page_idx in (8, 9):
    slide = prs.slides[page_idx]
    title = f'Page {page_idx+1}'
    print(f'\n=== {title} ===')

    chart_shape = None
    tables = []
    for shp in slide.shapes:
        if getattr(shp, 'has_chart', False):
            chart_shape = shp
        if getattr(shp, 'has_table', False):
            tables.append(shp)

    if chart_shape:
        chart = chart_shape.chart
        cats = list(chart.plots[0].categories)
        chart_xml = etree.tostring(chart._chartSpace).decode()
        layout_match = re.search(r'<c:plotArea>(.*?)</c:plotArea>', chart_xml, re.DOTALL)
        pa_x = pa_w = 0
        if layout_match:
            x_m = re.search(r'<c:x val="([^"]+)"', layout_match.group(1))
            w_m = re.search(r'<c:w val="([^"]+)"', layout_match.group(1))
            if x_m and w_m:
                pa_x = float(x_m.group(1))
                pa_w = float(w_m.group(1))
        plot_left = chart_shape.left + int(pa_x * chart_shape.width)
        plot_width = int(pa_w * chart_shape.width)
        print(f'  Chart: left={chart_shape.left}, width={chart_shape.width}')
        print(f'  PlotArea: left={plot_left} ({round(plot_left/914400*2.54, 2)}cm), width={plot_width} ({round(plot_width/914400*2.54, 2)}cm)')
        print(f'  cats[0]={cats[0]}, cats[-1]={cats[-1]}')

        # 计算06:00和10:00在图表中的位置
        def parse_min(t):
            t = str(t).strip()
            if ' ' in t:
                t = t.split(' ')[1]
            parts = t.replace('.', ':').split(':')
            h, m = int(parts[0]), int(parts[1]) if len(parts) > 1 else 0
            if h < 5:
                h += 24
            return h * 60 + m

        chart_start = parse_min(cats[0])
        chart_end = parse_min(cats[-1]) + 1
        chart_dur = chart_end - chart_start

        time_marks = [360, 600, 720, 810, 1020, 1110, 1230, 1350, 1500]  # 06:00,10:00,12:00,13:30,17:00,18:30,20:30,22:30,25:00
        labels = ['06:00', '10:00', '12:00', '13:30', '17:00', '18:30', '20:30', '22:30', '25:00']
        print(f'  时段边界在图表中的位置:')
        for lbl, t in zip(labels, time_marks):
            pos = plot_left + int((t - chart_start) / chart_dur * plot_width)
            print(f'    {lbl} → {pos} EMU ({round(pos/914400*2.54, 2)}cm)')

    for shp in tables:
        tbl = shp.table
        ncols = len(tbl.columns)
        texts = [tbl.rows[0].cells[i].text.strip()[:6] for i in range(min(ncols, 10))]
        print(f'\n  Table ({ncols} cols): left={shp.left} ({round(shp.left/914400*2.54, 2)}cm), width={shp.width}')
        cum = shp.left
        for ci in range(min(ncols, 10)):
            w = tbl.columns[ci].width
            print(f'    col[{ci}] "{texts[ci]}": left={cum}, width={w}, right={cum+w}')
            cum += w
