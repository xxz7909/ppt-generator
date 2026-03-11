"""验证修复：1) 表头列分隔线对齐 2) 节目表格z-order"""
from pptx import Presentation
from lxml import etree
import re

def parse_min(t):
    t = str(t).strip()
    if ' ' in t:
        t = t.split(' ')[1]
    parts = t.replace('.', ':').split(':')
    h, m = int(parts[0]), int(parts[1]) if len(parts) > 1 else 0
    if h < 5:
        h += 24
    return h * 60 + m

prs = Presentation('output_ppt/test_0222_v2.pptx')

for page_idx in (8, 9):
    slide = prs.slides[page_idx]
    print(f'\n=== Page {page_idx+1} ===')

    chart_shape = None
    header_tbl = None
    divider_tbl = None
    prog_tbl = None

    for shp in slide.shapes:
        if getattr(shp, 'has_chart', False):
            chart_shape = shp
        if getattr(shp, 'has_table', False):
            tbl = shp.table
            ncols = len(tbl.columns)
            if ncols in (20, 21):
                prog_tbl = shp
            elif ncols == 8:
                if tbl.rows[0].cells[0].text.strip():
                    header_tbl = shp
                else:
                    divider_tbl = shp

    # Chart plot area
    chart = chart_shape.chart
    cats = list(chart.plots[0].categories)
    chart_xml = etree.tostring(chart._chartSpace).decode()
    layout_match = re.search(r'<c:plotArea>(.*?)</c:plotArea>', chart_xml, re.DOTALL)
    pa_x = pa_w = 0
    if layout_match:
        x_m = re.search(r'<c:x val="([^"]+)"', layout_match.group(1))
        w_m = re.search(r'<c:w val="([^"]+)"', layout_match.group(1))
        if x_m and w_m:
            pa_x = float(x_m.group(1))
            pa_w = float(w_m.group(1))
    plot_left = chart_shape.left + int(pa_x * chart_shape.width)
    plot_width = int(pa_w * chart_shape.width)
    chart_start = parse_min(cats[0])
    chart_end = parse_min(cats[-1]) + 1
    chart_dur = chart_end - chart_start

    # Expected boundary positions
    boundaries_min = [360, 600, 720, 810, 1020, 1110, 1230, 1350, 1500]
    boundary_pos = [plot_left + int((t - chart_start) / chart_dur * plot_width) for t in boundaries_min]
    labels = ['06:00', '10:00', '12:00', '13:30', '17:00', '18:30', '20:30', '22:30', '25:00']

    # Check 1: header table unchanged position
    if header_tbl:
        print(f'  Header table: left={header_tbl.left}, width={header_tbl.width}')
        cum = header_tbl.left
        tbl_obj = header_tbl.table
        for ci in range(len(tbl_obj.columns)):
            w = tbl_obj.columns[ci].width
            right = cum + w
            # Find which boundary this right edge should align with
            if ci < len(tbl_obj.columns) - 1 and ci + 1 < len(boundary_pos):
                expected = boundary_pos[ci + 1]
                diff = right - expected
                print(f'    col[{ci}] right={right}, expected {labels[ci+1]}={expected}, diff={diff} EMU')
            else:
                print(f'    col[{ci}] right={right} (table end)')
            cum = right

    # Check 2: z-order
    parent = slide.shapes._spTree
    shape_elems = list(parent)
    if prog_tbl and divider_tbl:
        prog_idx = shape_elems.index(prog_tbl._element)
        div_idx = shape_elems.index(divider_tbl._element)
        if prog_idx > div_idx:
            print(f'  ✅ Z-order OK: prog(idx={prog_idx}) above divider(idx={div_idx})')
        else:
            print(f'  ❌ Z-order BAD: prog(idx={prog_idx}) below divider(idx={div_idx})')
