# -*- coding: utf-8 -*-
"""
PPT报告生成器 - 幻灯片构建工具
提供图表、表格、文本框等构建辅助函数
"""
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
import copy

from ppt_config import Colors, Fonts, Layout


# ═══════════════════ 文本工具 ═══════════════════

def add_textbox(slide, left, top, width, height, text,
                font_name=Fonts.MAIN, font_size=Fonts.BODY_SIZE,
                font_color=Colors.DARK_GRAY, bold=False,
                alignment=PP_ALIGN.LEFT, word_wrap=True,
                anchor=MSO_ANCHOR.TOP):
    """添加文本框"""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = word_wrap

    try:
        tf.auto_size = None
    except:
        pass

    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = alignment
    return txbox


def add_rich_textbox(slide, left, top, width, height, segments,
                     alignment=PP_ALIGN.LEFT, line_spacing=1.2):
    """
    添加富文本框
    segments: list of (text, font_size, color, bold) tuples
    """
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = alignment
    if line_spacing:
        p.line_spacing = Pt(line_spacing * 12)

    for i, (text, font_size, color, bold) in enumerate(segments):
        if i == 0:
            run = p.runs[0] if p.runs else p.add_run()
        else:
            run = p.add_run()
        run.text = text
        run.font.name = Fonts.MAIN
        run.font.size = font_size
        run.font.color.rgb = color
        run.font.bold = bold

    return txbox


def add_title(slide, text, font_size=Fonts.TITLE_SIZE):
    """添加标题"""
    return add_textbox(
        slide, Layout.TITLE_LEFT, Layout.TITLE_TOP,
        Layout.TITLE_WIDTH, Layout.TITLE_HEIGHT,
        text, font_size=font_size, font_color=Colors.PRIMARY,
        bold=True, alignment=PP_ALIGN.LEFT
    )


def add_subtitle(slide, text, top=None):
    """添加副标题/描述文本"""
    return add_textbox(
        slide, Layout.DESC_LEFT, top or Layout.DESC_TOP,
        Layout.DESC_WIDTH, Layout.DESC_HEIGHT,
        text, font_size=Fonts.BODY_SIZE, font_color=Colors.MEDIUM_GRAY,
        bold=False, alignment=PP_ALIGN.LEFT
    )


def add_page_number(slide, num):
    """添加页码"""
    add_textbox(
        slide, Layout.PAGE_NUM_LEFT, Layout.PAGE_NUM_TOP,
        Layout.PAGE_NUM_WIDTH, Layout.PAGE_NUM_HEIGHT,
        str(num), font_size=Fonts.PAGE_NUM_SIZE,
        font_color=Colors.PRIMARY, bold=True,
        alignment=PP_ALIGN.CENTER
    )


def add_note(slide, text, top=None):
    """添加底部注释"""
    return add_textbox(
        slide, Layout.NOTE_LEFT, top or Layout.NOTE_BOTTOM,
        Layout.NOTE_WIDTH, Layout.NOTE_HEIGHT,
        text, font_size=Fonts.TINY_SIZE,
        font_color=Colors.MEDIUM_GRAY, bold=False
    )


# ═══════════════════ 表格工具 ═══════════════════

def add_styled_table(slide, left, top, width, height,
                     headers, data_rows, highlight_keywords=None,
                     col_widths=None):
    """
    添加美化表格
    headers: list of str
    data_rows: list of list of str
    highlight_keywords: list of str - 包含这些关键字的行高亮
    col_widths: list of Emu - 列宽
    """
    n_rows = len(data_rows) + 1  # +1 for header
    n_cols = len(headers)

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    # 设置列宽
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    # 表头样式
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        _style_cell(cell, Fonts.TABLE_HEADER_SIZE, Colors.TABLE_HEADER_FG,
                    bold=True, bg_color=Colors.TABLE_HEADER_BG,
                    alignment=PP_ALIGN.CENTER)

    # 数据行
    for i, row in enumerate(data_rows):
        is_highlight = False
        if highlight_keywords:
            row_text = ' '.join(str(v) for v in row)
            is_highlight = any(kw in row_text for kw in highlight_keywords)

        bg = Colors.TABLE_HIGHLIGHT if is_highlight else \
             (Colors.TABLE_ROW_EVEN if i % 2 == 0 else Colors.TABLE_ROW_ODD)

        for j, val in enumerate(row):
            if j >= n_cols:
                break
            cell = table.cell(i + 1, j)
            cell.text = str(val) if val is not None else ""
            _style_cell(cell, Fonts.TABLE_BODY_SIZE, Colors.DARK_GRAY,
                       bold=is_highlight, bg_color=bg,
                       alignment=PP_ALIGN.CENTER)

    return table_shape


def _style_cell(cell, font_size, font_color, bold=False,
                bg_color=None, alignment=PP_ALIGN.CENTER):
    """设置单元格样式"""
    # 设置背景色
    if bg_color:
        tcPr = cell._tc.get_or_add_tcPr()
        solidFill = tcPr.makeelement(qn('a:solidFill'), {})
        srgbClr = solidFill.makeelement(qn('a:srgbClr'),
                                        {'val': f'{bg_color}'})
        solidFill.append(srgbClr)
        tcPr.append(solidFill)

    # 设置文本样式
    for paragraph in cell.text_frame.paragraphs:
        paragraph.alignment = alignment
        paragraph.font.name = Fonts.MAIN
        paragraph.font.size = font_size
        paragraph.font.color.rgb = font_color
        paragraph.font.bold = bold

    # 单元格内边距
    cell.margin_top = Pt(2)
    cell.margin_bottom = Pt(2)
    cell.margin_left = Pt(4)
    cell.margin_right = Pt(4)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_ranking_table(slide, left, top, width, height,
                      prev_date, curr_date, rankings,
                      max_rows=20, highlight_keywords=None):
    """
    添加排名对比表格（左右两列对比）
    rankings: list of ChannelRankItem
    """
    # 按当日份额排序
    sorted_curr = sorted(rankings, key=lambda x: x.current_share, reverse=True)
    sorted_prev = sorted(rankings, key=lambda x: x.prev_share, reverse=True)

    display_count = min(max_rows, len(sorted_curr))

    headers = [prev_date, '', '', curr_date, '', '']
    sub_headers = ['排名', '频道', '市场份额%', '排名', '频道', '市场份额%']

    n_rows = display_count + 2  # 2 header rows + data
    n_cols = 6

    if display_count == 0:
        return None

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    # 第一行表头 (日期)
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        _style_cell(cell, Fonts.TABLE_HEADER_SIZE, Colors.TABLE_HEADER_FG,
                    bold=True, bg_color=Colors.TABLE_HEADER_BG)

    # 合并第一行的日期单元格
    table.cell(0, 0).merge(table.cell(0, 2))
    table.cell(0, 3).merge(table.cell(0, 5))
    table.cell(0, 0).text = prev_date
    table.cell(0, 3).text = curr_date
    _style_cell(table.cell(0, 0), Fonts.TABLE_HEADER_SIZE, Colors.TABLE_HEADER_FG,
                bold=True, bg_color=Colors.TABLE_HEADER_BG)
    _style_cell(table.cell(0, 3), Fonts.TABLE_HEADER_SIZE, Colors.TABLE_HEADER_FG,
                bold=True, bg_color=Colors.TABLE_HEADER_BG)

    # 第二行子表头
    for j, h in enumerate(sub_headers):
        cell = table.cell(1, j)
        cell.text = h
        _style_cell(cell, Fonts.TABLE_BODY_SIZE, Colors.PRIMARY,
                    bold=True, bg_color=Colors.TABLE_ROW_EVEN)

    # 数据行
    for i in range(display_count):
        prev_item = sorted_prev[i] if i < len(sorted_prev) else None
        curr_item = sorted_curr[i] if i < len(sorted_curr) else None

        row_idx = i + 2
        is_highlight = False
        if highlight_keywords:
            for item in [prev_item, curr_item]:
                if item:
                    if any(kw in item.name or kw in item.short_name for kw in highlight_keywords):
                        is_highlight = True
                        break

        bg = Colors.TABLE_HIGHLIGHT if is_highlight else \
             (Colors.TABLE_ROW_EVEN if i % 2 == 0 else Colors.TABLE_ROW_ODD)

        # 前日数据
        if prev_item:
            row_data = [str(i+1),
                       prev_item.short_name or prev_item.name,
                       f'{prev_item.prev_share:.3f}']
        else:
            row_data = ['', '', '']

        for j, val in enumerate(row_data):
            cell = table.cell(row_idx, j)
            cell.text = val
            _style_cell(cell, Fonts.TABLE_BODY_SIZE, Colors.DARK_GRAY,
                       bold=is_highlight, bg_color=bg)

        # 当日数据
        if curr_item:
            row_data = [str(i+1),
                       curr_item.short_name or curr_item.name,
                       f'{curr_item.current_share:.3f}']
        else:
            row_data = ['', '', '']

        for j, val in enumerate(row_data):
            cell = table.cell(row_idx, j + 3)
            cell.text = val
            _style_cell(cell, Fonts.TABLE_BODY_SIZE, Colors.DARK_GRAY,
                       bold=is_highlight, bg_color=bg)

    return table_shape


# ═══════════════════ 图表工具 ═══════════════════

def add_column_chart(slide, left, top, width, height,
                     categories, series_data, chart_title=None,
                     series_colors=None, show_data_labels=True,
                     gap_width=100):
    """
    添加柱状图
    categories: list of str
    series_data: dict of {series_name: [values]}
    """
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series_data.items():
        chart_data.add_series(name, values)

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        left, top, width, height, chart_data
    )
    chart = chart_shape.chart

    # 基本样式
    chart.has_legend = len(series_data) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.name = Fonts.MAIN
        chart.legend.font.size = Fonts.SMALL_SIZE

    if chart_title:
        chart.has_title = True
        chart.chart_title.text_frame.text = chart_title
        chart.chart_title.text_frame.paragraphs[0].font.name = Fonts.MAIN
        chart.chart_title.text_frame.paragraphs[0].font.size = Fonts.BODY_SIZE
    else:
        chart.has_title = False

    # 系列颜色
    plot = chart.plots[0]
    plot.gap_width = gap_width

    if series_colors:
        for i, color in enumerate(series_colors):
            if i < len(plot.series):
                series = plot.series[i]
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = color

    # 数据标签
    if show_data_labels:
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.name = Fonts.MAIN
        data_labels.font.size = Fonts.TINY_SIZE
        data_labels.font.color.rgb = Colors.DARK_GRAY
        data_labels.number_format = '0.000'
        data_labels.number_format_is_linked = False

    # 分类轴样式
    cat_axis = chart.category_axis
    cat_axis.has_minor_gridlines = False
    cat_axis.has_major_gridlines = False
    cat_axis.tick_labels.font.name = Fonts.MAIN
    cat_axis.tick_labels.font.size = Fonts.TINY_SIZE

    # 值轴样式
    val_axis = chart.value_axis
    val_axis.has_minor_gridlines = False
    val_axis.major_gridlines.format.line.color.rgb = Colors.LIGHT_GRAY
    val_axis.tick_labels.font.name = Fonts.MAIN
    val_axis.tick_labels.font.size = Fonts.TINY_SIZE

    return chart_shape


def add_bar_chart(slide, left, top, width, height,
                  categories, series_data, series_colors=None,
                  show_data_labels=True):
    """添加条形图（水平柱状图）"""
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series_data.items():
        chart_data.add_series(name, values)

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        left, top, width, height, chart_data
    )
    chart = chart_shape.chart
    chart.has_legend = False
    chart.has_title = False

    plot = chart.plots[0]
    if series_colors:
        for i, color in enumerate(series_colors):
            if i < len(plot.series):
                series = plot.series[i]
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = color

    if show_data_labels:
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.name = Fonts.MAIN
        data_labels.font.size = Fonts.TINY_SIZE
        data_labels.number_format = '0.000'
        data_labels.number_format_is_linked = False

    # 轴样式
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.name = Fonts.MAIN
    cat_axis.tick_labels.font.size = Fonts.TINY_SIZE

    val_axis = chart.value_axis
    val_axis.has_minor_gridlines = False
    val_axis.tick_labels.font.name = Fonts.MAIN
    val_axis.tick_labels.font.size = Fonts.TINY_SIZE

    return chart_shape


def add_area_chart(slide, left, top, width, height,
                   categories, series_data, series_colors=None):
    """添加面积图"""
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series_data.items():
        chart_data.add_series(name, values)

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.AREA,
        left, top, width, height, chart_data
    )
    chart = chart_shape.chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.name = Fonts.MAIN
    chart.legend.font.size = Fonts.SMALL_SIZE
    chart.has_title = False

    # 系列颜色和透明度
    plot = chart.plots[0]
    if series_colors:
        for i, color in enumerate(series_colors):
            if i < len(plot.series):
                series = plot.series[i]
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = color

    # 轴样式
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.name = Fonts.MAIN
    cat_axis.tick_labels.font.size = Fonts.TINY_SIZE
    try:
        from pptx.enum.chart import XL_TICK_LABEL_POSITION
        cat_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
    except:
        pass

    val_axis = chart.value_axis
    val_axis.has_minor_gridlines = False
    val_axis.major_gridlines.format.line.color.rgb = Colors.LIGHT_GRAY
    val_axis.tick_labels.font.name = Fonts.MAIN
    val_axis.tick_labels.font.size = Fonts.TINY_SIZE

    return chart_shape


# ═══════════════════ 形状工具 ═══════════════════

def add_rect(slide, left, top, width, height, fill_color=None,
             border_color=None, border_width=Pt(0)):
    """添加矩形"""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color=None):
    """添加圆角矩形"""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    shape.line.fill.background()
    return shape


def format_change_str(value, suffix='%', positive_prefix='↑', negative_prefix='↓'):
    """格式化变化值字符串"""
    if value > 0:
        return f'{positive_prefix}{abs(value):.0f}{suffix}'
    elif value < 0:
        return f'{negative_prefix}{abs(value):.0f}{suffix}'
    else:
        return f'持平'


def format_change_color(value):
    """根据变化值返回颜色"""
    if value > 0:
        return Colors.ACCENT_GREEN
    elif value < 0:
        return Colors.ACCENT_RED
    else:
        return Colors.MEDIUM_GRAY
